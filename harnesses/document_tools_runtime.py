#!/usr/bin/env python3
from __future__ import annotations

import argparse
import copy
import json
import os
import shutil
import tempfile
from pathlib import Path
from typing import Any

from docx import Document
from docx.oxml.ns import qn
from openpyxl import load_workbook
from openpyxl.styles import PatternFill
from pptx import Presentation
from pypdf import PdfReader, PdfWriter
import pdfplumber


def _ensure_parent(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)


def _save_atomic(src_path: Path, save_fn) -> None:
    fd, tmp_name = tempfile.mkstemp(suffix=src_path.suffix)
    os.close(fd)
    tmp_path = Path(tmp_name)
    try:
        save_fn(tmp_path)
        shutil.move(str(tmp_path), str(src_path))
    finally:
        if tmp_path.exists():
            tmp_path.unlink(missing_ok=True)


def tool_get_task_metadata() -> dict[str, Any]:
    meta_path = Path('/root/task_metadata.json')
    if not meta_path.exists():
        return {'ok': False, 'error': 'task metadata not found at /root/task_metadata.json'}
    return {'ok': True, 'metadata': json.loads(meta_path.read_text())}


def _metadata_dict() -> dict[str, Any]:
    payload = tool_get_task_metadata()
    metadata = payload.get('metadata') if payload.get('ok') else {}
    return metadata if isinstance(metadata, dict) else {}


def _default_paths() -> tuple[str | None, str | None]:
    metadata = _metadata_dict()
    input_path = metadata.get('input_path')
    output_path = metadata.get('output_path')
    return (input_path if isinstance(input_path, str) else None, output_path if isinstance(output_path, str) else None)


def _default_document_path(prefer_output: bool = True) -> str | None:
    input_path, output_path = _default_paths()
    if prefer_output and output_path and Path(output_path).exists():
        return output_path
    if input_path and Path(input_path).exists():
        return input_path
    if output_path:
        return output_path
    return input_path


def tool_copy_file(src: str, dest: str) -> dict[str, Any]:
    if not src or not dest:
        default_src, default_dest = _default_paths()
        src = src or default_src or ''
        dest = dest or default_dest or ''
    src_path = Path(src)
    dest_path = Path(dest)
    if not src_path.exists():
        return {'ok': False, 'error': f'source file not found: {src}'}
    _ensure_parent(dest_path)
    shutil.copy2(src_path, dest_path)
    return {'ok': True, 'copied_to': str(dest_path), 'size_bytes': dest_path.stat().st_size}


def tool_inspect_document(path: str, max_items: int = 8) -> dict[str, Any]:
    file_path = Path(path)
    if not file_path.exists():
        return {'ok': False, 'error': f'file not found: {path}'}
    suffix = file_path.suffix.lower()
    if suffix == '.docx':
        doc = Document(file_path)
        paragraphs = []
        for i, para in enumerate(doc.paragraphs[:max_items]):
            paragraphs.append({'index': i, 'style': para.style.name if para.style else '', 'text': para.text})
        tables = []
        for ti, table in enumerate(doc.tables[:max_items]):
            rows = []
            for ri, row in enumerate(table.rows[:max_items]):
                rows.append([cell.text for cell in row.cells[:max_items]])
            tables.append({'index': ti, 'rows': len(table.rows), 'cols': len(table.columns), 'sample_rows': rows})
        return {'ok': True, 'type': 'docx', 'paragraphs': paragraphs, 'tables': tables}
    if suffix in {'.xlsx', '.xlsm'}:
        wb = load_workbook(file_path)
        sheets = []
        for ws in wb.worksheets[:max_items]:
            sample = []
            for row in ws.iter_rows(min_row=1, max_row=min(ws.max_row, max_items), min_col=1, max_col=min(ws.max_column, max_items), values_only=True):
                sample.append(list(row))
            sheets.append({'title': ws.title, 'max_row': ws.max_row, 'max_col': ws.max_column, 'sample': sample})
        return {'ok': True, 'type': 'xlsx', 'sheet_order': wb.sheetnames, 'sheets': sheets}
    if suffix == '.pptx':
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides[:max_items], start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text:
                    texts.append(shape.text)
            slides.append({'slide_index': i, 'texts': texts[:max_items]})
        return {'ok': True, 'type': 'pptx', 'slide_count': len(prs.slides), 'slides': slides}
    if suffix == '.pdf':
        reader = PdfReader(str(file_path))
        pages = []
        with pdfplumber.open(str(file_path)) as pdf:
            for i, page in enumerate(pdf.pages[:max_items], start=1):
                text = (page.extract_text() or '')[:1000]
                pages.append({'page_index': i, 'text_preview': text})
        return {'ok': True, 'type': 'pdf', 'page_count': len(reader.pages), 'pages': pages}
    return {'ok': False, 'error': f'unsupported file type: {suffix}'}


def tool_word_replace_text(path: str, replacements: list[dict[str, str]]) -> dict[str, Any]:
    doc = Document(path)
    replaced = 0
    for para in doc.paragraphs:
        for item in replacements:
            old = item['old']
            new = item['new']
            if old in para.text:
                para.text = para.text.replace(old, new)
                replaced += 1
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    for item in replacements:
                        old = item['old']
                        new = item['new']
                        if old in para.text:
                            para.text = para.text.replace(old, new)
                            replaced += 1
    doc.save(path)
    return {'ok': True, 'replaced_count': replaced}


def tool_word_insert_table_column(path: str, table_index: int, insert_after_column: int, header_text: str = '', fill_text: str = '') -> dict[str, Any]:
    doc = Document(path)
    table = doc.tables[table_index]
    target_idx = insert_after_column
    for row_idx, row in enumerate(table.rows):
        ref_tc = row.cells[target_idx]._tc
        template_tc = copy.deepcopy(row.cells[min(target_idx + 1, len(row.cells) - 1)]._tc)
        for text_node in template_tc.findall('.//' + qn('w:t')):
            text_node.text = ''
        if row_idx == 0 and header_text:
            if template_tc.findall('.//' + qn('w:t')):
                template_tc.findall('.//' + qn('w:t'))[0].text = header_text
            else:
                from docx.oxml import OxmlElement
                p = OxmlElement('w:p')
                r = OxmlElement('w:r')
                t = OxmlElement('w:t')
                t.text = header_text
                r.append(t)
                p.append(r)
                template_tc.append(p)
        elif fill_text:
            from docx.oxml import OxmlElement
            p = OxmlElement('w:p')
            r = OxmlElement('w:r')
            t = OxmlElement('w:t')
            t.text = fill_text
            r.append(t)
            p.append(r)
            template_tc.append(p)
        ref_tc.addnext(template_tc)
    doc.save(path)
    return {'ok': True, 'table_index': table_index, 'new_column_position': insert_after_column + 1}


def tool_xlsx_read_sheet(path: str, sheet_name: str, min_row: int = 1, max_row: int = 10, min_col: int = 1, max_col: int = 10) -> dict[str, Any]:
    wb = load_workbook(path)
    ws = wb[sheet_name]
    values = []
    for row in ws.iter_rows(min_row=min_row, max_row=max_row, min_col=min_col, max_col=max_col, values_only=True):
        values.append(list(row))
    return {'ok': True, 'sheet_name': sheet_name, 'values': values}


def tool_xlsx_write_cells(path: str, sheet_name: str, cells: list[dict[str, Any]]) -> dict[str, Any]:
    wb = load_workbook(path)
    ws = wb[sheet_name]
    for item in cells:
        ws[item['cell']] = item.get('value')
    wb.save(path)
    return {'ok': True, 'written': len(cells)}


def tool_xlsx_set_formula(path: str, sheet_name: str, cell: str, formula: str) -> dict[str, Any]:
    wb = load_workbook(path)
    ws = wb[sheet_name]
    ws[cell] = formula
    wb.save(path)
    return {'ok': True, 'cell': cell, 'formula': formula}


def tool_xlsx_highlight_cells(path: str, sheet_name: str, cells: list[str], fill_color: str = 'FFFF00') -> dict[str, Any]:
    wb = load_workbook(path)
    ws = wb[sheet_name]
    fill = PatternFill(fill_type='solid', fgColor=fill_color)
    for cell in cells:
        ws[cell].fill = fill
    wb.save(path)
    return {'ok': True, 'highlighted': len(cells), 'fill_color': fill_color}


def tool_xlsx_delete_sheet(path: str, sheet_name: str) -> dict[str, Any]:
    wb = load_workbook(path)
    del wb[sheet_name]
    wb.save(path)
    return {'ok': True, 'deleted_sheet': sheet_name, 'remaining_sheets': wb.sheetnames}


def tool_xlsx_reorder_sheets(path: str, sheet_order: list[str]) -> dict[str, Any]:
    wb = load_workbook(path)
    mapping = {ws.title: ws for ws in wb.worksheets}
    wb._sheets = [mapping[name] for name in sheet_order]
    wb.save(path)
    return {'ok': True, 'sheet_order': wb.sheetnames}


def _slide_id_list(prs: Presentation):
    return prs.slides._sldIdLst  # type: ignore[attr-defined]


def tool_ppt_list_slides(path: str) -> dict[str, Any]:
    prs = Presentation(path)
    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                texts.append(shape.text)
        slides.append({'slide_index': idx, 'texts': texts})
    return {'ok': True, 'slide_count': len(prs.slides), 'slides': slides}


def tool_ppt_replace_text(path: str, replacements: list[dict[str, str]]) -> dict[str, Any]:
    prs = Presentation(path)
    replaced = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for item in replacements:
                            if item['old'] in run.text:
                                run.text = run.text.replace(item['old'], item['new'])
                                replaced += 1
    prs.save(path)
    return {'ok': True, 'replaced_count': replaced}


def tool_ppt_reorder_slides(path: str, order: list[int]) -> dict[str, Any]:
    prs = Presentation(path)
    sldIdLst = _slide_id_list(prs)
    ids = list(sldIdLst)
    reordered = [ids[i - 1] for i in order]
    for el in list(sldIdLst):
        sldIdLst.remove(el)
    for el in reordered:
        sldIdLst.append(el)
    prs.save(path)
    return {'ok': True, 'order': order}


def tool_ppt_delete_slides(path: str, slide_indices: list[int]) -> dict[str, Any]:
    prs = Presentation(path)
    sldIdLst = _slide_id_list(prs)
    for idx in sorted(slide_indices, reverse=True):
        del sldIdLst[idx - 1]
    prs.save(path)
    return {'ok': True, 'deleted': slide_indices, 'remaining': len(prs.slides)}


def tool_ppt_set_bullets(path: str, slide_index: int, bullets: list[str], shape_index: int | None = None) -> dict[str, Any]:
    prs = Presentation(path)
    slide = prs.slides[slide_index - 1]
    target = None
    text_shapes = [shape for shape in slide.shapes if hasattr(shape, 'text_frame') and shape.has_text_frame]
    if shape_index is not None:
        target = text_shapes[shape_index]
    else:
        target = text_shapes[-1] if text_shapes else None
    if target is None:
        return {'ok': False, 'error': 'no text shape found on target slide'}
    tf = target.text_frame
    tf.clear()
    first = True
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
    prs.save(path)
    return {'ok': True, 'slide_index': slide_index, 'bullet_count': len(bullets)}


def tool_pdf_extract_text(path: str, page_numbers: list[int] | None = None, max_chars: int = 1500) -> dict[str, Any]:
    page_numbers = page_numbers or []
    with pdfplumber.open(path) as pdf:
        selected = pdf.pages if not page_numbers else [pdf.pages[i - 1] for i in page_numbers]
        pages = []
        for i, page in enumerate(selected, start=1):
            text = (page.extract_text() or '')[:max_chars]
            pages.append({'page_index': (page_numbers[i - 1] if page_numbers else i), 'text': text})
    return {'ok': True, 'pages': pages}


def tool_pdf_reorder_pages(path: str, order: list[int]) -> dict[str, Any]:
    reader = PdfReader(path)
    writer = PdfWriter()
    for idx in order:
        writer.add_page(reader.pages[idx - 1])
    with open(path, 'wb') as f:
        writer.write(f)
    return {'ok': True, 'page_order': order}


def tool_pdf_delete_pages(path: str, page_numbers: list[int]) -> dict[str, Any]:
    reader = PdfReader(path)
    writer = PdfWriter()
    delete = set(page_numbers)
    for idx, page in enumerate(reader.pages, start=1):
        if idx not in delete:
            writer.add_page(page)
    with open(path, 'wb') as f:
        writer.write(f)
    return {'ok': True, 'deleted_pages': sorted(page_numbers), 'remaining_pages': len(writer.pages)}




def _normalize_args(tool_name: str, kwargs: dict[str, Any]) -> dict[str, Any]:
    kwargs = dict(kwargs)
    kwargs.pop('reason', None)
    kwargs.pop('notes', None)

    if tool_name == 'copy_file':
        if 'src' not in kwargs:
            kwargs['src'] = kwargs.get('source_path') or kwargs.get('source')
        if 'dest' not in kwargs:
            kwargs['dest'] = kwargs.get('destination_path') or kwargs.get('dst') or kwargs.get('output_path')
        default_src, default_dest = _default_paths()
        kwargs['src'] = kwargs.get('src') or default_src
        kwargs['dest'] = kwargs.get('dest') or default_dest
    elif tool_name == 'word_insert_table_column':
        if 'insert_after_column' not in kwargs:
            kwargs['insert_after_column'] = kwargs.get('after_col_index') or kwargs.get('column_index')
        if 'header_text' not in kwargs:
            kwargs['header_text'] = kwargs.get('header') or kwargs.get('column_name') or ''
        if 'fill_text' not in kwargs:
            kwargs['fill_text'] = kwargs.get('fill_value') or kwargs.get('default_value') or ''

    inspect_tools = {'inspect_document', 'xlsx_read_sheet', 'ppt_list_slides', 'pdf_extract_text'}
    mutating_tools = {
        'word_replace_text',
        'word_insert_table_column',
        'xlsx_write_cells',
        'xlsx_set_formula',
        'xlsx_highlight_cells',
        'xlsx_delete_sheet',
        'xlsx_reorder_sheets',
        'ppt_replace_text',
        'ppt_reorder_slides',
        'ppt_delete_slides',
        'ppt_set_bullets',
        'pdf_reorder_pages',
        'pdf_delete_pages',
    }
    if tool_name in inspect_tools and not kwargs.get('path'):
        kwargs['path'] = _default_document_path(prefer_output=True)
    if tool_name in mutating_tools and not kwargs.get('path'):
        kwargs['path'] = _default_document_path(prefer_output=True)
    return kwargs

TOOLS = {
    'get_task_metadata': tool_get_task_metadata,
    'copy_file': tool_copy_file,
    'inspect_document': tool_inspect_document,
    'word_replace_text': tool_word_replace_text,
    'word_insert_table_column': tool_word_insert_table_column,
    'xlsx_read_sheet': tool_xlsx_read_sheet,
    'xlsx_write_cells': tool_xlsx_write_cells,
    'xlsx_set_formula': tool_xlsx_set_formula,
    'xlsx_highlight_cells': tool_xlsx_highlight_cells,
    'xlsx_delete_sheet': tool_xlsx_delete_sheet,
    'xlsx_reorder_sheets': tool_xlsx_reorder_sheets,
    'ppt_list_slides': tool_ppt_list_slides,
    'ppt_replace_text': tool_ppt_replace_text,
    'ppt_reorder_slides': tool_ppt_reorder_slides,
    'ppt_delete_slides': tool_ppt_delete_slides,
    'ppt_set_bullets': tool_ppt_set_bullets,
    'pdf_extract_text': tool_pdf_extract_text,
    'pdf_reorder_pages': tool_pdf_reorder_pages,
    'pdf_delete_pages': tool_pdf_delete_pages,
}


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument('--tool', required=True)
    parser.add_argument('--args-json', default='{}')
    args = parser.parse_args()

    kwargs = _normalize_args(args.tool, json.loads(args.args_json))
    fn = TOOLS.get(args.tool)
    if fn is None:
        print(json.dumps({'ok': False, 'error': f'unknown tool: {args.tool}'}))
        raise SystemExit(1)
    try:
        result = fn(**kwargs)
    except Exception as exc:
        print(json.dumps({'ok': False, 'error': str(exc)}))
        raise
    print(json.dumps(result, ensure_ascii=False))


if __name__ == '__main__':
    main()
