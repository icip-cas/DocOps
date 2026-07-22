import json
import os
import re
import sys
from pathlib import Path
from datetime import date, datetime

from openpyxl import load_workbook
from openpyxl.utils.cell import get_column_letter, range_boundaries
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

sys.path.insert(0, str(Path(__file__).parent))
from verifier_utils import *  # noqa: F401,F403

META_PATH = Path(os.environ.get('TASK_METADATA_PATH', '/tests/task_metadata.json'))
META = json.loads(META_PATH.read_text(encoding='utf-8'))
EXPECT = META['verifier_expectations']


def _vopt_norm_text(value):
    if value is None:
        return ''
    if isinstance(value, datetime):
        if value.hour or value.minute or value.second:
            return value.strftime('%Y-%m-%d %H:%M')
        return value.strftime('%Y-%m-%d')
    if isinstance(value, date):
        return value.strftime('%Y-%m-%d')
    if isinstance(value, float) and value.is_integer():
        value = int(value)
    text = str(value)
    text = text.replace('\u2013', '-').replace('\u2014', '-').replace('\u2212', '-')
    text = text.replace('\xa0', ' ')
    text = re.sub(r'\s+', ' ', text).strip()
    return text


def _vopt_norm_formula(value):
    return re.sub(r'\s+', '', _vopt_norm_text(value)).upper()


def _vopt_assert_formula(actual, expected, label='formula'):
    assert _vopt_norm_formula(actual) == _vopt_norm_formula(expected), f"{label}: expected {expected!r}, found {actual!r}"


def _vopt_norm_rows(rows):
    return [[_vopt_norm_text(cell) for cell in row] for row in rows]


def _vopt_rows_equal(actual, expected):
    return _vopt_norm_rows(actual) == _vopt_norm_rows(expected)


def _vopt_rows_text(rows):
    return '\n'.join('|'.join(_vopt_norm_text(cell) for cell in row) for row in rows)


def _vopt_header_has(header, required):
    actual = {_vopt_norm_text(cell).lower() for cell in header}
    missing = [cell for cell in required if _vopt_norm_text(cell).lower() not in actual]
    assert not missing, f"Missing expected header columns {missing!r}; actual={header!r}"


def _vopt_section_text(section_part):
    return '\n'.join(p.text for p in getattr(section_part, 'paragraphs', []) if p.text is not None)


def _vopt_header_text(doc):
    return '\n'.join(_vopt_section_text(section.header) for section in doc.sections)


def _vopt_footer_text(doc):
    return '\n'.join(_vopt_section_text(section.footer) for section in doc.sections)


def _vopt_table_ref(ws, expected_name=None):
    tables = ws.tables
    if expected_name:
        for name in tables.keys():
            if str(name).lower() == str(expected_name).lower():
                return tables[name].ref
    vals = list(tables.values())
    assert vals, f"No Excel native table found on sheet {ws.title!r}"
    return vals[0].ref


def _vopt_clean_area(value):
    text = _vopt_norm_text(value).replace("'", '').replace('$', '').replace(' ', '')
    if '!' in text:
        text = text.split('!', 1)[1]
    return text.upper()


def _vopt_assert_print_area(ws, expected):
    actual = _vopt_clean_area(ws.print_area)
    target = _vopt_clean_area(expected)
    assert target in actual or actual in target, f"{ws.title}: expected print area {expected!r}, found {ws.print_area!r}"


def _vopt_freeze_pane(value):
    return getattr(value, 'coordinate', value)


def _vopt_validated_cells(ws):
    cells = set()
    for dv in ws.data_validations.dataValidation:
        for rng in dv.cells.ranges:
            text = str(rng).replace('$', '')
            try:
                min_col, min_row, max_col, max_row = range_boundaries(text)
            except ValueError:
                cells.add(text)
                continue
            if max_row > 10000:
                max_row = max(min_row, 200)
            for row in range(min_row, max_row + 1):
                for col in range(min_col, max_col + 1):
                    cells.add(f"{get_column_letter(col)}{row}")
    return cells


class _VoptValidationRanges(list):
    def __init__(self, ws, ranges):
        super().__init__(ranges)
        self.ws = ws

    def __contains__(self, expected):
        expected = str(expected).replace('$', '')
        if super().__contains__(expected):
            return True
        try:
            min_col, min_row, max_col, max_row = range_boundaries(expected)
        except ValueError:
            return super().__contains__(expected)
        target = {f"{get_column_letter(col)}{row}" for row in range(min_row, max_row + 1) for col in range(min_col, max_col + 1)}
        return target.issubset(_vopt_validated_cells(self.ws))


def _vopt_dv_ranges(ws):
    ranges = []
    for dv in ws.data_validations.dataValidation:
        ranges.extend(str(rng).replace('$', '') for rng in dv.cells.ranges)
    return _VoptValidationRanges(ws, ranges)


def _vopt_ordered_subset(expected, actual):
    actual_iter = iter([_vopt_norm_text(item) for item in actual])
    for item in [_vopt_norm_text(value) for value in expected]:
        for candidate in actual_iter:
            if item == candidate:
                break
        else:
            return False
    return True



def _resolve_output(kind):
    if kind == 'deck':
        return Path(os.environ.get('PPT_OUTPUT_PATH', EXPECT['deck_output']))
    if kind == 'tracker':
        return Path(os.environ.get('XLSX_OUTPUT_PATH', EXPECT['tracker_output']))
    raise KeyError(kind)


def _slide_title(slide):
    for shape in slide.shapes:
        if getattr(shape, 'is_placeholder', False) and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
            text = ' '.join(shape.text.split())
            if text and normalize_text(text) not in {'draft', 'final'}:
                return text
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            text = ' '.join(shape.text.split())
            if text and normalize_text(text) not in {'draft', 'final'} and not normalize_text(text).startswith('status:'):
                return text
    return ''


def _slide_titles(prs):
    return [_slide_title(slide) for slide in prs.slides]


def _slide_by_title(prs, title):
    for slide in prs.slides:
        if normalize_text(title) in normalize_text(_slide_title(slide)):
            return slide
    raise AssertionError(f'Slide not found: {title}')


def _slide_by_keywords(prs, *keywords):
    targets = [normalize_text(keyword) for keyword in keywords]
    for slide in prs.slides:
        text = normalize_text(_slide_text(slide))
        if all(keyword in text for keyword in targets):
            return slide
    raise AssertionError(f"Slide not found with keywords: {keywords!r}")


def _slide_text(slide):
    parts = []
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            text = ' '.join(shape.text.split())
            if text:
                parts.append(text)
    return '\n'.join(parts)


def _deck_text(prs):
    return '\n'.join(_slide_text(slide) for slide in prs.slides)


def _first_table(slide):
    for shape in slide.shapes:
        if getattr(shape, 'has_table', False):
            return shape.table
    raise AssertionError('Expected a native slide table')


def _table_rows(table):
    return [[cell.text.strip() for cell in row.cells] for row in table.rows]


def _nonempty_rows(rows):
    return [row for row in rows if any(_vopt_norm_text(cell) for cell in row)]


def _row_values(ws, row, max_col):
    return [ws.cell(row, col).value for col in range(1, max_col + 1)]


def _headers(ws):
    return [_vopt_norm_text(cell.value) for cell in ws[1]]


def _header_map(ws):
    return {name.lower(): idx + 1 for idx, name in enumerate(_headers(ws)) if name}


def _find_col(mapping, *candidates):
    for candidate in candidates:
        key = candidate.lower()
        if key in mapping:
            return mapping[key]
    for key, idx in mapping.items():
        if any(candidate.lower() in key for candidate in candidates):
            return idx
    raise AssertionError(f'Missing expected column among {candidates!r}; actual={list(mapping)!r}')


def _records_by_id(ws, id_col='CAPA ID'):
    mapping = _header_map(ws)
    id_idx = _find_col(mapping, id_col)
    records = {}
    for row in range(2, ws.max_row + 1):
        item_id = _vopt_norm_text(ws.cell(row, id_idx).value)
        if not item_id:
            continue
        records[item_id] = {
            name: ws.cell(row, idx).value for name, idx in mapping.items()
        }
        records[item_id]['_row'] = row
    return records


def _assert_formula_tokens(value, label, *tokens):
    norm = _vopt_norm_formula(value)
    assert norm.startswith('='), f'{label}: expected live formula, found {value!r}'
    for token in tokens:
        assert _vopt_norm_formula(token) in norm, f'{label}: expected formula to reference {token!r}, found {value!r}'


def _get(record, *names):
    lowered = {key.lower(): value for key, value in record.items()}
    for name in names:
        if name.lower() in lowered:
            return lowered[name.lower()]
    for key, value in lowered.items():
        if any(name.lower() in key for name in names):
            return value
    raise AssertionError(f'Missing field {names!r} in record keys={list(record)!r}')


def _get_exact(record, *names):
    lowered = {key.lower(): value for key, value in record.items()}
    for name in names:
        if name.lower() in lowered:
            return lowered[name.lower()]
    raise AssertionError(f'Missing exact field {names!r} in record keys={list(record)!r}')


def _maybe_get(record, *names):
    try:
        return _get(record, *names)
    except AssertionError:
        return ''


def _assert_contains(actual, expected, label):
    actual_text = _vopt_norm_text(actual).lower()
    expected_text = _vopt_norm_text(expected).lower()
    assert expected_text in actual_text, f'{label}: expected to contain {expected!r}, found {actual!r}'


def _assert_tracker_records(ws):
    records = _records_by_id(ws)
    assert set(records) == {'CAPA-17', 'CAPA-18', 'CAPA-21', 'CAPA-24'}, f'Unexpected CAPA rows: {set(records)!r}'
    expected = {
        'CAPA-17': ('D-203', 'ICH E6(R2) 5.1', 'Priya Shah', '2026-08-14', 'E-17A; E-17B', 'At Risk', 'Yes'),
        'CAPA-18': ('D-219', '21 CFR Part 11.10(e)', 'Mateo Cruz', '2026-09-01', 'E-18A', 'Ready', 'Yes'),
        'CAPA-21': ('D-226', 'ICH Q9(R1) 5.4', 'Noah Patel', '2026-07-20', 'E-21A', 'Closed', 'No'),
        'CAPA-24': ('D-241', 'ICH E6(R2) 4.9', 'Lena Ortiz', '2026-08-05', 'E-24A', 'Blocked', 'Yes'),
    }
    for capa_id, (dev, clause, owner, due, evidence, status, publish) in expected.items():
        rec = records[capa_id]
        assert _vopt_norm_text(_get(rec, 'Deviation')) == dev, f'{capa_id}: wrong deviation'
        _assert_contains(_get(rec, 'Clause'), clause, f'{capa_id}: wrong clause')
        assert _vopt_norm_text(_get(rec, 'Owner')) == owner, f'{capa_id}: wrong owner'
        assert _vopt_norm_text(_get(rec, 'Due Date')) == due, f'{capa_id}: wrong due date'
        try:
            tracker_evidence = _get_exact(rec, 'Evidence IDs')
        except AssertionError:
            tracker_evidence = ''
        if tracker_evidence:
            assert evidence in _vopt_norm_text(tracker_evidence), f'{capa_id}: wrong evidence'
        assert _vopt_norm_text(_get(rec, 'Public Status')) == status, f'{capa_id}: wrong status'
        publish_value = _vopt_norm_text(_maybe_get(rec, 'Publish', 'Publish Scope', 'Public Boundary', 'Deck Inclusion'))
        if publish == 'Yes':
            assert publish_value in {'Yes', 'Include'} or publish_value.lower().startswith('publish'), f'{capa_id}: wrong publish flag'
        else:
            assert publish_value in {'No', 'Exclude'} or 'internal' in publish_value.lower(), f'{capa_id}: wrong publish flag'
    return records


def _assert_evidence_records(ws):
    records = _records_by_id(ws, 'Evidence ID')
    expected = {
        'E-17A': ('CAPA-17', 'Retraining roster', 'Yes', 'Verified'),
        'E-17B': ('CAPA-17', 'Monitoring checklist revision', 'Yes', 'Verified'),
        'E-18A': ('CAPA-18', 'eTMF upload audit trail', 'Yes', 'Verified'),
        'E-21A': ('CAPA-21', 'Unblinded lot investigation notes', 'No', 'Verified'),
        'E-24A': ('CAPA-24', 'Lab certificate replacement request', 'Yes', 'Missing'),
    }
    for evidence_id, (capa_id, evidence, public, verification) in expected.items():
        rec = records[evidence_id]
        assert _vopt_norm_text(_get(rec, 'CAPA ID')) == capa_id, f'{evidence_id}: wrong CAPA ID'
        assert _vopt_norm_text(_get_exact(rec, 'Evidence', 'Evidence Name', 'Evidence Summary')) == evidence, f'{evidence_id}: wrong evidence'
        assert _vopt_norm_text(_get(rec, 'Public')) == public, f'{evidence_id}: wrong public flag'
        assert _vopt_norm_text(_get(rec, 'Verification Status')) == verification, f'{evidence_id}: wrong verification status'
    return records


def _assert_live_formula(value, label):
    assert _vopt_norm_formula(value).startswith('='), f'{label}: expected live formula, found {value!r}'


def _deck_matrix_rows(prs):
    rows = _nonempty_rows(_table_rows(_first_table(_slide_by_title(prs, 'Public CAPA Matrix'))))
    assert rows, 'Public CAPA Matrix table is empty'
    header = [_vopt_norm_text(cell).lower() for cell in rows[0]]
    required = ['capa id', 'deviation', 'clause', 'owner', 'due date', 'evidence']
    missing = [col for col in required if col not in header]
    assert not missing, f'Public CAPA Matrix should preserve template columns; missing={missing!r}, header={rows[0]!r}'
    status_present = 'public status' in header or 'status' in header
    assert status_present, f'Public CAPA Matrix missing status column; header={rows[0]!r}'
    data_rows = rows[1:]
    by_id = {_vopt_norm_text(row[0]): row for row in data_rows if row}
    assert set(by_id) == set(EXPECT['public_capa_ids']), f'Unexpected public CAPA IDs in matrix: {set(by_id)!r}'
    for expected in EXPECT['public_capa_rows']:
        capa_id, deviation, clause, owner, due, evidence, status = expected
        row = by_id[capa_id]
        row_text = [_vopt_norm_text(cell) for cell in row]
        require_all(' | '.join(row_text), [deviation, clause, owner, due, evidence, status], f'public matrix row {capa_id}')
    return [row[0] for row in data_rows if row]


def _workbook_public_sets(ws):
    records = _records_by_id(ws)
    public_ids = []
    internal_ids = []
    for capa_id in ['CAPA-17', 'CAPA-18', 'CAPA-21', 'CAPA-24']:
        rec = records[capa_id]
        publish_value = _vopt_norm_text(_maybe_get(rec, 'Publish', 'Publish Scope', 'Public Boundary', 'Deck Inclusion')).lower()
        if publish_value in {'yes', 'include'} or publish_value.startswith('publish'):
            public_ids.append(capa_id)
        elif publish_value in {'no', 'exclude'} or 'internal' in publish_value:
            internal_ids.append(capa_id)
        else:
            raise AssertionError(f'{capa_id}: cannot derive public/internal boundary from {publish_value!r}')
    return public_ids, internal_ids


def _row_highlighted(ws, row):
    for col in range(1, ws.max_column + 1):
        rgb = cell_fill_rgb(ws.cell(row, col))
        if rgb and rgb not in ('FFFFFF', '000000', '00000000'):
            return True
    return False


def test_outputs_exist():
    deck = _resolve_output('deck')
    tracker = _resolve_output('tracker')
    assert deck.exists(), f'Missing deck output: {deck}'
    assert tracker.exists(), f'Missing tracker output: {tracker}'
    assert deck.suffix.lower() == '.pptx'
    assert tracker.suffix.lower() == '.xlsx'
    assert deck.stat().st_size > 0
    assert tracker.stat().st_size > 0


def test_deck_order_public_boundary_and_tables():
    prs = Presentation(_resolve_output('deck'))
    titles = _slide_titles(prs)
    for title in EXPECT['absent_slide_titles']:
        assert title not in titles, f'Obsolete/internal slide still present: {title}'
    full_text = _deck_text(prs)
    require_all(full_text, ['CAPA-17', 'CAPA-18', 'CAPA-24', 'CAPA-24', 'Blocked', 'CAPA-17', 'At Risk'], 'deck')
    forbid_any(full_text, EXPECT['deck_forbidden'], 'deck')

    _deck_matrix_rows(prs)
    reference_text = _slide_text(_slide_by_keywords(prs, 'reference', 'matrix template'))
    require_all(reference_text, ['LOCKED TEMPLATE', 'public-status vocabulary'], 'reference slide')


def test_tracker_workbook_structure_values_and_controls():
    wb = load_workbook(_resolve_output('tracker'), data_only=False)
    for sheet in EXPECT['tracker_sheet_order']:
        assert sheet in wb.sheetnames, f'Missing required sheet: {sheet}'
    ws = wb['CAPA Tracker']
    _assert_tracker_records(ws)
    assert _vopt_table_ref(ws, 'capa_audit_tracker') == EXPECT['tracker_table_ref']

    ev = wb['Evidence Readiness']
    _assert_evidence_records(ev)
    assert _vopt_table_ref(ev, 'evidence_readiness') == EXPECT['evidence_table_ref']

    for row in range(2, 6):
        _assert_live_formula(ws[f'J{row}'].value, f'CAPA Tracker!J{row}')
        _assert_live_formula(ws[f'K{row}'].value, f'CAPA Tracker!K{row}')
    for row in range(2, 7):
        _assert_live_formula(ev[f'F{row}'].value, f'Evidence Readiness!F{row}')
    for sheet in EXPECT['hidden_sheets']:
        assert wb[sheet].sheet_state in ('hidden', 'veryHidden'), f'{sheet} should be hidden'
    names = {dn.name for dn in wb.defined_names.values()}
    for name in EXPECT['defined_names']:
        assert name in names, f'Missing defined name: {name}'
    actual_ranges = _vopt_dv_ranges(ws)
    for expected in EXPECT['data_validation_ranges']['CAPA Tracker']:
        assert expected in actual_ranges, f'Missing data validation range: {expected}'
    for row in EXPECT['highlight_rows']:
        assert _row_highlighted(ws, row), f'Expected highlighted tracker row {row}'


def test_cross_output_public_capa_consistency():
    prs = Presentation(_resolve_output('deck'))
    deck_ids = _deck_matrix_rows(prs)

    wb = load_workbook(_resolve_output('tracker'), data_only=False)
    ws = wb['CAPA Tracker']
    workbook_public_ids, workbook_internal_ids = _workbook_public_sets(ws)

    assert deck_ids == EXPECT['public_capa_ids'], f"Unexpected public CAPA IDs in deck: {deck_ids!r}"
    assert workbook_public_ids == EXPECT['public_capa_ids'], f"Workbook public IDs diverge: {workbook_public_ids!r}"
    assert workbook_internal_ids == EXPECT['internal_capa_ids'], f"Workbook internal IDs diverge: {workbook_internal_ids!r}"

    full_text = _deck_text(prs)
    for internal_id in EXPECT['internal_capa_ids']:
        assert internal_id not in full_text, f'{internal_id} leaked to public deck'
    for capa_id in EXPECT['public_capa_ids']:
        assert capa_id in deck_ids, f'{capa_id} missing from public deck matrix'
