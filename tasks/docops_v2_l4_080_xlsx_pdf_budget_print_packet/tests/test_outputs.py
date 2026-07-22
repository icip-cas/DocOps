import json
import os
import re
import sys
import zipfile
from pathlib import Path
from datetime import date, datetime

import pdfplumber
from docx import Document
from openpyxl import load_workbook
from openpyxl.utils.cell import get_column_letter, range_boundaries
from pptx import Presentation
from pypdf import PdfReader

sys.path.insert(0, str(Path(__file__).parent))
from verifier_utils import *  # noqa: F401,F403

META_PATH = Path(os.environ.get('TASK_METADATA_PATH', '/tests/task_metadata.json'))
META = json.loads(META_PATH.read_text(encoding='utf-8'))
EXPECT = META['verifier_expectations']


def _vopt_norm_text(value):
    if value is None:
        return ''
    if isinstance(value, datetime):
        if value.hour or value.minute or value.second:
            return value.strftime('%Y-%m-%d %H:%M')
        return value.strftime('%Y-%m-%d')
    if isinstance(value, date):
        return value.strftime('%Y-%m-%d')
    if isinstance(value, float) and value.is_integer():
        value = int(value)
    text = str(value)
    text = text.replace('\u2013', '-').replace('\u2014', '-').replace('\u2212', '-')
    text = text.replace('\xa0', ' ')
    text = re.sub(r'\s+', ' ', text).strip()
    return text


def _vopt_norm_formula(value):
    return re.sub(r'\s+', '', _vopt_norm_text(value)).upper()


def _vopt_assert_formula(actual, expected, label='formula'):
    assert _vopt_norm_formula(actual) == _vopt_norm_formula(expected), f"{label}: expected {expected!r}, found {actual!r}"


def _vopt_norm_rows(rows):
    return [[_vopt_norm_text(cell) for cell in row] for row in rows]


def _vopt_rows_equal(actual, expected):
    return _vopt_norm_rows(actual) == _vopt_norm_rows(expected)


def _vopt_rows_text(rows):
    return '\n'.join('|'.join(_vopt_norm_text(cell) for cell in row) for row in rows)


def _vopt_header_has(header, required):
    actual = {_vopt_norm_text(cell).lower() for cell in header}
    missing = [cell for cell in required if _vopt_norm_text(cell).lower() not in actual]
    assert not missing, f"Missing expected header columns {missing!r}; actual={header!r}"


def _vopt_section_text(section_part):
    return '\n'.join(p.text for p in getattr(section_part, 'paragraphs', []) if p.text is not None)


def _vopt_header_text(doc):
    return '\n'.join(_vopt_section_text(section.header) for section in doc.sections)


def _vopt_footer_text(doc):
    return '\n'.join(_vopt_section_text(section.footer) for section in doc.sections)


def _vopt_table_ref(ws, expected_name=None):
    tables = ws.tables
    if expected_name:
        for name in tables.keys():
            if str(name).lower() == str(expected_name).lower():
                return tables[name].ref
    vals = list(tables.values())
    assert vals, f"No Excel native table found on sheet {ws.title!r}"
    return vals[0].ref


def _vopt_clean_area(value):
    text = _vopt_norm_text(value).replace("'", '').replace('$', '').replace(' ', '')
    if '!' in text:
        text = text.split('!', 1)[1]
    return text.upper()


def _vopt_assert_print_area(ws, expected):
    actual = _vopt_clean_area(ws.print_area)
    target = _vopt_clean_area(expected)
    assert target in actual or actual in target, f"{ws.title}: expected print area {expected!r}, found {ws.print_area!r}"


def _vopt_freeze_pane(value):
    return getattr(value, 'coordinate', value)


def _vopt_validated_cells(ws):
    cells = set()
    for dv in ws.data_validations.dataValidation:
        for rng in dv.cells.ranges:
            text = str(rng).replace('$', '')
            try:
                min_col, min_row, max_col, max_row = range_boundaries(text)
            except ValueError:
                cells.add(text)
                continue
            if max_row > 10000:
                max_row = max(min_row, 200)
            for row in range(min_row, max_row + 1):
                for col in range(min_col, max_col + 1):
                    cells.add(f"{get_column_letter(col)}{row}")
    return cells


class _VoptValidationRanges(list):
    def __init__(self, ws, ranges):
        super().__init__(ranges)
        self.ws = ws

    def __contains__(self, expected):
        expected = str(expected).replace('$', '')
        if super().__contains__(expected):
            return True
        try:
            min_col, min_row, max_col, max_row = range_boundaries(expected)
        except ValueError:
            return super().__contains__(expected)
        target = {f"{get_column_letter(col)}{row}" for row in range(min_row, max_row + 1) for col in range(min_col, max_col + 1)}
        return target.issubset(_vopt_validated_cells(self.ws))


def _vopt_dv_ranges(ws):
    ranges = []
    for dv in ws.data_validations.dataValidation:
        ranges.extend(str(rng).replace('$', '') for rng in dv.cells.ranges)
    return _VoptValidationRanges(ws, ranges)


def _vopt_ordered_subset(expected, actual):
    actual_iter = iter([_vopt_norm_text(item) for item in actual])
    for item in [_vopt_norm_text(value) for value in expected]:
        for candidate in actual_iter:
            if item == candidate:
                break
        else:
            return False
    return True



def _out(key):
    env = {
        'document': 'DOCX_OUTPUT_PATH',
        'deck': 'PPT_OUTPUT_PATH',
        'workbook': 'XLSX_OUTPUT_PATH',
        'packet': 'PDF_OUTPUT_PATH',
    }[key]
    return Path(os.environ.get(env, EXPECT['outputs'][key]))


def _doc_text(doc):
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append('|'.join(cell.text.strip() for cell in row.cells))
    return '\n'.join(parts)


def _doc_headings(doc):
    out = []
    for p in doc.paragraphs:
        style = p.style.name if p.style else ''
        if p.text.strip() and (style.startswith('Heading') or style == 'Title'):
            out.append(p.text.strip())
    return out


def _has_toc(path):
    with zipfile.ZipFile(path) as zf:
        xml = zf.read('word/document.xml').decode('utf-8', errors='ignore')
    return bool(re.search(r'TOC\s+(?:\\)?o|TOC\s*(?:&quot;|")', xml))


def _table_rows_doc(table):
    return [[cell.text.strip() for cell in row.cells] for row in table.rows]


def _highlighted_para(doc, phrase):
    for p in doc.paragraphs:
        if phrase in p.text:
            return any(run.font.highlight_color is not None for run in p.runs) or '<w:highlight' in p._p.xml
    return False


def _slide_title(slide):
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            text = ' '.join(shape.text.split())
            if text:
                return text
    return ''


def _slide_text(slide):
    parts = []
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            text = ' '.join(shape.text.split())
            if text:
                parts.append(text)
        if getattr(shape, 'has_table', False):
            for row in shape.table.rows:
                parts.append('|'.join(cell.text.strip() for cell in row.cells))
    return '\n'.join(parts)


def _first_table(slide):
    for shape in slide.shapes:
        if getattr(shape, 'has_table', False):
            return shape.table
    raise AssertionError('Missing native PPT table')


def _ppt_table_rows(table):
    return [[cell.text.strip() for cell in row.cells] for row in table.rows]


def _pdf_texts(path):
    with pdfplumber.open(str(path)) as pdf:
        return [page.extract_text() or '' for page in pdf.pages]


def _assert_status_formula(cell, row_num):
    formula = _vopt_norm_formula(cell.value)
    assert formula.startswith('='), f'{cell.coordinate}: expected live formula, found {cell.value!r}'
    assert f'D{row_num}' in formula, f'{cell.coordinate}: formula must reference the row Status cell D{row_num}, found {cell.value!r}'
    assert any(token in formula for token in ('IF(', 'IFS(', 'SWITCH(', 'CHOOSE(')), (
        f'{cell.coordinate}: Style Check should be conditional on Status, found {cell.value!r}'
    )


def _assert_pdf_required_semantics(text):
    norm_key = re.sub(r'[^A-Z0-9]+', '', _vopt_norm_text(text).upper())
    assert 'FINAL' in norm_key and 'BUDGETPRINTPACKET' in norm_key, 'pdf: missing final budget print packet status'
    assert 'F4A261' in norm_key, 'pdf: missing accent color reference F4A261'


def _hex_color(value):
    text = str(value or '').upper()
    match = re.search(r'[0-9A-F]{6,8}', text)
    if not match:
        return None
    rgb = match.group(0)
    return rgb[-6:]


def _workbook_style_colors(wb):
    colors = set()
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                for color in (
                    getattr(cell.fill, 'fgColor', None),
                    getattr(cell.fill, 'start_color', None),
                    getattr(cell.font, 'color', None),
                ):
                    rgb = _hex_color(getattr(color, 'rgb', None))
                    if rgb:
                        colors.add(rgb)
                for side in (cell.border.left, cell.border.right, cell.border.top, cell.border.bottom):
                    rgb = _hex_color(getattr(getattr(side, 'color', None), 'rgb', None))
                    if rgb:
                        colors.add(rgb)
    return colors


def _assert_required_style_colors(colors, label):
    expected = {
        _hex_color(EXPECT.get('primary_color')),
        _hex_color(EXPECT.get('accent_color')),
    }
    expected.discard(None)
    missing = expected - set(colors)
    assert not missing, f'{label}: missing required style colors {sorted(missing)!r}; found {sorted(colors)!r}'


def test_outputs_exist():
    for key, expected in EXPECT['outputs'].items():
        path = _out(key)
        assert path.exists(), f'Missing output {key}: {path}'
        assert path.stat().st_size > 0, f'Empty output {key}: {path}'


def test_docx_style_when_present():
    if 'docx' not in EXPECT:
        return
    path = _out('document')
    doc = Document(path)
    assert _doc_headings(doc) == EXPECT['docx']['headings'], f"Unexpected headings: {_doc_headings(doc)!r}"
    text = _doc_text(doc)
    require_all(text, EXPECT['docx']['required'], 'docx')
    forbid_any(text, EXPECT['forbidden'], 'docx')
    assert _has_toc(path), 'Expected live TOC field'
    assert EXPECT['docx']['header'] in _vopt_header_text(doc)
    assert EXPECT['docx']['footer'] in _vopt_footer_text(doc)
    assert _vopt_rows_equal(_table_rows_doc(doc.tables[0])[1:], EXPECT['docx']['rows'])
    assert _highlighted_para(doc, 'Style exception resolved'), 'Expected highlighted style exception paragraph'


def test_xlsx_style_when_present():
    if 'xlsx' not in EXPECT:
        return
    wb = load_workbook(_out('workbook'), data_only=False)
    assert wb.sheetnames == EXPECT['xlsx']['sheet_order']
    _assert_required_style_colors(_workbook_style_colors(wb), 'xlsx style')
    ws = wb['Styled Register']
    rows = [[ws.cell(r, c).value for c in range(1, 5)] for r in range(2, 2 + len(EXPECT['xlsx']['rows']))]
    assert _vopt_rows_equal(rows, EXPECT['xlsx']['rows']), f"Unexpected rows: {rows!r}"
    assert _vopt_table_ref(ws, EXPECT['xlsx']['table_name']) == EXPECT['xlsx']['table_ref']
    _vopt_assert_print_area(ws, EXPECT['xlsx']['print_area'])
    assert _vopt_freeze_pane(ws.freeze_panes) == EXPECT['xlsx']['freeze_panes']
    for row_num in range(2, 2 + len(EXPECT['xlsx']['rows'])):
        _assert_status_formula(ws[f'E{row_num}'], row_num)
    names = {dn.name for dn in wb.defined_names.values()}
    for name in EXPECT['xlsx']['defined_names']:
        assert name in names, f'Missing defined name: {name}'
    for sheet in EXPECT['xlsx']['hidden_sheets']:
        assert wb[sheet].sheet_state in ('hidden', 'veryHidden'), f'{sheet} should be hidden'
    actual_ranges = _vopt_dv_ranges(ws)
    for rng in EXPECT['xlsx']['dv_ranges']:
        assert rng in actual_ranges, f'Missing data validation range: {rng}'
    assert len(ws.conditional_formatting) >= 1, 'Expected conditional formatting'


def test_pptx_style_when_present():
    if 'pptx' not in EXPECT:
        return
    prs = Presentation(_out('deck'))
    for expected_title, slide in zip(EXPECT['pptx']['titles'], prs.slides):
        assert expected_title in _slide_text(slide), f"Missing slide title {expected_title!r}"
    full = '\n'.join(_slide_text(slide) for slide in prs.slides)
    require_all(full, EXPECT['pptx']['required'], 'pptx')
    forbid_any(full, EXPECT['forbidden'], 'pptx')
    table_rows = None
    for slide in prs.slides:
        try:
            rows = _ppt_table_rows(_first_table(slide))
            if len(rows) > 1:
                table_rows = rows
                break
        except AssertionError:
            pass
    assert table_rows == [[str(c) for c in row] for row in EXPECT['pptx']['rows']], f"Unexpected PPT table rows: {table_rows!r}"


def test_pdf_style_when_present():
    if 'pdf' not in EXPECT:
        return
    path = _out('packet')
    reader = PdfReader(str(path))
    assert len(reader.pages) == len(EXPECT['pdf']['titles'])
    texts = _pdf_texts(path)
    for expected_title, page_text in zip(EXPECT['pdf']['titles'], texts):
        assert expected_title in page_text, f"Missing PDF title {expected_title!r} on its expected page"
    full = '\n'.join(texts)
    _assert_pdf_required_semantics(full)
    _assert_required_style_colors(set(pdf_theme_rect_colors(path)), 'pdf style')
    forbid_any(full, EXPECT['forbidden'], 'pdf')
    outline = [title for _level, title in flatten_outline(reader.outline)]
    assert len(outline) >= len(EXPECT['pdf']['titles']) + 1, f"Expected root bookmark plus page bookmarks, found {outline!r}"
    for title in EXPECT['pdf']['titles']:
        assert title in outline, f'Missing PDF bookmark: {title}'
