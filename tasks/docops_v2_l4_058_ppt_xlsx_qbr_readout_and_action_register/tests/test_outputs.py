import json
import os
import re
import sys
from pathlib import Path
from datetime import date, datetime

from openpyxl import load_workbook
from openpyxl.utils.cell import get_column_letter, range_boundaries
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

sys.path.insert(0, str(Path(__file__).parent))
from verifier_utils import *  # noqa: F401,F403

META_PATH = Path(os.environ.get('TASK_METADATA_PATH', '/tests/task_metadata.json'))
META = json.loads(META_PATH.read_text(encoding='utf-8'))
INPUT_PATH = Path(os.environ.get('INPUT_PATH', META['input_path']))
EXPECT = META['verifier_expectations']


def _vopt_norm_text(value):
    if value is None:
        return ''
    if isinstance(value, datetime):
        if value.hour or value.minute or value.second:
            return value.strftime('%Y-%m-%d %H:%M')
        return value.strftime('%Y-%m-%d')
    if isinstance(value, date):
        return value.strftime('%Y-%m-%d')
    if isinstance(value, float) and value.is_integer():
        value = int(value)
    text = str(value)
    text = text.replace('\u2013', '-').replace('\u2014', '-').replace('\u2212', '-')
    text = text.replace('\xa0', ' ')
    text = re.sub(r'\s+', ' ', text).strip()
    return text


def _vopt_norm_formula(value):
    return re.sub(r'\s+', '', _vopt_norm_text(value)).upper()


def _vopt_assert_formula(actual, expected, label='formula'):
    assert _vopt_norm_formula(actual) == _vopt_norm_formula(expected), f"{label}: expected {expected!r}, found {actual!r}"


def _vopt_norm_rows(rows):
    return [[_vopt_norm_text(cell) for cell in row] for row in rows]


def _vopt_rows_equal(actual, expected):
    return _vopt_norm_rows(actual) == _vopt_norm_rows(expected)


def _vopt_rows_text(rows):
    return '\n'.join('|'.join(_vopt_norm_text(cell) for cell in row) for row in rows)


def _vopt_header_has(header, required):
    actual = {_vopt_norm_text(cell).lower() for cell in header}
    missing = [cell for cell in required if _vopt_norm_text(cell).lower() not in actual]
    assert not missing, f"Missing expected header columns {missing!r}; actual={header!r}"


def _vopt_section_text(section_part):
    return '\n'.join(p.text for p in getattr(section_part, 'paragraphs', []) if p.text is not None)


def _vopt_header_text(doc):
    return '\n'.join(_vopt_section_text(section.header) for section in doc.sections)


def _vopt_footer_text(doc):
    return '\n'.join(_vopt_section_text(section.footer) for section in doc.sections)


def _vopt_table_ref(ws, expected_name=None):
    tables = ws.tables
    if expected_name:
        for name in tables.keys():
            if str(name).lower() == str(expected_name).lower():
                return tables[name].ref
    vals = list(tables.values())
    assert vals, f"No Excel native table found on sheet {ws.title!r}"
    return vals[0].ref


def _vopt_clean_area(value):
    text = _vopt_norm_text(value).replace("'", '').replace('$', '').replace(' ', '')
    if '!' in text:
        text = text.split('!', 1)[1]
    return text.upper()


def _vopt_assert_print_area(ws, expected):
    actual = _vopt_clean_area(ws.print_area)
    target = _vopt_clean_area(expected)
    assert target in actual or actual in target, f"{ws.title}: expected print area {expected!r}, found {ws.print_area!r}"


def _vopt_freeze_pane(value):
    return getattr(value, 'coordinate', value)


def _vopt_validated_cells(ws):
    cells = set()
    for dv in ws.data_validations.dataValidation:
        for rng in dv.cells.ranges:
            text = str(rng).replace('$', '')
            try:
                min_col, min_row, max_col, max_row = range_boundaries(text)
            except ValueError:
                cells.add(text)
                continue
            if max_row > 10000:
                max_row = max(min_row, 200)
            for row in range(min_row, max_row + 1):
                for col in range(min_col, max_col + 1):
                    cells.add(f"{get_column_letter(col)}{row}")
    return cells


class _VoptValidationRanges(list):
    def __init__(self, ws, ranges):
        super().__init__(ranges)
        self.ws = ws

    def __contains__(self, expected):
        expected = str(expected).replace('$', '')
        if super().__contains__(expected):
            return True
        try:
            min_col, min_row, max_col, max_row = range_boundaries(expected)
        except ValueError:
            return super().__contains__(expected)
        target = {f"{get_column_letter(col)}{row}" for row in range(min_row, max_row + 1) for col in range(min_col, max_col + 1)}
        return target.issubset(_vopt_validated_cells(self.ws))


def _vopt_dv_ranges(ws):
    ranges = []
    for dv in ws.data_validations.dataValidation:
        ranges.extend(str(rng).replace('$', '') for rng in dv.cells.ranges)
    return _VoptValidationRanges(ws, ranges)


def _vopt_ordered_subset(expected, actual):
    actual_iter = iter([_vopt_norm_text(item) for item in actual])
    for item in [_vopt_norm_text(value) for value in expected]:
        for candidate in actual_iter:
            if item == candidate:
                break
        else:
            return False
    return True



def _resolve_output(kind):
    if kind == 'deck':
        return Path(os.environ.get('PPT_OUTPUT_PATH', EXPECT['deck_output']))
    if kind == 'register':
        return Path(os.environ.get('XLSX_OUTPUT_PATH', EXPECT['register_output']))
    raise KeyError(kind)


def _slide_title(slide):
    for shape in slide.shapes:
        if getattr(shape, 'is_placeholder', False) and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
            text = ' '.join(shape.text.split())
            if text and normalize_text(text) not in {'draft', 'final'}:
                return text
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            text = ' '.join(shape.text.split())
            if text and normalize_text(text) not in {'draft', 'final'} and not normalize_text(text).startswith('status:'):
                return text
    return ''


EXPECTED_RISK = {
    row[0]: {
        'account': row[0],
        'arr': row[1],
        'renewal window': row[2],
        'renewal date': row[2],
        'date': row[2],
        'status': row[3],
        'risk': row[3],
        'owner': row[4],
    }
    for row in EXPECT['risk_rows']
}


EXPECTED_ACTION = {
    row[1]: {
        'action id': row[0],
        'id': row[0],
        'account': row[1],
        'action': row[2],
        'owner': row[3],
        'status': row[4],
        'risk': row[4],
    }
    for row in EXPECT['action_rows']
}

for account, action_expectation in EXPECTED_ACTION.items():
    if account in EXPECTED_RISK:
        EXPECTED_RISK[account]['action'] = action_expectation['action']


def _norm_key(value):
    text = normalize_text(_vopt_norm_text(value))
    aliases = {
        'base arr': 'arr',
        'revised arr': 'arr',
        'risk status': 'status',
        'status check': 'check',
        'validation': 'check',
        'expected status': 'expected status',
    }
    return aliases.get(text, text)


def _row_maps(rows):
    assert rows and len(rows) > 1, f'Expected a table with header and data rows, found {rows!r}'
    headers = [_norm_key(cell) for cell in rows[0]]
    mapped = []
    for row in rows[1:]:
        item = {}
        for idx, header in enumerate(headers):
            if header:
                item[header] = _vopt_norm_text(row[idx] if idx < len(row) else '')
        if any(item.values()):
            mapped.append(item)
    return mapped


def _find_field(row, names):
    for name in names:
        key = _norm_key(name)
        if key in row:
            return row[key]
    return ''


def _assert_expected_rows_by_header(rows, expected_by_account, required_fields, label, optional_fields=()):
    actual = {row.get('account', ''): row for row in _row_maps(rows)}
    missing = [account for account in expected_by_account if account not in actual]
    assert not missing, f'{label}: missing accounts {missing!r}; actual={actual!r}'
    for account, expected in expected_by_account.items():
        row = actual[account]
        for field in required_fields:
            actual_value = _find_field(row, [field])
            assert _vopt_norm_text(actual_value) == _vopt_norm_text(expected[field]), (
                f'{label} {account} {field}: expected {expected[field]!r}, found {actual_value!r}; row={row!r}'
            )
        for field in optional_fields:
            if _norm_key(field) in row:
                actual_value = row[_norm_key(field)]
                assert _vopt_norm_text(actual_value) == _vopt_norm_text(expected[field]), (
                    f'{label} {account} {field}: expected {expected[field]!r}, found {actual_value!r}; row={row!r}'
                )


def _workbook_rows_by_header(ws):
    rows = [[ws.cell(r, c).value for c in range(1, ws.max_column + 1)] for r in range(1, 5)]
    return _row_maps(rows)


def _assert_status_check_formula(formula, row, label):
    norm = _vopt_norm_formula(formula)
    assert norm.startswith('='), f'{label}: expected a live formula, found {formula!r}'
    assert f'E{row}' in norm, f'{label}: formula must check the row status cell E{row}, found {formula!r}'
    assert 'OK' in norm, f'{label}: formula must return OK for valid status rows, found {formula!r}'
    assert any(token in norm for token in ('CHECK', 'MISMATCH', 'ERR', 'ERROR', 'FALSE', 'INVALID')), (
        f'{label}: formula must have a non-OK branch for invalid status rows, found {formula!r}'
    )


def _assert_summary_semantics(text):
    norm = normalize_text(text)
    require_all(text, ['Northstar Bank', 'Watch', 'Zenith Foods', '$1.40M'], 'Executive Summary')
    assert 'legal' in norm or 'addendum' in norm or 'cleared' in norm, (
        'Executive Summary: must explain why Northstar changed to Watch'
    )
    forbid_any(text, ['Northstar Bank currently marked Hold', '$1.20M', 'Replace with three summary bullets'], 'Executive Summary')


def _assert_callout_semantics(text):
    forbid_any(text, ['placeholder', 'LOCKED LAYOUT'], 'Risk Callout')
    norm = normalize_text(text)
    assert 'northstar bank' in norm or 'zenith foods' in norm, 'Risk Callout: must identify a real account'
    assert '$1.40m' in norm or '$2.10m' in norm or 'watch' in norm or 'hold' in norm, (
        'Risk Callout: must include a real risk/status/ARR fact'
    )


def _slide_titles(prs):
    return [_slide_title(slide) for slide in prs.slides]


def _slide_by_title(prs, title):
    target = normalize_text(title)
    for slide in prs.slides:
        if normalize_text(_slide_title(slide)) == target:
            return slide
    for slide in prs.slides:
        if target in normalize_text(_slide_text(slide)):
            return slide
    raise AssertionError(f'Slide not found: {title}')


def _slide_text(slide):
    parts = []
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            text = ' '.join(shape.text.split())
            if text:
                parts.append(text)
    return '\n'.join(parts)


def _first_table(slide):
    for shape in slide.shapes:
        if getattr(shape, 'has_table', False):
            return shape.table
    raise AssertionError('Expected a native slide table')


def _table_rows(table):
    rows = []
    for row in table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    return rows


def test_outputs_exist():
    deck = _resolve_output('deck')
    register = _resolve_output('register')
    assert deck.exists(), f'Missing deck output: {deck}'
    assert register.exists(), f'Missing register output: {register}'
    assert deck.suffix.lower() == '.pptx'
    assert register.suffix.lower() == '.xlsx'
    assert deck.stat().st_size > 0
    assert register.stat().st_size > 0


def test_deck_content_and_order():
    prs = Presentation(_resolve_output('deck'))
    assert _vopt_ordered_subset(EXPECT['titles_order'], _slide_titles(prs)), f"Unexpected slide order: {_slide_titles(prs)!r}"
    for title in EXPECT['absent_slide_titles']:
        assert title not in _slide_titles(prs), f'Scratch/obsolete slide still present: {title}'

    summary_text = _slide_text(_slide_by_title(prs, 'Executive Summary'))
    _assert_summary_semantics(summary_text)

    risk_table = _first_table(_slide_by_title(prs, 'Renewal Risk Table'))
    _assert_expected_rows_by_header(
        _table_rows(risk_table),
        EXPECTED_RISK,
        required_fields=('account', 'arr', 'status', 'owner'),
        optional_fields=('renewal window', 'action'),
        label='Renewal Risk Table',
    )

    action_table = _first_table(_slide_by_title(prs, 'Action Register'))
    _assert_expected_rows_by_header(
        _table_rows(action_table),
        EXPECTED_ACTION,
        required_fields=('action id', 'account', 'action', 'owner', 'status'),
        label='Action Register slide',
    )

    callout_text = _slide_text(_slide_by_title(prs, 'Risk Callout'))
    _assert_callout_semantics(callout_text)

    for title in EXPECT['reference_titles']:
        slide = _slide_by_title(prs, title)
        if 'Risk' in title:
            require_all(_slide_text(slide), ['LOCKED LAYOUT'], f'{title} reference')
        else:
            assert _table_rows(_first_table(slide))[0][0] == 'Action ID', f'{title} reference table not preserved'


def test_register_workbook_structure_and_values():
    wb = load_workbook(_resolve_output('register'), data_only=False)
    for sheet in EXPECT['register_sheet_order']:
        assert sheet in wb.sheetnames, f'Missing expected sheet: {sheet}'
    ws = wb['Action Register']
    _assert_expected_rows_by_header(
        [[ws.cell(r, c).value for c in range(1, ws.max_column + 1)] for r in range(1, 5)],
        EXPECTED_ACTION,
        required_fields=('action id', 'account', 'action', 'owner', 'status'),
        label='Action Register workbook',
    )

    assert _vopt_table_ref(ws, 'renewal_action_register') == EXPECT['register_table_ref']
    for cell in EXPECT['register_formula_cells']:
        _assert_status_check_formula(ws[cell].value, int(cell[1:]), cell)

    for sheet in EXPECT['hidden_sheets']:
        assert wb[sheet].sheet_state in ('hidden', 'veryHidden'), f'{sheet} must remain hidden'

    names = {dn.name for dn in wb.defined_names.values()}
    for name in EXPECT['defined_names']:
        assert name in names, f'Missing defined name: {name}'

    actual_ranges = _vopt_dv_ranges(ws)
    for expected in EXPECT['data_validation_ranges']['Action Register']:
        assert expected in actual_ranges, f'Missing data validation range: {expected}'


def test_cross_output_consistency():
    prs = Presentation(_resolve_output('deck'))
    ppt_rows = _row_maps(_table_rows(_first_table(_slide_by_title(prs, 'Action Register'))))
    wb = load_workbook(_resolve_output('register'), data_only=False)
    ws = wb['Action Register']
    xlsx_rows = _workbook_rows_by_header(ws)
    by_account = {row.get('account', ''): row for row in xlsx_rows}
    for ppt_row in ppt_rows:
        account = ppt_row.get('account', '')
        assert account in by_account, f'PPT account missing from XLSX action register: {account!r}'
        xlsx_row = by_account[account]
        for field in ('owner', 'action', 'status'):
            assert _find_field(ppt_row, [field]) == _find_field(xlsx_row, [field]), (
                f'PPT and XLSX action register {account} {field} diverge: {ppt_row!r} vs {xlsx_row!r}'
            )
