import json
import os
import re
from datetime import date, datetime
from pathlib import Path

from openpyxl import load_workbook
from openpyxl.utils.cell import get_column_letter, range_boundaries
from pptx import Presentation

from verifier_utils import *  # noqa: F401,F403

META_PATH = Path(os.environ.get("TASK_METADATA_PATH", "/tests/task_metadata.json"))
META = json.loads(META_PATH.read_text(encoding="utf-8"))
EXPECT = META["verifier_expectations"]


def _norm(value):
    if value is None:
        return ""
    if isinstance(value, datetime):
        return value.strftime("%Y-%m-%d")
    if isinstance(value, date):
        return value.strftime("%Y-%m-%d")
    if isinstance(value, float) and value.is_integer():
        value = int(value)
    text = str(value).replace("\u2013", "-").replace("\u2014", "-").replace("\xa0", " ")
    return re.sub(r"\s+", " ", text).strip()


def _norm_formula(value):
    return re.sub(r"\s+", "", _norm(value)).upper()


def _path(kind):
    env = {"pptx": "PPTX_OUTPUT_PATH", "xlsx": "XLSX_OUTPUT_PATH"}[kind]
    key = {"pptx": "pptx_output", "xlsx": "xlsx_output"}[kind]
    return Path(os.environ.get(env, EXPECT[key]))


def _ppt_text_and_titles(path):
    prs = Presentation(path)
    all_text = []
    titles = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                for line in shape.text.splitlines():
                    line = line.strip()
                    if line:
                        texts.append(line)
                        all_text.append(line)
        titles.append(texts[0] if texts else "")
    return "\n".join(all_text), titles, len(prs.slides)


def _rows(ws, start, end, max_col):
    return [[_norm(ws.cell(r, c).value) for c in range(1, max_col + 1)] for r in range(start, end + 1)]


def _expected(rows):
    return [[_norm(c) for c in row] for row in rows]


def _table_ref(ws, expected_name):
    for name in ws.tables.keys():
        if str(name).lower() == str(expected_name).lower():
            return ws.tables[name].ref
    raise AssertionError(f"Missing table {expected_name!r} on {ws.title!r}")


def _clean_area(value):
    text = _norm(value).replace("'", "").replace("$", "").replace(" ", "")
    if "!" in text:
        text = text.split("!", 1)[1]
    return text.upper()


def _range_covered(ws, expected):
    expected = expected.replace("$", "")
    min_col, min_row, max_col, max_row = range_boundaries(expected)
    target = {f"{get_column_letter(col)}{row}" for row in range(min_row, max_row + 1) for col in range(min_col, max_col + 1)}
    cells = set()
    for dv in ws.data_validations.dataValidation:
        for rng in dv.cells.ranges:
            min_c, min_r, max_c, max_r = range_boundaries(str(rng).replace("$", ""))
            cells.update(f"{get_column_letter(col)}{row}" for row in range(min_r, max_r + 1) for col in range(min_c, max_c + 1))
    return target.issubset(cells)


def _sheet_text(wb, sheets):
    parts = []
    for sheet in sheets:
        for row in wb[sheet].iter_rows():
            for cell in row:
                if cell.value is not None:
                    parts.append(str(cell.value))
    return "\n".join(parts)


def test_outputs_exist():
    assert _path("pptx").exists()
    assert _path("xlsx").exists()
    assert _path("pptx").suffix.lower() == ".pptx"
    assert _path("xlsx").suffix.lower() == ".xlsx"


def test_pptx_structure_content_and_privacy():
    text, titles, count = _ppt_text_and_titles(_path("pptx"))
    assert count == 6
    assert titles == EXPECT["slide_titles"]
    require_all(text, EXPECT["ppt_required"], "summary deck")
    forbid_any(text, EXPECT["ppt_forbidden"], "summary deck")


def test_workbook_structure_formulas_controls():
    wb = load_workbook(_path("xlsx"), data_only=False)
    assert wb.sheetnames == EXPECT["sheet_order"]
    for sheet, (table_name, ref) in EXPECT["tables"].items():
        assert _table_ref(wb[sheet], table_name) == ref
    for sheet in EXPECT["hidden_sheets"]:
        assert wb[sheet].sheet_state in ("hidden", "veryHidden")
    names = {dn.name for dn in wb.defined_names.values()}
    for name in EXPECT["defined_names"]:
        assert name in names
    for ref, expected in EXPECT["formula_cells"].items():
        sheet, cell = ref.split("!")
        assert _norm_formula(wb[sheet][cell].value) == _norm_formula(expected), f"{ref}: expected {expected!r}, found {wb[sheet][cell].value!r}"
    for sheet, ranges in EXPECT["data_validation"].items():
        for rng in ranges:
            assert _range_covered(wb[sheet], rng), f"{sheet}: missing validation over {rng}"
    for sheet, area in EXPECT["print_areas"].items():
        actual = _clean_area(wb[sheet].print_area)
        target = _clean_area(area)
        assert target in actual or actual in target


def test_workbook_values_and_privacy():
    wb = load_workbook(_path("xlsx"), data_only=False)
    assert _rows(wb["Project Summary"], 1, 11, 4) == _expected(EXPECT["summary_rows"])
    assert _rows(wb["Measure Calculations"], 1, 5, 12) == _expected(EXPECT["measure_rows"])
    assert _rows(wb["Eligibility Review"], 1, 4, 5) == _expected(EXPECT["eligibility_rows"])
    assert _rows(wb["Excluded Measures"], 1, 3, 4) == _expected(EXPECT["excluded_rows"])
    assert _rows(wb["Submission Checklist"], 1, 6, 4) == _expected(EXPECT["checklist_rows"])
    assert _rows(wb["Raw Inputs"], 1, 6, 9) == _expected(EXPECT["raw_rows"])
    assert _rows(wb["Program Rules"], 1, 5, 2) == _expected(EXPECT["rules_rows"])
    public = _sheet_text(wb, ["Project Summary", "Measure Calculations", "Eligibility Review", "Excluded Measures", "Submission Checklist"])
    forbid_any(public, EXPECT["forbidden_public"], "public workbook sheets")


def test_cross_output_consistency():
    ppt_text, _, _ = _ppt_text_and_titles(_path("pptx"))
    wb = load_workbook(_path("xlsx"), data_only=False)
    for anchor in ["Harbor Point Medical Office", "BrightGrid Commercial Custom Rebate", "L-101", "H-202", "L-102", "2026-08-14"]:
        assert anchor in ppt_text
        assert anchor in _sheet_text(wb, ["Project Summary", "Measure Calculations", "Eligibility Review", "Submission Checklist"])
    assert wb["Measure Calculations"]["G5"].value == "=SUM(G2:G4)"
    assert wb["Measure Calculations"]["L5"].value == "=SUM(L2:L4)"
    assert wb["Project Summary"]["B11"].value == '=IF(AND(B4=3,B5=2,B8>7000),"Ready","Hold")'
    for row in range(2, 5):
        assert wb["Eligibility Review"].cell(row, 4).value == "Eligible"
    assert [wb["Excluded Measures"].cell(row, 4).value for row in range(2, 4)] == ["Exclude", "Exclude"]
    for forbidden in EXPECT["forbidden_public"]:
        assert normalize_text(forbidden) not in normalize_text(ppt_text)
