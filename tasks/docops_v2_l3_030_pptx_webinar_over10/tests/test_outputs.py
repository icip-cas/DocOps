
import json, sys
from pathlib import Path
from pptx import Presentation
sys.path.insert(0, str(Path(__file__).parent))
from verifier_utils import *

META_PATH=Path(__import__('os').environ.get('TASK_METADATA_PATH','/tests/task_metadata.json'))
if not META_PATH.exists():
    META_PATH=Path(__file__).parent/'task_metadata.json'
META=json.loads(META_PATH.read_text())
INPUT_PATH=Path(META['input_path']); OUTPUT_PATH=Path(META['output_path'])


def slide(prs,title):
    for s in prs.slides:
        if any(getattr(x,'has_text_frame',False) and x.text.strip().splitlines()[0].strip()==title for x in s.shapes):
            return s
    raise AssertionError(title)

def text(s):
    return '\n'.join(shape.text for shape in s.shapes if getattr(shape,'has_text_frame',False))

def deck_text(prs):
    parts=[]
    for s in prs.slides:
        parts.append(text(s))
        parts.extend(ppt_notes_lines(s))
    return normalize_text('\n'.join(parts))

def assert_no_internal_draft_leak(prs):
    all_text=deck_text(prs)
    for forbidden in ['Internal Only Draft','Draft Mutation Leak','Internal Only Draft Mutation Leak']:
        assert normalize_text(forbidden) not in all_text, f'Forbidden internal/draft leak remained: {forbidden}'

def table_matrix(table):
    return [[table.cell(r,c).text.strip() for c in range(len(table.columns))] for r in range(len(table.rows))]

def test_webinar_over10():
    run_preflight(META, INPUT_PATH, OUTPUT_PATH)
    prs=Presentation(OUTPUT_PATH); inp=Presentation(INPUT_PATH)
    assert_ppt_titles_equal(prs, ['Title','Agenda','Setup','Scenario','Deep Dive','Quiz','Actions','Reference - Legal'])
    assert 'Scratch Backup' not in ppt_slide_titles(prs) and 'Deprecated Poll' not in ppt_slide_titles(prs)
    assert all(not t.startswith('Draft:') for t in ppt_slide_titles(prs))
    assert_no_internal_draft_leak(prs)
    agenda=ppt_find_slide(prs,'Agenda')
    tables=[sh.table for sh in agenda.shapes if getattr(sh,'has_table',False)]
    assert tables
    assert table_matrix(tables[0])==[
        ['Section','Slide'],
        ['Setup','3'],
        ['Scenario','4'],
        ['Deep Dive','5'],
        ['Quiz','6'],
        ['Actions','7'],
    ]
    quiz=ppt_find_slide(prs,'Quiz')
    qtables=[sh.table for sh in quiz.shapes if getattr(sh,'has_table',False)]
    assert qtables
    assert table_matrix(qtables[0])==[
        ['Question','Required Action'],
        ['How should participants respond to the live check-in?','Open poll'],
    ]
    for title in ['Setup','Scenario','Deep Dive','Quiz','Actions']:
        assert 'Section: Delivery' in text(ppt_find_slide(prs,title))
    for idx,title in enumerate(['Title','Agenda','Setup','Scenario','Deep Dive','Quiz','Actions'], start=1):
        assert f'Slide {idx} of 8' in text(ppt_find_slide(prs,title))
    assert ppt_notes_lines(ppt_find_slide(prs,'Actions'))==['Facilitator note: send recap within 24 hours.']
    ref=ppt_find_slide(prs,'Reference - Legal'); inp_ref=ppt_find_slide(inp,'Reference - Legal')
    assert text(ref)==text(inp_ref)
    assert ppt_notes_lines(ref)==ppt_notes_lines(inp_ref)
    assert len(ref.shapes)==len(inp_ref.shapes)
