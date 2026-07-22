import json
import sys
from pathlib import Path

from pptx import Presentation

sys.path.insert(0, str(Path(__file__).parent))
from verifier_utils import *  # noqa: F401,F403

META_PATH = Path(__import__('os').environ.get('TASK_METADATA_PATH', '/tests/task_metadata.json'))
if not META_PATH.exists():
    META_PATH = Path(__file__).parent / 'task_metadata.json'
META = json.loads(META_PATH.read_text())
INPUT_PATH = Path(META['input_path'])
OUTPUT_PATH = Path(META['output_path'])


def slide_text(slide):
    return '\n'.join(
        ' '.join(shape.text.split())
        for shape in slide.shapes
        if getattr(shape, 'has_text_frame', False) and shape.text.strip()
    )


def deck_text(prs):
    parts = []
    for slide in prs.slides:
        parts.append(slide_text(slide))
        parts.extend(ppt_notes_lines(slide))
    return normalize_text('\n'.join(parts))


def assert_no_internal_draft_leak(prs):
    text = deck_text(prs)
    for forbidden in [
        'Internal Only Draft',
        'Draft Mutation Leak',
        'Internal Only Draft Mutation Leak',
    ]:
        assert normalize_text(forbidden) not in text, f'Forbidden internal/draft leak remained: {forbidden}'


def table_on(slide, index=0):
    tables = [shape.table for shape in slide.shapes if getattr(shape, 'has_table', False)]
    assert len(tables) > index, f'Missing table {index} on slide'
    return tables[index]


def table_matrix(table):
    return [
        [norm_cell(table.cell(r, c).text) for c in range(len(table.columns))]
        for r in range(len(table.rows))
    ]


def slide_title(slide):
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False) and shape.text.strip():
            return ' '.join(shape.text.strip().splitlines()[0].split())
    return ''


def slide_by_exact_title(prs, title):
    for slide in prs.slides:
        if normalize_text(slide_title(slide)) == normalize_text(title):
            return slide
    raise AssertionError(f'Slide not found by exact title: {title}')


def test_qbr_release_deck_longflow():
    run_preflight(META, INPUT_PATH, OUTPUT_PATH)
    prs = Presentation(OUTPUT_PATH)
    src = Presentation(INPUT_PATH)

    expected_titles = [
        'Title',
        'Agenda',
        'Market Context',
        'KPI Snapshot',
        'Customer Risks',
        'Mitigation Plan',
        'Decision Request',
        'Appendix - Data Notes',
        'Reference - Locked Legal',
    ]
    assert_ppt_titles_equal(prs, expected_titles)
    for forbidden in ['Draft:', 'Scratch Backup', 'Deprecated Roadmap', 'Working Copy']:
        assert normalize_text(forbidden) not in normalize_text('\n'.join(ppt_slide_titles(prs)))
    assert_no_internal_draft_leak(prs)
    assert len(prs.slides) == 9

    agenda = table_matrix(table_on(slide_by_exact_title(prs, 'Agenda')))
    assert agenda == [
        ['Section', 'Slide', 'Owner'],
        ['Market Context', '3', 'Maya Chen'],
        ['KPI Snapshot', '4', 'Jon Bell'],
        ['Customer Risks', '5', 'Priya Rao'],
        ['Mitigation Plan', '6', 'Maya Chen'],
        ['Decision Request', '7', 'Jon Bell'],
        ['Appendix - Data Notes', '8', 'Priya Rao'],
    ]

    for idx, title in enumerate(expected_titles, start=1):
        if title in ['Title', 'Agenda', 'Reference - Locked Legal']:
            continue
        text = slide_text(slide_by_exact_title(prs, title))
        assert f'Slide {idx} of 9' in text, f'Missing slide number on {title}'
        assert 'QBR Release | 2026-06-05' in text, f'Missing release footer on {title}'

    kpi = table_matrix(table_on(slide_by_exact_title(prs, 'KPI Snapshot')))
    assert kpi == [
        ['Metric', 'Value', 'QBR Readout'],
        ['ARR', '3.2M', 'Up 10% QoQ'],
        ['Churn risk', '4 accounts', 'Requires executive review'],
        ['NRR', '94%', 'Below target'],
        ['Expansion pipeline', '1.1M', 'Two deals require support plan'],
    ]

    risks = table_matrix(table_on(slide_by_exact_title(prs, 'Customer Risks')))
    assert risks == [
        ['Risk', 'Account Count', 'Owner', 'Mitigation'],
        ['Renewal delay', '2', 'Maya Chen', 'Executive outreach by 2026-06-10'],
        ['Adoption gap', '1', 'Priya Rao', 'Enablement plan by 2026-06-12'],
        ['Support escalation', '1', 'Jon Bell', 'Daily incident review'],
    ]

    actions = table_matrix(table_on(slide_by_exact_title(prs, 'Mitigation Plan')))
    assert actions == [
        ['Action', 'Owner', 'Due Date', 'Status'],
        ['Confirm renewal-risk owners', 'Maya Chen', '2026-06-07', 'Open'],
        ['Publish enablement plan', 'Priya Rao', '2026-06-12', 'Open'],
        ['Start daily incident review', 'Jon Bell', '2026-06-06', 'Ready'],
    ]

    decision_slide = slide_by_exact_title(prs, 'Decision Request')
    decision_text = slide_text(decision_slide)
    for required in [
        'Approve executive outreach for 4 renewal-risk accounts.',
        'Authorize enablement plan for NRR recovery.',
        'Decision needed by 2026-06-07.',
    ]:
        assert required in decision_text
    assert norm_cells(ppt_notes_lines(decision_slide)) == [
        'Speaker note: ask for decision on renewal-risk outreach before reviewing appendix.'
    ]

    appendix_text = slide_text(slide_by_exact_title(prs, 'Appendix - Data Notes'))
    assert 'Public source: QBR metrics workbook, 2026-06-05.' in appendix_text
    assert 'Internal:' not in slide_text(slide_by_exact_title(prs, 'Customer Risks'))

    ref = slide_by_exact_title(prs, 'Reference - Locked Legal')
    src_ref = slide_by_exact_title(src, 'Reference - Locked Legal')
    assert slide_text(ref) == slide_text(src_ref)
    assert norm_cells(ppt_notes_lines(ref)) == norm_cells(ppt_notes_lines(src_ref))
    assert len(ref.shapes) == len(src_ref.shapes)
