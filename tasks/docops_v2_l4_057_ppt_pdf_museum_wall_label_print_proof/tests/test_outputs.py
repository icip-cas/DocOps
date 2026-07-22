import json
import os
import re
import sys
from pathlib import Path
from datetime import date, datetime

import pdfplumber
from pptx import Presentation
from pypdf import PdfReader

sys.path.insert(0, str(Path(__file__).parent))
from verifier_utils import *  # noqa: F401,F403

META_PATH = Path(os.environ.get('TASK_METADATA_PATH', '/tests/task_metadata.json'))
META = json.loads(META_PATH.read_text(encoding='utf-8'))
EXPECT = META['verifier_expectations']


def _vopt_norm_text(value):
    if value is None:
        return ''
    if isinstance(value, datetime):
        if value.hour or value.minute or value.second:
            return value.strftime('%Y-%m-%d %H:%M')
        return value.strftime('%Y-%m-%d')
    if isinstance(value, date):
        return value.strftime('%Y-%m-%d')
    if isinstance(value, float) and value.is_integer():
        value = int(value)
    text = str(value)
    text = text.replace('\u2013', '-').replace('\u2014', '-').replace('\u2212', '-')
    text = text.replace('\xa0', ' ')
    text = re.sub(r'\s+', ' ', text).strip()
    return text


def _vopt_norm_formula(value):
    return re.sub(r'\s+', '', _vopt_norm_text(value)).upper()


def _vopt_assert_formula(actual, expected, label='formula'):
    assert _vopt_norm_formula(actual) == _vopt_norm_formula(expected), f"{label}: expected {expected!r}, found {actual!r}"


def _vopt_norm_rows(rows):
    return [[_vopt_norm_text(cell) for cell in row] for row in rows]


def _vopt_rows_equal(actual, expected):
    return _vopt_norm_rows(actual) == _vopt_norm_rows(expected)


def _vopt_rows_text(rows):
    return '\n'.join('|'.join(_vopt_norm_text(cell) for cell in row) for row in rows)


def _vopt_header_has(header, required):
    actual = {_vopt_norm_text(cell).lower() for cell in header}
    missing = [cell for cell in required if _vopt_norm_text(cell).lower() not in actual]
    assert not missing, f"Missing expected header columns {missing!r}; actual={header!r}"


def _vopt_section_text(section_part):
    return '\n'.join(p.text for p in getattr(section_part, 'paragraphs', []) if p.text is not None)


def _vopt_header_text(doc):
    return '\n'.join(_vopt_section_text(section.header) for section in doc.sections)


def _vopt_footer_text(doc):
    return '\n'.join(_vopt_section_text(section.footer) for section in doc.sections)


def _vopt_table_ref(ws, expected_name=None):
    tables = ws.tables
    if expected_name:
        for name in tables.keys():
            if str(name).lower() == str(expected_name).lower():
                return tables[name].ref
    vals = list(tables.values())
    assert vals, f"No Excel native table found on sheet {ws.title!r}"
    return vals[0].ref


def _vopt_clean_area(value):
    text = _vopt_norm_text(value).replace("'", '').replace('$', '').replace(' ', '')
    if '!' in text:
        text = text.split('!', 1)[1]
    return text.upper()


def _vopt_assert_print_area(ws, expected):
    actual = _vopt_clean_area(ws.print_area)
    target = _vopt_clean_area(expected)
    assert target in actual or actual in target, f"{ws.title}: expected print area {expected!r}, found {ws.print_area!r}"


def _vopt_freeze_pane(value):
    return getattr(value, 'coordinate', value)


def _vopt_validated_cells(ws):
    cells = set()
    for dv in ws.data_validations.dataValidation:
        for rng in dv.cells.ranges:
            text = str(rng).replace('$', '')
            try:
                min_col, min_row, max_col, max_row = range_boundaries(text)
            except ValueError:
                cells.add(text)
                continue
            if max_row > 10000:
                max_row = max(min_row, 200)
            for row in range(min_row, max_row + 1):
                for col in range(min_col, max_col + 1):
                    cells.add(f"{get_column_letter(col)}{row}")
    return cells


class _VoptValidationRanges(list):
    def __init__(self, ws, ranges):
        super().__init__(ranges)
        self.ws = ws

    def __contains__(self, expected):
        expected = str(expected).replace('$', '')
        if super().__contains__(expected):
            return True
        try:
            min_col, min_row, max_col, max_row = range_boundaries(expected)
        except ValueError:
            return super().__contains__(expected)
        target = {f"{get_column_letter(col)}{row}" for row in range(min_row, max_row + 1) for col in range(min_col, max_col + 1)}
        return target.issubset(_vopt_validated_cells(self.ws))


def _vopt_dv_ranges(ws):
    ranges = []
    for dv in ws.data_validations.dataValidation:
        ranges.extend(str(rng).replace('$', '') for rng in dv.cells.ranges)
    return _VoptValidationRanges(ws, ranges)


def _vopt_ordered_subset(expected, actual):
    actual_iter = iter([_vopt_norm_text(item) for item in actual])
    for item in [_vopt_norm_text(value) for value in expected]:
        for candidate in actual_iter:
            if item == candidate:
                break
        else:
            return False
    return True



def _deck_path():
    return Path(os.environ.get('PPT_OUTPUT_PATH', EXPECT['deck_output']))


def _pdf_path():
    return Path(os.environ.get('PDF_OUTPUT_PATH', EXPECT['packet_output']))


def _slide_title(slide):
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            text = ' '.join(shape.text.split())
            if text:
                return text
    return ''


def _slide_text(slide):
    parts = []
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            text = ' '.join(shape.text.split())
            if text:
                parts.append(text)
        if getattr(shape, 'has_table', False):
            for row in shape.table.rows:
                parts.append('|'.join(cell.text.strip() for cell in row.cells))
    return '\n'.join(parts)


def _norm_lower(value):
    return _vopt_norm_text(value).lower()


def _has_terms(text, terms):
    lowered = _norm_lower(text)
    return all(term.lower() in lowered for term in terms)


def _has_omit_note(text, obj_id):
    lowered = _norm_lower(text)
    return obj_id.lower() in lowered and any(term in lowered for term in ['omit', 'exclude', 'restricted', 'do not publish'])


def _has_card_size(text):
    normalized = _norm_lower(text).replace('×', 'x')
    normalized = re.sub(r'\s+', ' ', normalized)
    return bool(re.search(r'6\s*x\s*4\s*(inch|inches|in)', normalized))


def _pdf_texts(path):
    with pdfplumber.open(str(path)) as pdf:
        return [page.extract_text() or '' for page in pdf.pages]


def test_outputs_exist():
    assert _deck_path().exists()
    assert _pdf_path().exists()
    assert _deck_path().suffix.lower() == '.pptx'
    assert _pdf_path().suffix.lower() == '.pdf'


def test_ppt_label_geometry_order_and_rights_boundary():
    prs = Presentation(_deck_path())
    assert abs(prs.slide_width - EXPECT['slide_size'][0]) < 100
    assert abs(prs.slide_height - EXPECT['slide_size'][1]) < 100
    assert len(prs.slides) in (4, 5), f'Unexpected slide count: {len(prs.slides)}'
    slide_texts = [_slide_text(slide) for slide in prs.slides]
    reference_text = slide_texts[-1]
    assert _has_terms(reference_text, ['reference', 'grid']), 'Missing retained reference grid slide'
    label_slides = slide_texts[:-1]
    if len(label_slides) == 4 and not EXPECT['public_ids'][0] in label_slides[0]:
        assert _has_terms(label_slides[0], ['wall', 'label']) or _has_terms(label_slides[0], ['print', 'deck']), 'PPTX cover slide is not recognizable'
        label_slides = label_slides[1:]
    assert len(label_slides) == 3, f'Expected three public label slides plus reference grid, found {len(label_slides)} label slides'
    for slide_text, obj in zip(label_slides, EXPECT['public_objects']):
        for value in obj[:5]:
            assert str(value) in slide_text, f'{value} missing from PPTX label slide for {obj[0]}'
    full = '\n'.join(_slide_text(slide) for slide in prs.slides)
    forbid_any(full, EXPECT['deck_forbidden'], 'deck')
    for rid in EXPECT['restricted_ids']:
        if rid in full:
            assert _has_omit_note(full, rid), 'Restricted ID may be referenced only in an omission statement'
        assert not any(rid in text for text in label_slides), f'{rid} should not appear on public label slides'
        assert f'{rid} |' not in full, f'{rid} should not have a label slide'


def test_pdf_print_proof_pages_bookmarks_and_boundary():
    path = _pdf_path()
    reader = PdfReader(str(path))
    assert len(reader.pages) == EXPECT['pdf_page_count']
    texts = _pdf_texts(path)
    assert _has_terms(texts[0], ['museum', 'wall label', 'print proof']), 'Missing PDF cover title'
    assert _has_card_size(texts[0]), 'PDF cover must state the 6 x 4 inch card size'
    for page_text, obj in zip(texts[1:], EXPECT['public_objects']):
        for value in obj[:5]:
            assert str(value) in page_text, f'{value} missing from PDF proof page for {obj[0]}'
    full = '\n'.join(texts)
    forbid_any(full, EXPECT['pdf_forbidden'], 'pdf')
    for rid in EXPECT['restricted_ids']:
        if rid in full:
            assert rid in texts[0] and _has_omit_note(texts[0], rid), 'Restricted ID may appear only as a cover exclusion note'
            assert not any(rid in page_text for page_text in texts[1:]), f'{rid} should not appear on public proof pages'
    outline = [title for _level, title in flatten_outline(reader.outline)]
    assert any(_has_terms(title, ['cover']) or _has_terms(title, ['wall label']) for title in outline), 'Missing PDF cover bookmark'
    for obj_id in EXPECT['public_ids']:
        assert any(obj_id in title for title in outline), f'Missing PDF bookmark for {obj_id}'


def test_cross_output_public_object_consistency():
    prs = Presentation(_deck_path())
    deck_text = '\n'.join(_slide_text(slide) for slide in prs.slides)
    pdf_text = '\n'.join(_pdf_texts(_pdf_path()))
    for obj_id in EXPECT['public_ids']:
        assert obj_id in deck_text, f'{obj_id} missing from PPTX'
        assert obj_id in pdf_text, f'{obj_id} missing from PDF'
    for obj_id in EXPECT['restricted_ids']:
        assert f'{obj_id} Label Proof' not in pdf_text
        assert f'{obj_id} |' not in deck_text
