from __future__ import annotations

import math
import re
from pathlib import Path
from typing import Iterable

import pdfplumber
import pytest
from docx import Document
from openpyxl import load_workbook
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pypdf import PdfReader


def run_preflight(meta: dict, input_path: Path, output_path: Path) -> None:
    if not input_path.exists():
        raise AssertionError(f'Missing input document: {input_path}')
    if not output_path.exists():
        raise AssertionError(f'Expected output not found: {output_path}')
    assert output_path.stat().st_size > 0, f'Output file is empty: {output_path}'
    assert output_path.suffix.lower() == meta['output_extension'].lower(), (
        f"Output suffix {output_path.suffix} does not match expected {meta['output_extension']}"
    )


def normalize_text(text: str) -> str:
    text = text.replace('\u2019', "'").replace('\u2013', '-').replace('\u2014', '-')
    return ' '.join(text.lower().split())


def answer_text(output_path: Path) -> str:
    return output_path.read_text(encoding='utf-8', errors='ignore').strip()


def require_all(text: str, phrases: Iterable[str], label: str) -> None:
    norm = normalize_text(text)
    missing = [p for p in phrases if normalize_text(p) not in norm]
    assert not missing, f'{label}: missing required phrases: {missing}'


def require_any_group(text: str, groups: list[list[str]], label: str) -> None:
    norm = normalize_text(text)
    for group in groups:
        if any(normalize_text(p) in norm for p in group):
            return
    raise AssertionError(f'{label}: none of the accepted phrase groups were found: {groups}')


def forbid_any(text: str, phrases: Iterable[str], label: str) -> None:
    norm = normalize_text(text)
    hits = [p for p in phrases if normalize_text(p) in norm]
    assert not hits, f'{label}: forbidden phrases still present: {hits}'


def assert_blueish(rgb: str, label: str) -> None:
    assert rgb and len(rgb) == 6, f'{label}: missing RGB color'
    r = int(rgb[0:2], 16)
    g = int(rgb[2:4], 16)
    b = int(rgb[4:6], 16)
    assert b >= r and b >= g, f'{label}: color {rgb} is not blue-dominant'


def flatten_outline(outline) -> list[tuple[int, str]]:
    items = []

    def walk(nodes, level):
        for node in nodes:
            if isinstance(node, list):
                walk(node, level + 1)
            else:
                title = ''
                try:
                    title = str(node.get('/Title', ''))
                except Exception:
                    title = str(node)
                items.append((level, title))

    if isinstance(outline, list):
        walk(outline, 0)
    return items


# -----------------------------
# Word helpers
# -----------------------------

def docx_nonempty_paragraphs(doc: Document):
    return [(i, p) for i, p in enumerate(doc.paragraphs) if p.text.strip()]


def docx_texts(doc: Document) -> list[str]:
    return [p.text.strip() for p in doc.paragraphs if p.text.strip()]


def docx_para_by_text(doc: Document, text: str):
    for p in doc.paragraphs:
        if p.text.strip() == text:
            return p
    raise AssertionError(f'DOCX paragraph not found: {text}')


def first_nonempty_run(paragraph):
    for run in paragraph.runs:
        if run.text.strip():
            return run
    return paragraph.runs[0] if paragraph.runs else None


def docx_run_signature(paragraph):
    run = first_nonempty_run(paragraph)
    if run is None:
        return None
    font = run.font
    color = None
    try:
        if font.color and font.color.rgb:
            color = str(font.color.rgb)
    except Exception:
        color = None
    return {
        'style': paragraph.style.name if paragraph.style else '',
        'name': font.name,
        'size': int(font.size) if font.size else None,
        'bold': bool(font.bold) if font.bold is not None else None,
        'italic': bool(font.italic) if font.italic is not None else None,
        'underline': bool(font.underline) if font.underline is not None else None,
        'color': color,
    }


def docx_para_has_highlight(paragraph) -> bool:
    for run in paragraph.runs:
        if run.font.highlight_color is not None:
            return True
        xml = run._r.xml
        if '<w:shd' in xml or '<w:highlight' in xml:
            return True
    if '<w:shd' in paragraph._p.xml or '<w:highlight' in paragraph._p.xml:
        return True
    return False


def docx_has_page_break(paragraph) -> bool:
    return 'w:type="page"' in paragraph._p.xml or "w:type='page'" in paragraph._p.xml


def docx_inline_shape_sizes(doc: Document):
    return [(int(s.width), int(s.height)) for s in doc.inline_shapes]


def coeff_of_var(nums: list[float]) -> float:
    if not nums:
        return 0.0
    mean = sum(nums) / len(nums)
    if mean == 0:
        return 0.0
    variance = sum((x - mean) ** 2 for x in nums) / len(nums)
    return math.sqrt(variance) / mean


# -----------------------------
# Excel helpers
# -----------------------------

def load_xlsx_pair(path: Path):
    return load_workbook(path), load_workbook(path, data_only=True)


def normalize_formula(formula):
    if formula is None:
        return ""
    text = str(formula).strip()
    if not text.startswith("="):
        return text
    text = text[1:]
    text = re.sub(r"\s+", "", text)
    text = text.replace("$", "")
    text = re.sub(r"'([^']+)'!", r"\1!", text)
    return text.upper()


def formula_equivalent(actual, expected, table_names=None):
    return normalize_formula(actual) == normalize_formula(expected)


def cell_fill_rgb(cell):
    fill = cell.fill
    if fill is None:
        return None
    color = fill.fgColor or fill.start_color
    if color is None:
        return None
    rgb = color.rgb
    if rgb:
        rgb = str(rgb)
    if rgb and len(rgb) == 8:
        rgb = rgb[2:]
    return rgb


def cell_font_rgb(cell):
    color = cell.font.color
    if color is None:
        return None
    rgb = color.rgb
    if rgb:
        rgb = str(rgb)
    if rgb and len(rgb) == 8:
        rgb = rgb[2:]
    return rgb


def style_signature(cell):
    return {
        'font_name': cell.font.name,
        'font_size': float(cell.font.sz) if cell.font.sz else None,
        'bold': bool(cell.font.bold),
        'italic': bool(cell.font.italic),
        'font_rgb': cell_font_rgb(cell),
        'fill_rgb': cell_fill_rgb(cell),
        'alignment_h': cell.alignment.horizontal,
        'alignment_v': cell.alignment.vertical,
        'wrap': bool(cell.alignment.wrap_text),
        'numfmt': cell.number_format,
        'border': str(cell.border),
    }


def row_style_signature(ws, row: int, cols: list[int]) -> list[dict]:
    return [style_signature(ws.cell(row, col)) for col in cols]


def row_values(ws, row: int, max_col: int) -> list:
    return [ws.cell(row, c).value for c in range(1, max_col + 1)]


def workbook_values_signature(wb) -> dict:
    out = {}
    for ws in wb.worksheets:
        out[ws.title] = [tuple(row) for row in ws.iter_rows(values_only=True)]
    return out


def all_cf_ranges(ws) -> list[str]:
    try:
        return [str(r) for r in ws.conditional_formatting]
    except Exception:
        return []


# -----------------------------
# PowerPoint helpers
# -----------------------------

def ppt_shapes_with_text(slide):
    return [shape for shape in slide.shapes if getattr(shape, 'has_text_frame', False)]


def ppt_texts(prs: Presentation) -> list[str]:
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            txt = ' '.join(getattr(shape, 'text', '').split())
            if txt:
                texts.append(txt)
    return texts


def ppt_first_run(paragraph):
    for run in paragraph.runs:
        if run.text.strip():
            return run
    return paragraph.runs[0] if paragraph.runs else None


def ppt_paragraph_signature(paragraph):
    run = ppt_first_run(paragraph)
    if run is None:
        return None
    font = run.font
    color = None
    try:
        if font.color and font.color.type is not None:
            color = str(font.color.rgb)
    except Exception:
        color = None
    return {
        'size': int(font.size) if font.size else None,
        'bold': font.bold,
        'italic': font.italic,
        'name': font.name,
        'color': color,
        'level': paragraph.level,
        'alignment': paragraph.alignment,
    }


def shape_fill_rgb(shape):
    try:
        if shape.fill.type is None:
            return None
        rgb = shape.fill.fore_color.rgb
        return str(rgb) if rgb else None
    except Exception:
        return None


def shape_line_rgb(shape):
    try:
        if shape.line.fill.type is None:
            return None
        rgb = shape.line.fill.fore_color.rgb
        return str(rgb) if rgb else None
    except Exception:
        return None


def slide_title_shape(slide):
    return slide.shapes[0]


# -----------------------------
# PDF helpers
# -----------------------------

def pdf_text(path: Path) -> str:
    reader = PdfReader(str(path))
    return '\n'.join(page.extract_text() or '' for page in reader.pages)


def pdf_page_texts(path: Path) -> list[str]:
    reader = PdfReader(str(path))
    return [(page.extract_text() or '') for page in reader.pages]


def pdf_body_lines(page):
    text = page.extract_text() or ''
    return [line.strip() for line in text.splitlines() if line.strip()]


def approx_equal(a: float, b: float, tol: float) -> bool:
    return abs(a - b) <= tol


def line_text_words(page, target_prefix: str):
    lines = defaultdict(list)
    for word in page.extract_words(extra_attrs=['fontname', 'size']):
        key = round(word['top'], 1)
        lines[key].append(word)
    for words in lines.values():
        words.sort(key=lambda w: w['x0'])
        line = ' '.join(w['text'] for w in words)
        if normalize_text(line).startswith(normalize_text(target_prefix)):
            return words
    return []


def chars_for_line(page, target_prefix: str):
    words = line_text_words(page, target_prefix)
    if not words:
        return []
    x0 = min(w['x0'] for w in words) - 2
    x1 = max(w['x1'] for w in words) + 2
    top = min(w['top'] for w in words) - 2
    bottom = max(w['bottom'] for w in words) + 2
    chars = []
    for ch in page.chars:
        if ch['x0'] >= x0 and ch['x1'] <= x1 and ch['top'] >= top and ch['bottom'] <= bottom:
            chars.append(ch)
    return chars


def line_has_nonblack_chars(page, target_prefix: str) -> bool:
    for ch in chars_for_line(page, target_prefix):
        color = ch.get('non_stroking_color')
        if color not in (None, 0, (0, 0, 0), (0.0, 0.0, 0.0), (0,)):
            return True
    return False


def line_has_colored_rect(page, target_prefix: str) -> bool:
    words = line_text_words(page, target_prefix)
    if not words:
        return False
    x0 = min(w['x0'] for w in words) - 6
    x1 = max(w['x1'] for w in words) + 6
    top = min(w['top'] for w in words) - 4
    bottom = max(w['bottom'] for w in words) + 4
    for rect in page.rects:
        rect_top = rect.get('top', 0)
        rect_bottom = rect.get('bottom', 0)
        rect_x0 = rect.get('x0', 0)
        rect_x1 = rect.get('x1', 0)
        color = rect.get('non_stroking_color')
        if color in (None, (1, 1, 1), (1.0, 1.0, 1.0)):
            continue
        if rect_x1 >= x0 and rect_x0 <= x1 and rect_bottom >= top and rect_top <= bottom:
            return True
    return False


def line_is_visually_highlighted(page, target_prefix: str) -> bool:
    return line_has_nonblack_chars(page, target_prefix) or line_has_colored_rect(page, target_prefix)


def extract_pdf_outline_titles(path: Path) -> list[tuple[int, str]]:
    reader = PdfReader(str(path))
    return flatten_outline(reader.outline)


def pdf_theme_rect_colors(path: Path) -> list[str]:
    colors = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            for rect in page.rects:
                color = rect.get('non_stroking_color')
                if isinstance(color, tuple) and len(color) >= 3:
                    rgb = ''.join(f'{round(max(0, min(1, c)) * 255):02X}' for c in color[:3])
                    colors.append(rgb)
    return colors

from collections import Counter, defaultdict
from pptx.enum.shapes import MSO_SHAPE_TYPE


def first_nonempty_run(paragraph):
    for run in paragraph.runs:
        if run.text.strip():
            return run
    return paragraph.runs[0] if paragraph.runs else None


def docx_run_signature(paragraph):
    run = first_nonempty_run(paragraph)
    if run is None:
        return None
    font = run.font
    color = None
    try:
        if font.color and font.color.rgb:
            color = str(font.color.rgb)
    except Exception:
        color = None
    return {
        'style': paragraph.style.name if paragraph.style else '',
        'name': font.name,
        'size': int(font.size) if font.size else None,
        'bold': bool(font.bold) if font.bold is not None else None,
        'italic': bool(font.italic) if font.italic is not None else None,
        'underline': bool(font.underline) if font.underline is not None else None,
        'color': color,
    }


def docx_signature_matches(target_paragraph, reference_paragraph) -> bool:
    target = docx_run_signature(target_paragraph)
    reference = docx_run_signature(reference_paragraph)
    if target is None or reference is None:
        return target == reference
    if target['style'] != reference['style']:
        return False
    for key in ('bold', 'italic', 'underline', 'color'):
        if target[key] is not None and reference[key] is not None and target[key] != reference[key]:
            return False
    for key in ('name', 'size'):
        if target[key] is not None and reference[key] is not None and target[key] != reference[key]:
            return False
    return True


def has_page_break_before(doc, heading_text):
    prev = None
    for p in doc.paragraphs:
        if p.text.strip() == heading_text:
            if prev is None:
                return False
            return docx_has_page_break(prev) or docx_has_page_break(p)
        prev = p
    raise AssertionError(f'Heading not found for page-break check: {heading_text}')


def docx_header_text(doc):
    parts = []
    for section in doc.sections:
        for p in section.header.paragraphs:
            if p.text.strip():
                parts.append(p.text.strip())
    return '\n'.join(parts)


def docx_footer_text(doc):
    parts = []
    for section in doc.sections:
        for p in section.footer.paragraphs:
            if p.text.strip():
                parts.append(p.text.strip())
    return '\n'.join(parts)


def sheet_data_validation_ranges(ws):
    out = []
    if ws.data_validations is None:
        return out
    for dv in ws.data_validations.dataValidation:
        out.append(str(dv.sqref))
    return out


def workbook_defined_names(wb):
    names = []
    for item in wb.defined_names.values():
        names.append(item.name)
    return sorted(names)


def pdf_page_image_count(path, page_index):
    with pdfplumber.open(path) as pdf:
        return len(pdf.pages[page_index].images)


def pdf_line_is_highlighted(path, page_index, prefix):
    with pdfplumber.open(path) as pdf:
        return line_is_visually_highlighted(pdf.pages[page_index], prefix)


def slide_picture_count(slide):
    return sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)


def ppt_notes_lines(slide):
    lines = []
    notes_slide = slide.notes_slide
    for shape in notes_slide.shapes:
        if not getattr(shape, 'has_text_frame', False):
            continue
        for para in shape.text_frame.paragraphs:
            text = ' '.join(para.text.split())
            if text:
                lines.append(text)
    return lines


def ppt_title_signature(slide):
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            for para in shape.text_frame.paragraphs:
                text = ' '.join(para.text.split())
                if text:
                    return ppt_paragraph_signature(para)
    return None


def ppt_first_body_signature(slide):
    title_seen = False
    for shape in slide.shapes:
        if not getattr(shape, 'has_text_frame', False):
            continue
        for para in shape.text_frame.paragraphs:
            text = ' '.join(para.text.split())
            if not text:
                continue
            if not title_seen:
                title_seen = True
                continue
            return ppt_paragraph_signature(para)
    return None
def docx_has_page_break(paragraph) -> bool:
    xml = paragraph._p.xml
    if 'w:type="page"' in xml or "w:type='page'" in xml:
        return True
    if '<w:pageBreakBefore' in xml:
        return True
    try:
        if paragraph.paragraph_format.page_break_before:
            return True
    except Exception:
        pass
    return False


def has_page_break_before(doc, heading_text):
    prev = None
    for p in doc.paragraphs:
        if p.text.strip() == heading_text:
            if prev is None:
                return docx_has_page_break(p)
            return docx_has_page_break(prev) or docx_has_page_break(p)
        prev = p
    raise AssertionError(f'Heading not found for page-break check: {heading_text}')


def docx_has_toc_field(doc) -> bool:
    xml = doc._element.xml
    normalized = ' '.join(xml.replace('\u201c', '"').replace('\u201d', '"').split())
    return (
        'TOC \\o "1-3" \\h \\z \\u' in normalized
        or 'TOC \\o "1-1" \\h \\z \\u' in normalized
        or 'TOC ' in normalized
    )


def docx_section_header_text_at(doc, idx):
    parts = []
    section = doc.sections[idx]
    for p in section.header.paragraphs:
        if p.text.strip():
            parts.append(p.text.strip())
    return '\n'.join(parts)


def docx_section_footer_text_at(doc, idx):
    parts = []
    section = doc.sections[idx]
    for p in section.footer.paragraphs:
        if p.text.strip():
            parts.append(p.text.strip())
    return '\n'.join(parts)


def sheet_row_highlighted(ws, row_idx, min_col, max_col):
    for col in range(min_col, max_col + 1):
        rgb = cell_fill_rgb(ws.cell(row_idx, col))
        if rgb and rgb not in ('FFFFFF', '000000', '00000000'):
            return True
    if min_col <= 6 <= max_col:
        status_value = ws.cell(row_idx, 6).value
        if status_value == 'Late':
            ranges = all_cf_ranges(ws)
            if any(r.startswith('<ConditionalFormatting A2:H') for r in ranges):
                return True
    return False


def expand_sqref_cells(sqrefs):
    from openpyxl.utils import get_column_letter, range_boundaries

    cells = set()
    for sqref in sqrefs or []:
        for token in str(sqref).replace(',', ' ').split():
            if ':' in token:
                min_col, min_row, max_col, max_row = range_boundaries(token)
                for row in range(min_row, max_row + 1):
                    for col in range(min_col, max_col + 1):
                        cells.add(f"{get_column_letter(col)}{row}")
            else:
                cells.add(token)
    return sorted(cells)

# -----------------------------
# Robust verifier helpers added for L3 validator hardening
# -----------------------------

def norm_cell(value):
    if value is None:
        return ''
    if hasattr(value, 'strftime'):
        try:
            if getattr(value, 'hour', 0) or getattr(value, 'minute', 0) or getattr(value, 'second', 0):
                return value.strftime('%Y-%m-%d %H:%M:%S')
            return value.strftime('%Y-%m-%d')
        except Exception:
            pass
    if isinstance(value, float) and value.is_integer():
        value = int(value)
    return ' '.join(str(value).replace('\u2013', '-').replace('\u2014', '-').split())


def norm_cells(values):
    return [norm_cell(v) for v in values]


def norm_table_rows(rows):
    return [norm_cells(row) for row in rows]


def norm_join(values):
    if isinstance(values, str):
        return normalize_text(values)
    return normalize_text('\n'.join(norm_cell(v) for v in values))


def assert_norm_equal(actual, expected, label='value'):
    assert norm_cell(actual) == norm_cell(expected), f'{label}: expected {expected!r}, found {actual!r}'


def assert_norm_list_equal(actual, expected, label='values'):
    actual_norm = norm_cells(actual)
    expected_norm = norm_cells(expected)
    assert actual_norm == expected_norm, f'{label}: expected {expected_norm!r}, found {actual_norm!r}'


def assert_required_terms(text, terms, label='text'):
    norm = normalize_text(text)
    missing = [term for term in terms if normalize_text(term) not in norm]
    assert not missing, f'{label}: missing required terms: {missing}'


def assert_forbidden_terms(text, terms, label='text'):
    norm = normalize_text(text)
    hits = [term for term in terms if normalize_text(term) in norm]
    assert not hits, f'{label}: forbidden terms present: {hits}'


def table_ref(ws, expected_name=None):
    tables = ws.tables
    if expected_name:
        try:
            return tables[expected_name].ref
        except Exception:
            expected_norm = str(expected_name).lower()
            for name, table in tables.items():
                if str(name).lower() == expected_norm:
                    return table.ref
    vals = list(tables.values())
    if vals:
        return vals[0].ref
    raise AssertionError(f'{ws.title}: expected a native Excel table')


def freeze_pane_ref(ws):
    value = ws.freeze_panes
    if value is None:
        return None
    return getattr(value, 'coordinate', str(value))


def sheet_hidden_state_ok(ws, expected='hidden'):
    if expected in ('hidden', 'veryHidden'):
        return ws.sheet_state in ('hidden', 'veryHidden')
    return ws.sheet_state == expected


def text_lines_normalized(text):
    return [' '.join(line.split()) for line in str(text or '').splitlines() if line.strip()]


def likely_pdf_title(page_text, expected):
    expected_norm = normalize_text(expected)
    lines = [line for line in text_lines_normalized(page_text) if not re.fullmatch(r'(page\s*)?\d+(\s+of\s+\d+)?', line, flags=re.I)]
    head = ' '.join(lines[:4])
    return expected_norm in normalize_text(head) or expected_norm in normalize_text(page_text)


def assert_pdf_titles_in_order(page_texts, expected_titles):
    assert len(page_texts) == len(expected_titles), f'Unexpected page count/order: {len(page_texts)} pages vs {len(expected_titles)} titles'
    for idx, (page_text, expected) in enumerate(zip(page_texts, expected_titles)):
        assert likely_pdf_title(page_text, expected), f'Page {idx + 1}: expected title {expected!r}'


def assert_pdf_outline_exact(path, expected_outline):
    actual = extract_pdf_outline_titles(path)
    actual_norm = [(level, normalize_text(title)) for level, title in actual]
    expected_norm = [(level, normalize_text(title)) for level, title in expected_outline]
    assert actual_norm == expected_norm, f'Unexpected outline: {actual!r}'


def ppt_slide_text(slide):
    parts = []
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            text = ' '.join(shape.text.split())
            if text:
                parts.append(text)
    return '\n'.join(parts)


def ppt_slide_title_text(slide):
    try:
        if slide.shapes.title and getattr(slide.shapes.title, 'has_text_frame', False):
            title = ' '.join(slide.shapes.title.text.split())
            if title:
                return title
    except Exception:
        pass
    candidates = []
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            lines = text_lines_normalized(shape.text)
            if lines:
                candidates.append(lines[0])
    return candidates[0] if candidates else ''


def ppt_slide_titles(prs):
    return [ppt_slide_title_text(slide) for slide in prs.slides]


def ppt_find_slide(prs, title):
    target = normalize_text(title)
    for slide in prs.slides:
        title_text = ppt_slide_title_text(slide)
        full_text = ppt_slide_text(slide)
        if normalize_text(title_text) == target or target in normalize_text(full_text):
            return slide
    raise AssertionError(f'Slide not found: {title}')


def assert_ppt_titles_equal(prs, expected_titles):
    actual = ppt_slide_titles(prs)
    actual_norm = [normalize_text(t) for t in actual]
    expected_norm = [normalize_text(t) for t in expected_titles]
    assert actual_norm == expected_norm, f'Unexpected slide order: {actual!r}'


def table_cell_text(table, row, col):
    return norm_cell(table.cell(row, col).text)


def assert_table_header_contains(table, cols, required, label='table'):
    header = [table_cell_text(table, 0, c) for c in range(cols)]
    header_norm = {normalize_text(x) for x in header}
    missing = [x for x in required if normalize_text(x) not in header_norm]
    assert not missing, f'{label}: missing header columns {missing}; found {header!r}'

