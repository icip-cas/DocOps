from __future__ import annotations

import math
import re
from pathlib import Path
from typing import Iterable

import pdfplumber
import pytest
from docx import Document
from openpyxl import load_workbook
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pypdf import PdfReader


def run_preflight(meta: dict, input_path: Path, output_path: Path) -> None:
    if not input_path.exists():
        raise AssertionError(f'Missing input document: {input_path}')
    if not output_path.exists():
        raise AssertionError(f'Expected output not found: {output_path}')
    assert output_path.stat().st_size > 0, f'Output file is empty: {output_path}'
    assert output_path.suffix.lower() == meta['output_extension'].lower(), (
        f"Output suffix {output_path.suffix} does not match expected {meta['output_extension']}"
    )


def normalize_text(text: str) -> str:
    text = text.replace('\u2019', "'").replace('\u2013', '-').replace('\u2014', '-')
    text = re.sub(r'(?<=\w)[-/](?=\w)', ' ', text)
    return ' '.join(text.lower().split())


def answer_text(output_path: Path) -> str:
    return output_path.read_text(encoding='utf-8', errors='ignore').strip()


def require_all(text: str, phrases: Iterable[str], label: str) -> None:
    norm = normalize_text(text)
    missing = [p for p in phrases if normalize_text(p) not in norm]
    assert not missing, f'{label}: missing required phrases: {missing}'


def require_any_group(text: str, groups: list[list[str]], label: str) -> None:
    norm = normalize_text(text)
    for group in groups:
        if any(normalize_text(p) in norm for p in group):
            return
    raise AssertionError(f'{label}: none of the accepted phrase groups were found: {groups}')


def require_group_hits(text: str, groups: list[list[str]], label: str) -> None:
    norm = normalize_text(text)
    missing = [group for group in groups if not any(normalize_text(p) in norm for p in group)]
    assert not missing, f'{label}: missing required concept groups: {missing}'


def require_ordered_anchors(text: str, anchors: list[str], label: str) -> None:
    norm = normalize_text(text)
    pos = -1
    for anchor in anchors:
        anchor_norm = normalize_text(anchor)
        next_pos = norm.find(anchor_norm, pos + 1)
        assert next_pos != -1, f'{label}: missing anchor {anchor!r}'
        assert next_pos > pos, f'{label}: anchor order is incorrect around {anchor!r}'
        pos = next_pos


def assert_prefixed_items(items: list[str], prefixes: list[str], label: str) -> None:
    assert len(items) == len(prefixes), f'{label}: expected {len(prefixes)} items, found {len(items)}'
    for idx, (item, prefix) in enumerate(zip(items, prefixes), start=1):
        item_norm = normalize_text(item)
        prefix_norm = normalize_text(prefix)
        assert item_norm.startswith(prefix_norm), (
            f'{label}: item {idx} must start with {prefix!r}, got {item!r}'
        )


def forbid_any(text: str, phrases: Iterable[str], label: str) -> None:
    norm = normalize_text(text)
    hits = [p for p in phrases if normalize_text(p) in norm]
    assert not hits, f'{label}: forbidden phrases still present: {hits}'


def assert_blueish(rgb: str, label: str) -> None:
    assert rgb and len(rgb) == 6, f'{label}: missing RGB color'
    r = int(rgb[0:2], 16)
    g = int(rgb[2:4], 16)
    b = int(rgb[4:6], 16)
    assert b >= r and b >= g, f'{label}: color {rgb} is not blue-dominant'


def flatten_outline(outline) -> list[tuple[int, str]]:
    items = []

    def walk(nodes, level):
        for node in nodes:
            if isinstance(node, list):
                walk(node, level + 1)
            else:
                title = ''
                try:
                    title = str(node.get('/Title', ''))
                except Exception:
                    title = str(node)
                items.append((level, title))

    if isinstance(outline, list):
        walk(outline, 0)
    return items


# -----------------------------
# Word helpers
# -----------------------------

def docx_nonempty_paragraphs(doc: Document):
    return [(i, p) for i, p in enumerate(doc.paragraphs) if p.text.strip()]


def docx_texts(doc: Document) -> list[str]:
    return [p.text.strip() for p in doc.paragraphs if p.text.strip()]


def docx_para_by_text(doc: Document, text: str):
    for p in doc.paragraphs:
        if p.text.strip() == text:
            return p
    raise AssertionError(f'DOCX paragraph not found: {text}')


def first_nonempty_run(paragraph):
    for run in paragraph.runs:
        if run.text.strip():
            return run
    return paragraph.runs[0] if paragraph.runs else None


def docx_run_signature(paragraph):
    run = first_nonempty_run(paragraph)
    if run is None:
        return None
    font = run.font
    color = None
    try:
        if font.color and font.color.rgb:
            color = str(font.color.rgb)
    except Exception:
        color = None
    return {
        'style': paragraph.style.name if paragraph.style else '',
        'name': font.name,
        'size': int(font.size) if font.size else None,
        'bold': bool(font.bold) if font.bold is not None else None,
        'italic': bool(font.italic) if font.italic is not None else None,
        'underline': bool(font.underline) if font.underline is not None else None,
        'color': color,
    }


def docx_para_has_highlight(paragraph) -> bool:
    for run in paragraph.runs:
        if run.font.highlight_color is not None:
            return True
        xml = run._r.xml
        if '<w:shd' in xml or '<w:highlight' in xml:
            return True
    if '<w:shd' in paragraph._p.xml or '<w:highlight' in paragraph._p.xml:
        return True
    return False


def docx_has_page_break(paragraph) -> bool:
    return 'w:type="page"' in paragraph._p.xml or "w:type='page'" in paragraph._p.xml


def docx_inline_shape_sizes(doc: Document):
    return [(int(s.width), int(s.height)) for s in doc.inline_shapes]


def coeff_of_var(nums: list[float]) -> float:
    if not nums:
        return 0.0
    mean = sum(nums) / len(nums)
    if mean == 0:
        return 0.0
    variance = sum((x - mean) ** 2 for x in nums) / len(nums)
    return math.sqrt(variance) / mean


# -----------------------------
# Excel helpers
# -----------------------------

def load_xlsx_pair(path: Path):
    return load_workbook(path), load_workbook(path, data_only=True)


def cell_fill_rgb(cell):
    fill = cell.fill
    if fill is None or fill.patternType is None:
        return None
    color = fill.fgColor or fill.start_color
    if color is None:
        return None
    rgb = color.rgb
    if rgb and len(rgb) == 8:
        rgb = rgb[2:]
    if not rgb or rgb.upper() == '000000':
        return None
    return rgb


def cell_font_rgb(cell):
    color = cell.font.color
    if color is None:
        return None
    rgb = color.rgb
    if rgb:
        rgb = str(rgb)
    if rgb and len(rgb) == 8:
        rgb = rgb[2:]
    return rgb


def style_signature(cell):
    return {
        'font_name': cell.font.name,
        'font_size': float(cell.font.sz) if cell.font.sz else None,
        'bold': bool(cell.font.bold),
        'italic': bool(cell.font.italic),
        'font_rgb': cell_font_rgb(cell),
        'fill_rgb': cell_fill_rgb(cell),
        'alignment_h': cell.alignment.horizontal,
        'alignment_v': cell.alignment.vertical,
        'wrap': bool(cell.alignment.wrap_text),
        'numfmt': cell.number_format,
        'border': str(cell.border),
    }


def row_style_signature(ws, row: int, cols: list[int]) -> list[dict]:
    return [style_signature(ws.cell(row, col)) for col in cols]


def row_values(ws, row: int, max_col: int) -> list:
    return [ws.cell(row, c).value for c in range(1, max_col + 1)]


def workbook_values_signature(wb) -> dict:
    out = {}
    for ws in wb.worksheets:
        out[ws.title] = [tuple(row) for row in ws.iter_rows(values_only=True)]
    return out


def all_cf_ranges(ws) -> list[str]:
    try:
        return [str(r) for r in ws.conditional_formatting]
    except Exception:
        return []


# -----------------------------
# PowerPoint helpers
# -----------------------------

def ppt_shapes_with_text(slide):
    return [shape for shape in slide.shapes if getattr(shape, 'has_text_frame', False)]


def ppt_texts(prs: Presentation) -> list[str]:
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            txt = ' '.join(getattr(shape, 'text', '').split())
            if txt:
                texts.append(txt)
    return texts


def ppt_first_run(paragraph):
    for run in paragraph.runs:
        if run.text.strip():
            return run
    return paragraph.runs[0] if paragraph.runs else None


def ppt_paragraph_signature(paragraph):
    run = ppt_first_run(paragraph)
    if run is None:
        return None
    font = run.font
    color = None
    try:
        if font.color and font.color.type is not None:
            color = str(font.color.rgb)
    except Exception:
        color = None
    return {
        'size': int(font.size) if font.size else None,
        'bold': font.bold,
        'italic': font.italic,
        'name': font.name,
        'color': color,
        'level': paragraph.level,
        'alignment': paragraph.alignment,
    }


def shape_fill_rgb(shape):
    try:
        if shape.fill.type is None:
            return None
        rgb = shape.fill.fore_color.rgb
        return str(rgb) if rgb else None
    except Exception:
        return None


def shape_line_rgb(shape):
    try:
        if shape.line.fill.type is None:
            return None
        rgb = shape.line.fill.fore_color.rgb
        return str(rgb) if rgb else None
    except Exception:
        return None


def slide_title_shape(slide):
    return slide.shapes[0]


# -----------------------------
# PDF helpers
# -----------------------------

def pdf_text(path: Path) -> str:
    reader = PdfReader(str(path))
    return '\n'.join(page.extract_text() or '' for page in reader.pages)


def pdf_page_texts(path: Path) -> list[str]:
    reader = PdfReader(str(path))
    return [(page.extract_text() or '') for page in reader.pages]


def pdf_body_lines(page):
    text = page.extract_text() or ''
    return [line.strip() for line in text.splitlines() if line.strip()]


def approx_equal(a: float, b: float, tol: float) -> bool:
    return abs(a - b) <= tol


def line_text_words(page, target_prefix: str):
    lines = defaultdict(list)
    for word in page.extract_words(extra_attrs=['fontname', 'size']):
        key = round(word['top'], 1)
        lines[key].append(word)
    for words in lines.values():
        words.sort(key=lambda w: w['x0'])
        line = ' '.join(w['text'] for w in words)
        if normalize_text(line).startswith(normalize_text(target_prefix)):
            return words
    return []


def chars_for_line(page, target_prefix: str):
    words = line_text_words(page, target_prefix)
    if not words:
        return []
    x0 = min(w['x0'] for w in words) - 2
    x1 = max(w['x1'] for w in words) + 2
    top = min(w['top'] for w in words) - 2
    bottom = max(w['bottom'] for w in words) + 2
    chars = []
    for ch in page.chars:
        if ch['x0'] >= x0 and ch['x1'] <= x1 and ch['top'] >= top and ch['bottom'] <= bottom:
            chars.append(ch)
    return chars


def line_has_nonblack_chars(page, target_prefix: str) -> bool:
    for ch in chars_for_line(page, target_prefix):
        color = ch.get('non_stroking_color')
        if color not in (None, 0, (0, 0, 0), (0.0, 0.0, 0.0), (0,)):
            return True
    return False


def line_has_colored_rect(page, target_prefix: str) -> bool:
    words = line_text_words(page, target_prefix)
    if not words:
        return False
    x0 = min(w['x0'] for w in words) - 6
    x1 = max(w['x1'] for w in words) + 6
    top = min(w['top'] for w in words) - 4
    bottom = max(w['bottom'] for w in words) + 4
    for rect in page.rects:
        rect_top = rect.get('top', 0)
        rect_bottom = rect.get('bottom', 0)
        rect_x0 = rect.get('x0', 0)
        rect_x1 = rect.get('x1', 0)
        color = rect.get('non_stroking_color')
        if color in (None, (1, 1, 1), (1.0, 1.0, 1.0)):
            continue
        if rect_x1 >= x0 and rect_x0 <= x1 and rect_bottom >= top and rect_top <= bottom:
            return True
    return False


def line_is_visually_highlighted(page, target_prefix: str) -> bool:
    return line_has_nonblack_chars(page, target_prefix) or line_has_colored_rect(page, target_prefix)


def extract_pdf_outline_titles(path: Path) -> list[tuple[int, str]]:
    reader = PdfReader(str(path))
    return flatten_outline(reader.outline)


def pdf_theme_rect_colors(path: Path) -> list[str]:
    colors = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            for rect in page.rects:
                color = rect.get('non_stroking_color')
                if isinstance(color, tuple) and len(color) >= 3:
                    rgb = ''.join(f'{round(max(0, min(1, c)) * 255):02X}' for c in color[:3])
                    colors.append(rgb)
    return colors
