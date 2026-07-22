import json
import sys
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.path.insert(0, str(Path(__file__).parent))
from verifier_utils import *  # noqa: F401,F403

META = json.loads(Path('/tests/task_metadata.json').read_text(encoding='utf-8'))
INPUT_PATH = Path(META['input_path'])
OUTPUT_PATH = Path(META['output_path'])
EXPECT = META['verifier_expectations']


def range_signatures(ws, refs):
    return [style_signature(ws[ref]) for ref in refs]


def sheet_row_highlighted(ws, row_idx, min_col, max_col):
    for col in range(min_col, max_col + 1):
        rgb = cell_fill_rgb(ws.cell(row_idx, col))
        if rgb and rgb not in ('FFFFFF', '000000', '00000000'):
            return True
    if min_col <= 6 <= max_col:
        status_value = ws.cell(row_idx, 6).value
        if status_value == 'Late':
            ranges = all_cf_ranges(ws)
            if any(r.startswith('<ConditionalFormatting A2:H') for r in ranges):
                return True
    return False


def heading_paragraphs(doc):
    out = []
    for p in doc.paragraphs:
        style = p.style.name if p.style else ''
        if p.text.strip() and style.startswith('Heading'):
            out.append((p.text.strip(), style))
    return out


def paragraphs_under_heading(doc, heading_text):
    paras = []
    capture = False
    for p in doc.paragraphs:
        text = p.text.strip()
        style = p.style.name if p.style else ''
        if text == heading_text:
            capture = True
            continue
        if capture and text and style.startswith('Heading'):
            break
        if capture and text:
            paras.append(text)
    return paras


def compact_text(text):
    return normalize_text(str(text)).replace('\xa0', ' ')


def assert_contains_all(text, phrases, label):
    norm = compact_text(text)
    missing = [phrase for phrase in phrases if compact_text(phrase) not in norm]
    assert not missing, f'{label}: missing required anchors: {missing}'


def slide_titles(prs):
    titles = []
    for slide in prs.slides:
        title = None
        for shape in slide.shapes:
            if getattr(shape, 'has_text_frame', False):
                txt = ' '.join(shape.text.split())
                if txt:
                    title = txt.split('\n')[0]
                    break
        titles.append(title or '')
    return titles


def slide_by_title(prs, title):
    norm = normalize_text(title)
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, 'has_text_frame', False):
                txt = ' '.join(shape.text.split())
                if txt and normalize_text(txt.split('\n')[0]) == norm:
                    return slide
    raise AssertionError(f'Slide not found: {title}')


def slide_text_lines(slide, skip_title=True):
    lines = []
    title_text = None
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            txt = ' '.join(shape.text.split())
            if txt:
                title_text = txt.split('\n')[0]
                break
    for shape in slide.shapes:
        if not getattr(shape, 'has_text_frame', False):
            continue
        for para in shape.text_frame.paragraphs:
            text = ' '.join(para.text.split())
            if not text:
                continue
            if skip_title and title_text and normalize_text(text) == normalize_text(title_text):
                continue
            lines.append(text)
    return lines


def slide_table(slide):
    for shape in slide.shapes:
        if getattr(shape, 'has_table', False):
            return shape.table
    raise AssertionError('No table found on slide.')


def page_titles_from_texts(page_texts):
    titles = []
    for text in page_texts:
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        titles.append(lines[0] if lines else '')
    return titles


def page_text_by_title(page_texts, title):
    titles = page_titles_from_texts(page_texts)
    for idx, current in enumerate(titles):
        if normalize_text(current) == normalize_text(title):
            return page_texts[idx]
    raise AssertionError(f'PDF page title not found: {title}')


def verify_word():
    doc = Document(OUTPUT_PATH)
    headings = heading_paragraphs(doc)

    if 'heading_sequence' in EXPECT:
        actual = headings
        if actual and actual[0][0] not in [item[0] for item in EXPECT['heading_sequence']]:
            actual = actual[1:]
        assert actual == [tuple(item) for item in EXPECT['heading_sequence']], (
            f'Unexpected heading sequence: {actual!r}'
        )

    if 'heading_order' in EXPECT:
        actual = [text for text, _style in headings]
        # Some templates style the document title as a heading even though it is
        # not part of the section-order requirement. Ignore that leading title.
        if actual and actual[0] not in EXPECT['heading_order']:
            actual = actual[1:]
        assert actual == EXPECT['heading_order'], f'Unexpected heading order: {actual!r}'

    for text in EXPECT.get('required_paragraphs', []):
        assert text in docx_texts(doc), f'Required paragraph missing: {text}'

    for text in EXPECT.get('forbidden_paragraphs', []):
        assert text not in docx_texts(doc), f'Forbidden paragraph present: {text}'

    for spec in EXPECT.get('paragraphs_under_headings', []):
        actual = paragraphs_under_heading(doc, spec['heading'])
        assert actual == spec['exact_lines'], (
            f"Unexpected content under heading {spec['heading']}: {actual!r}"
        )

    for table_spec in EXPECT.get('tables', []):
        table = doc.tables[table_spec['index']]
        for key, expected in table_spec['cells'].items():
            row_idx, col_idx = [int(x) for x in key.split(',')]
            actual = table.cell(row_idx, col_idx).text.strip()
            assert actual == expected, f"table[{table_spec['index']}][{key}] expected {expected!r}, found {actual!r}"

    for spec in EXPECT.get('table_sizes', []):
        table = doc.tables[spec['index']]
        assert len(table.rows) == spec['rows'], f"Unexpected row count on table {spec['index']}"
        assert len(table.columns) == spec['cols'], f"Unexpected col count on table {spec['index']}"

    for spec in EXPECT.get('style_matches', []):
        target = docx_para_by_text(doc, spec['target_text'])
        reference = docx_para_by_text(doc, spec['reference_text'])
        assert docx_signature_matches(target, reference), (
            f"Style mismatch between {spec['target_text']!r} and {spec['reference_text']!r}"
        )

    for heading in EXPECT.get('page_break_before', []):
        assert has_page_break_before(doc, heading), f'Missing page break before heading: {heading}'

    if 'section_count' in EXPECT:
        assert len(doc.sections) == EXPECT['section_count'], (
            f"Unexpected section count: {len(doc.sections)}"
        )

    if EXPECT.get('toc_present'):
        assert docx_has_toc_field(doc), 'Expected TOC field in document XML'

    for spec in EXPECT.get('section_headers', []):
        actual = docx_section_header_text_at(doc, spec['index'])
        require_all(actual, spec.get('contains', []), f"section {spec['index']} header")

    for spec in EXPECT.get('section_footers', []):
        actual = docx_section_footer_text_at(doc, spec['index'])
        require_all(actual, spec.get('contains', []), f"section {spec['index']} footer")

    for spec in EXPECT.get('section_links', []):
        section = doc.sections[spec['index']]
        assert section.header.is_linked_to_previous == spec['header'], (
            f"Section {spec['index']} header link mismatch"
        )
        assert section.footer.is_linked_to_previous == spec['footer'], (
            f"Section {spec['index']} footer link mismatch"
        )

    header_text = docx_header_text(doc)
    footer_text = docx_footer_text(doc)
    require_all(header_text, EXPECT.get('header_contains', []), 'docx header')
    require_all(footer_text, EXPECT.get('footer_contains', []), 'docx footer')

    if 'inline_shape_count' in EXPECT:
        assert len(doc.inline_shapes) == EXPECT['inline_shape_count'], (
            f"Unexpected inline shape count: {len(doc.inline_shapes)}"
        )

    for text in EXPECT.get('highlighted_paragraphs', []):
        para = docx_para_by_text(doc, text)
        assert docx_para_has_highlight(para), f'Expected highlighted paragraph: {text}'


def verify_excel():
    wb = load_workbook(OUTPUT_PATH)
    in_wb = load_workbook(INPUT_PATH)

    if 'sheet_order' in EXPECT:
        assert wb.sheetnames == EXPECT['sheet_order'], f'Unexpected sheet order: {wb.sheetnames!r}'

    for ref, expected in EXPECT.get('exact_cells', {}).items():
        sheet_name, cell_ref = ref.split('!')
        actual = wb[sheet_name][cell_ref].value
        actual = '' if actual is None else str(actual)
        assert actual == expected, f'{ref}: expected {expected!r}, found {actual!r}'

    for ref, expected in EXPECT.get('formula_cells', {}).items():
        sheet_name, cell_ref = ref.split('!')
        actual = wb[sheet_name][cell_ref].value
        assert actual == expected, f'{ref}: expected formula {expected!r}, found {actual!r}'

    for spec in EXPECT.get('header_style_matches', []):
        target = range_signatures(wb[spec['target_sheet']], spec['target_cells'])
        input_target = range_signatures(in_wb[spec['target_sheet']], spec['target_cells'])
        reference = input_target
        if spec.get('reference_sheet') and spec.get('reference_cells'):
            input_reference = range_signatures(in_wb[spec['reference_sheet']], spec['reference_cells'])
            if input_target == input_reference:
                reference = range_signatures(wb[spec['reference_sheet']], spec['reference_cells'])
        assert target == reference, f"Header style mismatch for {spec['target_sheet']}"

    for spec in EXPECT.get('highlight_rows', []):
        ws = wb[spec['sheet']]
        assert sheet_row_highlighted(ws, spec['row'], spec['min_col'], spec['max_col']), (
            f"Expected highlighted row {spec['sheet']}!{spec['row']}"
        )

    for spec in EXPECT.get('nonhighlight_rows', []):
        ws = wb[spec['sheet']]
        assert not sheet_row_highlighted(ws, spec['row'], spec['min_col'], spec['max_col']), (
            f"Row should not be highlighted: {spec['sheet']}!{spec['row']}"
        )

    for sheet_name in EXPECT.get('unchanged_sheets', []):
        assert workbook_values_signature(wb)[sheet_name] == workbook_values_signature(in_wb)[sheet_name], (
            f'Sheet values changed unexpectedly: {sheet_name}'
        )

    for ref, expected in EXPECT.get('table_refs', {}).items():
        sheet_name, table_name = ref.split('!')
        actual = wb[sheet_name].tables[table_name].ref
        assert actual == expected, f'{ref}: expected table ref {expected!r}, found {actual!r}'

    for sheet_name, expected in EXPECT.get('sheet_states', {}).items():
        assert wb[sheet_name].sheet_state == expected, (
            f'{sheet_name}: expected state {expected!r}, found {wb[sheet_name].sheet_state!r}'
        )

    for sheet_name, expected in EXPECT.get('freeze_panes', {}).items():
        actual = wb[sheet_name].freeze_panes
        actual = str(actual) if actual else None
        assert actual == expected, f'{sheet_name}: expected freeze panes {expected!r}, found {actual!r}'

    for sheet_name, expected in EXPECT.get('conditional_format_ranges', {}).items():
        actual = all_cf_ranges(wb[sheet_name])
        assert actual == expected, f'{sheet_name}: unexpected CF ranges {actual!r}'

    for sheet_name, expected in EXPECT.get('data_validation_ranges', {}).items():
        actual = sheet_data_validation_ranges(wb[sheet_name])
        assert actual == expected, f'{sheet_name}: unexpected validation ranges {actual!r}'

    if 'defined_names' in EXPECT:
        actual = workbook_defined_names(wb)
        assert actual == EXPECT['defined_names'], f'Unexpected defined names: {actual!r}'

    for ref, expected in EXPECT.get('comments_contain', {}).items():
        sheet_name, cell_ref = ref.split('!')
        cell = wb[sheet_name][cell_ref]
        text = cell.comment.text if cell.comment else ''
        assert expected in text, f'{ref}: expected comment containing {expected!r}, found {text!r}'


def verify_pdf():
    page_texts = pdf_page_texts(OUTPUT_PATH)
    titles = page_titles_from_texts(page_texts)

    if 'page_count' in EXPECT:
        assert len(page_texts) == EXPECT['page_count'], f'Unexpected page count: {len(page_texts)}'

    if 'page_titles_order' in EXPECT:
        assert titles == EXPECT['page_titles_order'], f'Unexpected page order: {titles!r}'

    for title in EXPECT.get('absent_page_titles', []):
        assert title not in titles, f'Page title should be absent: {title}'

    cover_text = page_texts[0] if page_texts else ''
    for line in EXPECT.get('cover_lines', []):
        assert normalize_text(line) in normalize_text(cover_text), f'Missing cover line: {line}'

    for spec in EXPECT.get('required_page_lines', []):
        text = page_text_by_title(page_texts, spec['page_title'])
        norm = normalize_text(text)
        for line in spec['lines']:
            assert normalize_text(line) in norm, f"Missing page line on {spec['page_title']}: {line}"

    for spec in EXPECT.get('forbidden_page_lines', []):
        text = page_text_by_title(page_texts, spec['page_title'])
        norm = normalize_text(text)
        for line in spec['lines']:
            assert normalize_text(line) not in norm, f"Forbidden page line remained on {spec['page_title']}: {line}"

    if 'outline' in EXPECT:
        outline = extract_pdf_outline_titles(OUTPUT_PATH)
        assert outline == [tuple(item) for item in EXPECT['outline']], f'Unexpected outline: {outline!r}'

    for spec in EXPECT.get('page_image_mins', []):
        page_index = spec.get('page_index')
        if page_index is None:
            page_index = titles.index(spec['page_title'])
        actual = pdf_page_image_count(OUTPUT_PATH, page_index)
        assert actual >= spec['min_images'], f'Expected at least {spec["min_images"]} images on page {page_index}'

    for spec in EXPECT.get('highlighted_page_lines', []):
        title = spec['page_title']
        idx = titles.index(title)
        assert pdf_line_is_highlighted(OUTPUT_PATH, idx, spec['prefix']), (
            f'Expected highlighted line on {title}: {spec["prefix"]}'
        )

    if 'theme_colors_any' in EXPECT:
        actual = pdf_theme_rect_colors(OUTPUT_PATH)
        for color in EXPECT['theme_colors_any']:
            assert color in actual, f'Expected theme color {color} in PDF rect palette'


def verify_ppt():
    prs = Presentation(OUTPUT_PATH)
    in_prs = Presentation(INPUT_PATH)

    if 'titles_order' in EXPECT:
        actual = slide_titles(prs)
        assert actual == EXPECT['titles_order'], f'Unexpected slide order: {actual!r}'

    for title in EXPECT.get('absent_slide_titles', []):
        assert title not in slide_titles(prs), f'Slide should be absent: {title}'

    for spec in EXPECT.get('slide_text', []):
        slide = slide_by_title(prs, spec['slide_title'])
        actual_lines = slide_text_lines(slide)
        if 'exact_lines' in spec:
            assert actual_lines == spec['exact_lines'], (
                f"Unexpected slide text on {spec['slide_title']}: {actual_lines!r}"
            )
        for line in spec.get('required_lines', []):
            assert line in actual_lines, f"Missing line on {spec['slide_title']}: {line}"

    for spec in EXPECT.get('slide_text_anchors', []):
        slide = slide_by_title(prs, spec['slide_title'])
        actual_text = ' '.join(slide_text_lines(slide, skip_title=False))
        assert_contains_all(actual_text, spec.get('required_all', []), f"slide {spec['slide_title']}")

    for spec in EXPECT.get('slide_tables', []):
        slide = slide_by_title(prs, spec['slide_title'])
        table = slide_table(slide)
        for key, expected in spec['cells'].items():
            row_idx, col_idx = [int(x) for x in key.split(',')]
            actual = table.cell(row_idx, col_idx).text.strip()
            assert actual == expected, (
                f"Unexpected table cell on {spec['slide_title']}[{key}]: {actual!r}"
            )

    for spec in EXPECT.get('picture_slides', []):
        slide = slide_by_title(prs, spec['slide_title'])
        assert slide_picture_count(slide) >= spec['min_pictures'], (
            f"Expected at least {spec['min_pictures']} pictures on {spec['slide_title']}"
        )

    for slide_title, expected in EXPECT.get('shape_counts', {}).items():
        slide = slide_by_title(prs, slide_title)
        assert len(slide.shapes) == expected, f'{slide_title}: unexpected shape count {len(slide.shapes)}'

    for spec in EXPECT.get('notes_exact_lines', []):
        slide = slide_by_title(prs, spec['slide_title'])
        actual = ppt_notes_lines(slide)
        assert actual == spec['exact_lines'], f'Unexpected notes on {spec["slide_title"]}: {actual!r}'

    for title in EXPECT.get('reference_slides_match_input', []):
        slide = slide_by_title(prs, title)
        ref = slide_by_title(in_prs, title)
        assert slide_text_lines(slide, skip_title=False) == slide_text_lines(ref, skip_title=False), (
            f'{title}: slide text changed relative to input reference'
        )
        assert ppt_notes_lines(slide) == ppt_notes_lines(ref), (
            f'{title}: slide notes changed relative to input reference'
        )
        assert len(slide.shapes) == len(ref.shapes), (
            f'{title}: shape count changed relative to input reference'
        )
        assert slide_picture_count(slide) == slide_picture_count(ref), (
            f'{title}: picture count changed relative to input reference'
        )

    for spec in EXPECT.get('shape_count_matches_input', []):
        slide = slide_by_title(prs, spec['slide_title'])
        ref = slide_by_title(in_prs, spec['reference_title'])
        assert len(slide.shapes) == len(ref.shapes), (
            f'{spec["slide_title"]}: shape count does not match input reference {spec["reference_title"]}'
        )

    for spec in EXPECT.get('picture_count_matches_input', []):
        slide = slide_by_title(prs, spec['slide_title'])
        ref = slide_by_title(in_prs, spec['reference_title'])
        assert slide_picture_count(slide) == slide_picture_count(ref), (
            f'{spec["slide_title"]}: picture count does not match input reference {spec["reference_title"]}'
        )

    for spec in EXPECT.get('title_signature_matches_input', []):
        slide = slide_by_title(prs, spec['slide_title'])
        ref = slide_by_title(in_prs, spec['reference_title'])
        assert ppt_title_signature(slide) == ppt_title_signature(ref), (
            f'{spec["slide_title"]}: title signature does not match input reference {spec["reference_title"]}'
        )

    for spec in EXPECT.get('body_signature_matches_input', []):
        slide = slide_by_title(prs, spec['slide_title'])
        ref = slide_by_title(in_prs, spec['reference_title'])
        assert ppt_first_body_signature(slide) == ppt_first_body_signature(ref), (
            f'{spec["slide_title"]}: body signature does not match input reference {spec["reference_title"]}'
        )


def test_semantic_verifier():
    run_preflight(META, INPUT_PATH, OUTPUT_PATH)
    if META['doc_type'] == 'excel':
        verify_excel()
    elif META['doc_type'] == 'word':
        verify_word()
    elif META['doc_type'] == 'ppt':
        verify_ppt()
    elif META['doc_type'] == 'pdf':
        verify_pdf()
    else:
        raise AssertionError(f"Unsupported doc type: {META['doc_type']}")
