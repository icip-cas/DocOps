import json
import os
import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from verifier_utils import normalize_text, run_preflight

META_PATH = Path(os.environ.get("TASK_METADATA_PATH", "/tests/task_metadata.json"))
if not META_PATH.exists():
    META_PATH = Path(__file__).parent / "task_metadata.json"
META = json.loads(META_PATH.read_text(encoding="utf-8"))
EXPECT = META["verifier_expectations"]
INPUT_PATH = Path(os.environ.get("INPUT_PATH", META["input_path"]))
OUTPUT_PATH = Path(os.environ.get("OUTPUT_PATH", META["output_path"]))


def rgb_str(rgb):
    return None if rgb is None else str(rgb).upper()


def slide_text(slide, include_notes=True):
    parts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            txt = shape.text.strip()
            if txt:
                parts.append(txt)
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    txt = cell.text.strip()
                    if txt:
                        parts.append(txt)
    if include_notes and getattr(slide, "has_notes_slide", False):
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            parts.append(notes)
    return "\n".join(parts)


def all_text(prs):
    return "\n".join(slide_text(slide, include_notes=True) for slide in prs.slides)


def require_any(text, phrases, label):
    text_norm = normalize_text(text)
    assert any(normalize_text(phrase) in text_norm for phrase in phrases), (
        f"{label}: expected at least one of {phrases!r}"
    )


def require_all_terms(text, phrases, label):
    text_norm = normalize_text(text)
    missing = [phrase for phrase in phrases if normalize_text(phrase) not in text_norm]
    assert not missing, f"{label}: missing required public terms: {missing}"


def slide_title(slide):
    candidates = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = shape.text.strip()
            if text:
                candidates.append((shape.top, shape.left, text.splitlines()[0]))
    assert candidates, "Slide has no title text"
    return sorted(candidates, key=lambda item: (item[0], item[1]))[0][2]


def title_shape(slide):
    title = slide_title(slide)
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = shape.text.strip()
            if text and text.splitlines()[0] == title:
                return shape
    raise AssertionError(f"Could not find title shape for {title!r}")


def table_rows(slide):
    tables = []
    for shape in slide.shapes:
        if getattr(shape, "has_table", False) and shape.has_table:
            rows = []
            for row in shape.table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            tables.append(rows)
    return tables


def background_rgb(slide):
    fill = slide.background.fill
    return rgb_str(fill.fore_color.rgb) if fill.type == MSO_FILL.SOLID else None


def has_accent_shape(slide):
    expected = EXPECT["style"]["accent_fill"]
    for shape in slide.shapes:
        fill = getattr(shape, "fill", None)
        if fill is not None and fill.type == MSO_FILL.SOLID and rgb_str(fill.fore_color.rgb) == expected:
            return True
    return False


def title_font_info(slide):
    shape = title_shape(slide)
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                return run.font.name, bool(run.font.bold), rgb_str(run.font.color.rgb)
    raise AssertionError("Title has no text run")


def speaker_notes(prs):
    return [
        slide.notes_slide.notes_text_frame.text.strip() if getattr(slide, "has_notes_slide", False) else ""
        for slide in prs.slides
    ]


def assert_table_header_style(slide):
    style = EXPECT["style"]
    for shape in slide.shapes:
        if getattr(shape, "has_table", False) and shape.has_table:
            for cell in shape.table.rows[0].cells:
                assert rgb_str(cell.fill.fore_color.rgb) == style["table_header_fill"]
                runs = [
                    run
                    for paragraph in cell.text_frame.paragraphs
                    for run in paragraph.runs
                    if run.text.strip()
                ]
                assert runs, "Table header cell has no text run"
                assert all(run.font.bold for run in runs)
                assert all(rgb_str(run.font.color.rgb) == style["table_header_font"] for run in runs)


def slide_alt_texts(slide):
    descrs = []
    for shape in slide.shapes:
        for node in shape._element.xpath(".//p:cNvPr"):
            descr = node.get("descr")
            if descr:
                descrs.append(descr)
    return descrs


def hyperlink_targets(slide):
    targets = []
    for rel in slide.part.rels.values():
        if rel.reltype == RT.HYPERLINK:
            targets.append(rel.target_ref)
    return targets


def allowed_removed_context(text, start, end):
    window = text[max(0, start - 150): min(len(text), end + 240)]
    return re.search(
        r"\b(no|not|removed|replaced|excluded|cleaned|cleared|without|only public|public)\b",
        window,
        re.IGNORECASE,
    ) is not None


def assert_no_private_or_stale_text(prs):
    raw_text = all_text(prs)
    text = normalize_text(raw_text)

    hard_forbidden = [
        "555-0199",
        "http://internal.example/eoc",
        "eoc.internal",
        "do not publish",
        "unverified closure rumor",
        "medical privacy note",
    ]
    hits = [phrase for phrase in hard_forbidden if normalize_text(phrase) in text]
    assert not hits, f"Forbidden private/source values still present: {hits}"

    stale_patterns = [
        r"16\s+(?:total\s+|public\s+)?shelters",
        r"shelter\s+count[^\n]{0,80}\b16\b",
        r"from\s+16\s+to\s+18",
        r"corrected\s+from\s+16",
        r"previous\s+(?:figure|figures|version)[^\n]{0,80}\b16\b",
        r"11\s+wheelchair",
        r"wheelchair\s+accessible[^\n]{0,100}\b11\b",
        r"from\s+11\s+to\s+14",
        r"corrected\s+from\s+11",
        r"previous\s+(?:figure|figures|version)[^\n]{0,100}\b11\b",
    ]
    stale_hits = [pattern for pattern in stale_patterns if re.search(pattern, raw_text, re.IGNORECASE)]
    assert not stale_hits, f"Stale source correction values remain in public deck: {stale_hits}"

    contextual_private = [
        "DRAFT",
        "PRIVATE",
        "INTERNAL",
        "staff-only",
        "off-slide",
        "personal phone",
        "emergency operations center",
        "EOC",
    ]
    contextual_hits = []
    for phrase in contextual_private:
        flags = 0 if phrase in {"DRAFT", "PRIVATE", "INTERNAL", "EOC"} else re.IGNORECASE
        pattern = re.compile(rf"\b{re.escape(phrase)}\b", flags)
        for match in pattern.finditer(raw_text):
            if not allowed_removed_context(raw_text, match.start(), match.end()):
                contextual_hits.append(phrase)
                break
    assert not contextual_hits, f"Private markers/material still present outside cleanup context: {contextual_hits}"


def offslide_text_shapes(slide, prs):
    offenders = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame and shape.text.strip():
            if shape.left < 0 or shape.top < 0 or shape.left > prs.slide_width or shape.top > prs.slide_height:
                offenders.append(shape.text.strip())
    return offenders


def has_meaningful_alt_text(slide, title):
    title_norm = normalize_text(title)
    footer_norm = normalize_text(EXPECT["style"]["footer_text"])
    for descr in slide_alt_texts(slide):
        descr_norm = normalize_text(descr)
        if len(descr_norm) < 24:
            continue
        if descr_norm in {title_norm, footer_norm, "amber accent bar"}:
            continue
        if descr_norm.startswith("slide title:"):
            continue
        if descr_norm == title_norm:
            continue
        return True
    return False


def test_output_exists_and_is_pptx():
    run_preflight(META, INPUT_PATH, OUTPUT_PATH)
    assert OUTPUT_PATH != INPUT_PATH


def test_slide_order_public_content_tables_and_privacy_cleanup():
    prs = Presentation(OUTPUT_PATH)
    assert [slide_title(slide) for slide in prs.slides] == EXPECT["slide_titles"]
    assert_no_private_or_stale_text(prs)
    by_title = {slide_title(slide): slide for slide in prs.slides}

    cover_text = slide_text(by_title["Emergency Shelter Public Briefing"], include_notes=False)
    require_all_terms(cover_text, ["2026-07-20", "18", "14"], "title slide")
    require_any(cover_text, ["public emergency shelter", "public shelter accessibility", "public distribution"], "title slide")

    network_text = slide_text(by_title["Corrected Shelter Network"], include_notes=False)
    require_all_terms(network_text, ["18", "14"], "corrected shelter network")
    require_any(network_text, ["wheelchair accessible", "accessible shelters"], "corrected shelter network")

    map_text = slide_text(by_title["Accessibility Map"], include_notes=False)
    require_any(map_text, ["map", "zone", "north", "central", "south", "locations"], "accessibility map")
    require_any(map_text, ["14", "accessible", "wheelchair"], "accessibility map")

    transit_text = slide_text(by_title["Transit Access Updates"], include_notes=False)
    require_any(transit_text, ["transit", "bus", "paratransit", "route", "shuttle"], "transit access updates")
    require_any(transit_text, ["accessible", "lift", "drop-off", "parking", "stop"], "transit access updates")

    service_text = slide_text(by_title["Service Animal and Medical Equipment"], include_notes=False)
    require_all_terms(service_text, ["Service animal", "medical equipment"], "service animal and medical equipment")
    require_any(service_text, ["charging", "power", "equipment", "medication"], "service animal and medical equipment")

    link_text = slide_text(by_title["Hotline and Web Links"], include_notes=False)
    require_all_terms(
        link_text,
        ["1-800-555-0148", "https://city.example/shelters/public", "accessibility@city.example"],
        "hotline and web links",
    )

    update_text = slide_text(by_title["Update Log"], include_notes=False)
    require_all_terms(update_text, ["2026-07-20"], "update log")
    require_any(update_text, ["shelter count", "18"], "update log")
    require_any(update_text, ["wheelchair", "14", "accessible"], "update log")

    checklist_text = slide_text(by_title["Accessibility Checklist"], include_notes=False)
    require_any(checklist_text, ["Alt text", "key visuals"], "accessibility checklist")
    require_any(checklist_text, ["Public hyperlinks", "links", "website"], "accessibility checklist")
    require_any(checklist_text, ["Speaker notes", "public notes"], "accessibility checklist")

    for slide in prs.slides:
        offenders = offslide_text_shapes(slide, prs)
        assert not offenders, f"{slide_title(slide)} has off-slide text shapes: {offenders}"


def test_alt_text_and_public_hyperlinks():
    prs = Presentation(OUTPUT_PATH)
    by_title = {slide_title(slide): slide for slide in prs.slides}
    for title in EXPECT["slide_titles"]:
        assert has_meaningful_alt_text(by_title[title], title), f"{title}: missing meaningful alt text"
    all_targets = []
    for slide in prs.slides:
        all_targets.extend(hyperlink_targets(slide))
    for target in EXPECT["hyperlinks"]:
        assert target in all_targets, f"Missing hyperlink target: {target}"
    forbidden_targets = [target for target in all_targets if "internal" in target.lower() or "eoc" in target.lower()]
    assert not forbidden_targets, f"Forbidden hyperlink targets remain: {forbidden_targets}"


def test_speaker_notes_cleaned_and_preserved():
    prs = Presentation(OUTPUT_PATH)
    notes = speaker_notes(prs)
    assert len(notes) == len(EXPECT["slide_titles"])
    for title, note in zip(EXPECT["slide_titles"], notes):
        assert note.strip(), f"{title}: missing public speaker notes"
    assert_no_private_or_stale_text(prs)


def test_style_migration_applied():
    prs = Presentation(OUTPUT_PATH)
    style = EXPECT["style"]
    for idx, slide in enumerate(prs.slides):
        expected_bg = style["title_background"] if idx == 0 else style["content_background"]
        expected_title_color = style["title_color_title_slide"] if idx == 0 else style["title_color_content"]
        assert background_rgb(slide) == expected_bg
        font_name, bold, color = title_font_info(slide)
        assert font_name == style["title_font"]
        assert bold
        assert color == expected_title_color
        assert has_accent_shape(slide)
        assert style["footer_text"] in slide_text(slide, include_notes=False)
        if slide_title(slide) in EXPECT["tables"]:
            assert_table_header_style(slide)


def test_source_artifact_was_not_modified():
    source = Presentation(INPUT_PATH)
    assert len(source.slides) > len(EXPECT["slide_titles"])
    text = normalize_text(all_text(source))
    missing = [p for p in EXPECT["source_must_contain"] if normalize_text(p) not in text]
    assert not missing, f"Source artifact no longer contains expected seed defects: {missing}"
