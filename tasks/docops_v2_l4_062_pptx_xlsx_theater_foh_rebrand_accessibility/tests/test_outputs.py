import json
import os
import re
from datetime import date, datetime
from pathlib import Path

from openpyxl import load_workbook
from openpyxl.utils.cell import get_column_letter, range_boundaries
from pptx import Presentation

from verifier_utils import *  # noqa: F401,F403

META_PATH = Path(os.environ.get("TASK_METADATA_PATH", "/tests/task_metadata.json"))
META = json.loads(META_PATH.read_text(encoding="utf-8"))
EXPECT = META["verifier_expectations"]


def _norm(value):
    if value is None:
        return ""
    if isinstance(value, datetime):
        return value.strftime("%Y-%m-%d")
    if isinstance(value, date):
        return value.strftime("%Y-%m-%d")
    if isinstance(value, float) and value.is_integer():
        value = int(value)
    return re.sub(r"\s+", " ", str(value).replace("\u2013", "-").replace("\u2014", "-").replace("\xa0", " ")).strip()


def _norm_formula(value):
    return re.sub(r"\s+", "", _norm(value)).upper()


def _norm_key(value):
    return re.sub(r"[^a-z0-9]+", "", _norm(value).lower())


def _contains_all(text, tokens):
    norm = normalize_text(text)
    return all(normalize_text(token) in norm for token in tokens)


def _add_rgb(colors, value):
    if value:
        colors.append(str(value).upper())


def _path(kind):
    env = {"pptx": "PPTX_OUTPUT_PATH", "xlsx": "XLSX_OUTPUT_PATH"}[kind]
    key = {"pptx": "pptx_output", "xlsx": "xlsx_output"}[kind]
    return Path(os.environ.get(env, EXPECT[key]))


def _ppt_text_titles_colors(path):
    prs = Presentation(path)
    text, titles, colors = [], [], []
    for slide in prs.slides:
        slide_text = []
        try:
            _add_rgb(colors, slide.background.fill.fore_color.rgb)
        except Exception:
            pass
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                for line in shape.text.splitlines():
                    line = line.strip()
                    if line:
                        text.append(line)
                        slide_text.append(line)
                try:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            _add_rgb(colors, run.font.color.rgb)
                except Exception:
                    pass
            try:
                _add_rgb(colors, shape.fill.fore_color.rgb)
            except Exception:
                pass
            try:
                _add_rgb(colors, shape.line.color.rgb)
            except Exception:
                pass
        titles.append(slide_text[0] if slide_text else "")
    return "\n".join(text), titles, colors, len(prs.slides)


def _rows(ws, start, end, max_col):
    return [[_norm(ws.cell(r, c).value) for c in range(1, max_col + 1)] for r in range(start, end + 1)]


def _expected(rows):
    return [[_norm(c) for c in row] for row in rows]


def _table_ref(ws, expected_name):
    for name in ws.tables.keys():
        if str(name).lower() == str(expected_name).lower():
            return ws.tables[name].ref
    raise AssertionError(f"Missing table {expected_name!r} on {ws.title!r}")


def _table_bounds_for(ws, expected_name):
    return range_boundaries(_table_ref(ws, expected_name).replace("$", ""))


def _table_data_text(ws, expected_name):
    min_col, min_row, max_col, max_row = _table_bounds_for(ws, expected_name)
    parts = []
    for row in range(min_row, max_row + 1):
        for col in range(min_col, max_col + 1):
            if ws.cell(row, col).value is not None:
                parts.append(str(ws.cell(row, col).value))
    return "\n".join(parts)


def _table_rows(ws, expected_name):
    min_col, min_row, max_col, max_row = _table_bounds_for(ws, expected_name)
    return [
        [_norm(ws.cell(row, col).value) for col in range(min_col, max_col + 1)]
        for row in range(min_row + 1, max_row + 1)
    ]


def _find_table_row(rows, tokens):
    for row in rows:
        text = " | ".join(row)
        if all(_norm_key(token) in _norm_key(text) for token in tokens):
            return row
    raise AssertionError(f"no table row contains {tokens!r}; rows={rows!r}")


def _range_covered(ws, expected):
    min_col, min_row, max_col, max_row = range_boundaries(expected.replace("$", ""))
    target = {f"{get_column_letter(col)}{row}" for row in range(min_row, max_row + 1) for col in range(min_col, max_col + 1)}
    cells = set()
    for dv in ws.data_validations.dataValidation:
        for rng in dv.cells.ranges:
            min_c, min_r, max_c, max_r = range_boundaries(str(rng).replace("$", ""))
            cells.update(f"{get_column_letter(col)}{row}" for row in range(min_r, max_r + 1) for col in range(min_c, max_c + 1))
    return target.issubset(cells)


def _clean_area(value):
    text = _norm(value).replace("'", "").replace("$", "").replace(" ", "")
    if "!" in text:
        text = text.split("!", 1)[1]
    return text.upper()


def _area_covers(actual, target):
    actual = _clean_area(actual)
    target = _clean_area(target)
    if target in actual or actual in target:
        return True
    actual_min_col, actual_min_row, actual_max_col, actual_max_row = range_boundaries(actual)
    target_min_col, target_min_row, target_max_col, target_max_row = range_boundaries(target)
    return (
        actual_min_col <= target_min_col
        and actual_min_row <= target_min_row
        and actual_max_col >= target_max_col
        and actual_max_row >= target_max_row
    )


def _sheet_text(wb, sheets):
    _assert_required_sheets(wb, sheets)
    parts = []
    for sheet in sheets:
        for row in wb[sheet].iter_rows():
            for cell in row:
                if cell.value is not None:
                    parts.append(str(cell.value))
    return "\n".join(parts)


def _assert_required_sheets(wb, sheets):
    missing = [sheet for sheet in sheets if sheet not in wb.sheetnames]
    assert not missing, f"Missing required workbook sheets: {missing}; found {wb.sheetnames}"


def _readiness_formula_ok(value):
    formula = _norm_formula(value)
    assert formula.startswith("="), "deck-ready control must be a formula"
    if "DECKREADY" in formula:
        return
    for token in ("IF", "AND", "READY", "HOLD"):
        assert token in formula, f"deck-ready formula missing {token}"
    assert ("B6=4" in formula or "ACCESSIBILITYCOUNT=4" in formula), "deck-ready formula must check four accessibility services"
    assert ("B7=2" in formula or "ADVISORYCOUNT=2" in formula), "deck-ready formula must check two advisories"
    assert "B8=4" in formula, "deck-ready formula must check four excluded private notes"


def test_outputs_exist():
    assert _path("pptx").exists()
    assert _path("xlsx").exists()
    assert _path("pptx").suffix.lower() == ".pptx"
    assert _path("xlsx").suffix.lower() == ".xlsx"


def test_pptx_retrieval_style_and_privacy():
    text, titles, colors, count = _ppt_text_titles_colors(_path("pptx"))
    assert count == 7
    assert titles == EXPECT["slide_titles"]
    required_checks = [
        ["Aurora Lines"],
        ["North Pier Theater"],
        ["Doors", "6:30 PM"],
        ["Curtain", "7:30 PM"],
        ["Act I", "52 minutes"],
        ["Intermission", "18 minutes"],
        ["Act II", "47 minutes"],
        ["Door 3"],
        ["Rows B-D"],
        ["Captioning screens"],
        ["Audio description", "Coat Check"],
        ["Quiet Room", "Studio C"],
        ["strobe", "Scene 5"],
        ["photography", "recording"],
        ["meet-and-greet", "cancelled"],
        ["Merch", "10:15 PM"],
    ]
    missing = [tokens for tokens in required_checks if not _contains_all(text, tokens)]
    assert not missing, f"signage deck: missing required semantics {missing!r}"
    forbid_any(text, EXPECT["ppt_forbidden"], "signage deck")
    for color in EXPECT["required_colors"]:
        assert color in colors, f"required brand color {color} missing; found {colors}"
    for color in EXPECT["forbidden_colors"]:
        assert color not in colors, f"legacy forbidden color {color} still present"


def test_workbook_structure_formulas_controls():
    wb = load_workbook(_path("xlsx"), data_only=False)
    assert wb.sheetnames == EXPECT["sheet_order"]
    for sheet, (table_name, ref) in EXPECT["tables"].items():
        min_col, min_row, max_col, max_row = range_boundaries(ref)
        actual_min_col, actual_min_row, actual_max_col, actual_max_row = _table_bounds_for(wb[sheet], table_name)
        assert actual_max_col - actual_min_col + 1 >= max_col - min_col + 1
        assert actual_max_row - actual_min_row + 1 >= max_row - min_row + 1
    for sheet in EXPECT["hidden_sheets"]:
        assert wb[sheet].sheet_state in ("hidden", "veryHidden")
    names = {dn.name for dn in wb.defined_names.values()}
    for name in EXPECT["defined_names"]:
        assert name in names
    deck_ready_row = _find_table_row(_table_rows(wb["Deck Summary"], "tblDeckSummary"), ["Deck", "Ready"])
    _readiness_formula_ok(deck_ready_row[1])
    for sheet, ranges in EXPECT["data_validation"].items():
        assert wb[sheet].data_validations.dataValidation, f"{sheet}: missing data validation"
    for sheet, area in EXPECT["print_areas"].items():
        target = _clean_area(_table_ref(wb[sheet], EXPECT["tables"][sheet][0]))
        assert _area_covers(wb[sheet].print_area, target)


def test_workbook_values_and_privacy():
    wb = load_workbook(_path("xlsx"), data_only=False)
    _assert_required_sheets(wb, ["Deck Summary", "Slide Source Trace", "Accessibility Map", "Advisories", "Style QA", "Private Notes"])
    summary = _table_data_text(wb["Deck Summary"], "tblDeckSummary")
    for tokens in [["Aurora Lines"], ["North Pier Theater"], ["6:30 PM"], ["7:30 PM"], ["Accessibility", "4"], ["Advisory", "2"], ["Private", "4"]]:
        assert _contains_all(summary, tokens), f"Deck Summary missing {tokens!r}"
    slide_trace = _table_data_text(wb["Slide Source Trace"], "tblSlideTrace")
    for title in EXPECT["slide_titles"]:
        assert normalize_text(title) in normalize_text(slide_trace), f"Slide Source Trace missing {title}"
    access = _table_data_text(wb["Accessibility Map"], "tblAccessibility")
    for tokens in [["Wheelchair", "Door 3", "Rows"], ["Captioning"], ["Audio", "Coat Check"], ["Quiet", "Studio C"]]:
        assert _contains_all(access, tokens), f"Accessibility Map missing {tokens!r}"
    advisories = _table_data_text(wb["Advisories"], "tblAdvisories")
    for tokens in [["Strobe", "Scene 5"], ["Photography", "Recording"]]:
        assert _contains_all(advisories, tokens), f"Advisories missing {tokens!r}"
    style = _table_data_text(wb["Style QA"], "tblStyleQA")
    for tokens in [["102A43", "Pass"], ["F2C94C", "Pass"], ["F7F1E1", "Pass"], ["red", "Pass"], ["placeholder", "Pass"]]:
        assert _contains_all(style, tokens), f"Style QA missing {tokens!r}"
    private = _table_data_text(wb["Private Notes"], "tblPrivateNotes")
    for tokens in [["VIP donor", "Excluded"], ["champagne", "Excluded"], ["artist hotel", "Excluded"], ["staffing shortage", "Excluded"]]:
        assert _contains_all(private, tokens), f"Private Notes missing {tokens!r}"
    public = _sheet_text(wb, ["Deck Summary", "Slide Source Trace", "Accessibility Map", "Advisories", "Style QA"])
    forbid_any(public, EXPECT["ppt_forbidden"], "public workbook sheets")


def test_cross_output_consistency():
    deck_text, _, _, _ = _ppt_text_titles_colors(_path("pptx"))
    wb = load_workbook(_path("xlsx"), data_only=False)
    public = _sheet_text(wb, ["Deck Summary", "Slide Source Trace", "Accessibility Map", "Advisories", "Style QA"])
    for tokens in [["Aurora Lines"], ["Door 3"], ["Rows B-D"], ["Captioning"], ["Audio description"], ["Quiet", "Studio C"], ["strobe", "Scene 5"]]:
        assert _contains_all(deck_text, tokens), f"deck missing cross-output anchor {tokens!r}"
        assert _contains_all(public, tokens), f"workbook missing cross-output anchor {tokens!r}"
    deck_ready_row = _find_table_row(_table_rows(wb["Deck Summary"], "tblDeckSummary"), ["Deck", "Ready"])
    _readiness_formula_ok(deck_ready_row[1])
