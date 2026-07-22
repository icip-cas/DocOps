import json
import os
import re
import sys
import zipfile
from pathlib import Path
from datetime import date, datetime

from docx import Document
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

sys.path.insert(0, str(Path(__file__).parent))
from verifier_utils import *  # noqa: F401,F403

META_PATH = Path(os.environ.get('TASK_METADATA_PATH', '/tests/task_metadata.json'))
META = json.loads(META_PATH.read_text(encoding='utf-8'))
EXPECT = META['verifier_expectations']


def _vopt_norm_text(value):
    if value is None:
        return ''
    if isinstance(value, datetime):
        if value.hour or value.minute or value.second:
            return value.strftime('%Y-%m-%d %H:%M')
        return value.strftime('%Y-%m-%d')
    if isinstance(value, date):
        return value.strftime('%Y-%m-%d')
    if isinstance(value, float) and value.is_integer():
        value = int(value)
    text = str(value)
    text = text.replace('\u2013', '-').replace('\u2014', '-').replace('\u2212', '-')
    text = text.replace('\xa0', ' ')
    text = re.sub(r'\s+', ' ', text).strip()
    return text


def _vopt_norm_formula(value):
    return re.sub(r'\s+', '', _vopt_norm_text(value)).upper()


def _vopt_assert_formula(actual, expected, label='formula'):
    assert _vopt_norm_formula(actual) == _vopt_norm_formula(expected), f"{label}: expected {expected!r}, found {actual!r}"


def _vopt_norm_rows(rows):
    return [[_vopt_norm_text(cell) for cell in row] for row in rows]


def _vopt_rows_equal(actual, expected):
    return _vopt_norm_rows(actual) == _vopt_norm_rows(expected)


def _vopt_rows_text(rows):
    return '\n'.join('|'.join(_vopt_norm_text(cell) for cell in row) for row in rows)


def _vopt_header_has(header, required):
    actual = {_vopt_norm_text(cell).lower() for cell in header}
    missing = [cell for cell in required if _vopt_norm_text(cell).lower() not in actual]
    assert not missing, f"Missing expected header columns {missing!r}; actual={header!r}"


def _vopt_section_text(section_part):
    return '\n'.join(p.text for p in getattr(section_part, 'paragraphs', []) if p.text is not None)


def _vopt_header_text(doc):
    return '\n'.join(_vopt_section_text(section.header) for section in doc.sections)


def _vopt_footer_text(doc):
    return '\n'.join(_vopt_section_text(section.footer) for section in doc.sections)


def _vopt_table_ref(ws, expected_name=None):
    tables = ws.tables
    if expected_name:
        for name in tables.keys():
            if str(name).lower() == str(expected_name).lower():
                return tables[name].ref
    vals = list(tables.values())
    assert vals, f"No Excel native table found on sheet {ws.title!r}"
    return vals[0].ref


def _vopt_clean_area(value):
    text = _vopt_norm_text(value).replace("'", '').replace('$', '').replace(' ', '')
    if '!' in text:
        text = text.split('!', 1)[1]
    return text.upper()


def _vopt_assert_print_area(ws, expected):
    actual = _vopt_clean_area(ws.print_area)
    target = _vopt_clean_area(expected)
    assert target in actual or actual in target, f"{ws.title}: expected print area {expected!r}, found {ws.print_area!r}"


def _vopt_freeze_pane(value):
    return getattr(value, 'coordinate', value)


def _vopt_validated_cells(ws):
    cells = set()
    for dv in ws.data_validations.dataValidation:
        for rng in dv.cells.ranges:
            text = str(rng).replace('$', '')
            try:
                min_col, min_row, max_col, max_row = range_boundaries(text)
            except ValueError:
                cells.add(text)
                continue
            if max_row > 10000:
                max_row = max(min_row, 200)
            for row in range(min_row, max_row + 1):
                for col in range(min_col, max_col + 1):
                    cells.add(f"{get_column_letter(col)}{row}")
    return cells


class _VoptValidationRanges(list):
    def __init__(self, ws, ranges):
        super().__init__(ranges)
        self.ws = ws

    def __contains__(self, expected):
        expected = str(expected).replace('$', '')
        if super().__contains__(expected):
            return True
        try:
            min_col, min_row, max_col, max_row = range_boundaries(expected)
        except ValueError:
            return super().__contains__(expected)
        target = {f"{get_column_letter(col)}{row}" for row in range(min_row, max_row + 1) for col in range(min_col, max_col + 1)}
        return target.issubset(_vopt_validated_cells(self.ws))


def _vopt_dv_ranges(ws):
    ranges = []
    for dv in ws.data_validations.dataValidation:
        ranges.extend(str(rng).replace('$', '') for rng in dv.cells.ranges)
    return _VoptValidationRanges(ws, ranges)


def _vopt_ordered_subset(expected, actual):
    actual_iter = iter([_vopt_norm_text(item) for item in actual])
    for item in [_vopt_norm_text(value) for value in expected]:
        for candidate in actual_iter:
            if item == candidate:
                break
        else:
            return False
    return True



def _resolve_output(kind):
    if kind == 'deck':
        return Path(os.environ.get('PPT_OUTPUT_PATH', EXPECT['deck_output']))
    if kind == 'memo':
        return Path(os.environ.get('DOCX_OUTPUT_PATH', EXPECT['memo_output']))
    raise KeyError(kind)


def _slide_title(slide):
    for shape in slide.shapes:
        if not getattr(shape, 'is_placeholder', False):
            continue
        try:
            if shape.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                text = ' '.join(shape.text.split())
                if text:
                    return text
        except Exception:
            continue
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            text = ' '.join(shape.text.split())
            if text and normalize_text(text) not in ('draft', 'final'):
                return text
    return ''


def _slide_titles(prs):
    return [_slide_title(slide) for slide in prs.slides]


def _slide_by_title(prs, title):
    for slide in prs.slides:
        if normalize_text(title) in normalize_text(_slide_text(slide)):
            return slide
    raise AssertionError(f'Slide not found: {title}')


def _slide_text(slide):
    parts = []
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            text = ' '.join(shape.text.split())
            if text:
                parts.append(text)
    return '\n'.join(parts)


def _first_table(slide):
    for shape in slide.shapes:
        if getattr(shape, 'has_table', False):
            return shape.table
    raise AssertionError('Expected native table')


def _table_rows(table):
    return [[cell.text.strip() for cell in row.cells] for row in table.rows]


def _presentation_text(prs):
    parts = []
    for slide in prs.slides:
        parts.append(_slide_text(slide))
        for shape in slide.shapes:
            if getattr(shape, 'has_table', False):
                parts.append(_vopt_rows_text(_table_rows(shape.table)))
    return '\n'.join(part for part in parts if part)


def _doc_heading_order(doc):
    out = []
    for p in doc.paragraphs:
        style = p.style.name if p.style else ''
        if p.text.strip() and (style.startswith('Heading') or style == 'Title'):
            out.append(p.text.strip())
    return out


def _norm_heading(text):
    text = _vopt_norm_text(text)
    text = re.sub(r'^\d+\.\s*', '', text)
    text = text.replace('Incident Overview', 'Executive Summary')
    text = text.replace('Incident Summary', 'Executive Summary')
    text = text.replace('Incident Timeline', 'Timeline')
    text = text.replace('Decisions', 'Decision Summary')
    text = text.replace('Incident Review: Executive Memo', 'Incident Review Executive Memo')
    text = text.replace('Incident Review - Executive Memo', 'Incident Review Executive Memo')
    text = text.replace('Incident Review — Executive Memo', 'Incident Review Executive Memo')
    text = text.replace('Incident Review – Executive Memo', 'Incident Review Executive Memo')
    return text


def _memo_heading_sequence(doc):
    headings = [_norm_heading(h) for h in _doc_heading_order(doc)]
    headings = [h for h in headings if normalize_text(h) != 'table of contents']
    first_text = next((p.text.strip() for p in doc.paragraphs if p.text.strip()), '')
    if first_text and normalize_text('Incident Review Executive Memo') in normalize_text(_norm_heading(first_text)):
        if not headings or headings[0] != 'Incident Review Executive Memo':
            headings.insert(0, 'Incident Review Executive Memo')
    return headings


def _assert_ordered_headings(doc, expected):
    actual = _memo_heading_sequence(doc)
    assert _vopt_ordered_subset(expected, actual), f"Unexpected memo headings: {actual!r}"


def _doc_text(doc):
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append('|'.join(cell.text.strip() for cell in row.cells))
    return '\n'.join(parts)


def _safe_section_part_text(doc, part_name):
    texts = []
    for idx, section in enumerate(doc.sections):
        try:
            part = getattr(section, part_name)
            texts.append(_vopt_section_text(part))
        except Exception as exc:
            texts.append(f'__{part_name.upper()}_ERROR_SECTION_{idx + 1}: {exc}__')
    return '\n'.join(texts)


def _package_part_text(path, prefix):
    texts = []
    try:
        with zipfile.ZipFile(path) as zf:
            for name in zf.namelist():
                if not name.startswith(prefix) or not name.endswith('.xml'):
                    continue
                xml = zf.read(name).decode('utf-8', errors='ignore')
                texts.extend(re.findall(r'<w:t[^>]*>(.*?)</w:t>', xml))
                texts.extend(re.findall(r'<w:instrText[^>]*>(.*?)</w:instrText>', xml))
    except Exception:
        pass
    return '\n'.join(texts)


def _has_toc_field(path):
    with zipfile.ZipFile(path) as zf:
        xml = zf.read('word/document.xml').decode('utf-8', errors='ignore')
    return bool(re.search(r'TOC\s+(?:\\)?o|TOC\s*(?:&quot;|")', xml))


def _customer_message_later_section_or_page(path, doc):
    if len(doc.sections) >= 2:
        return True
    try:
        with zipfile.ZipFile(path) as zf:
            xml = zf.read('word/document.xml').decode('utf-8', errors='ignore')
    except Exception:
        return False
    idx = xml.lower().find('customer message')
    if idx < 0:
        idx = xml.lower().find('customer')
    if idx < 0:
        return False
    before = xml[:idx]
    return 'w:type="page"' in before or "w:type='page'" in before or '<w:br' in before and 'page' in before


def _assert_release_facts(text, label):
    norm = normalize_text(text)
    groups = [
        ('severity path', ['s1', '09:42', 's2']),
        ('external incident decision', ['d-14', 'declare external incident', 'approved']),
        ('customer faq decision', ['d-15', 'customer faq', 'luis romero']),
        ('failover/recovery', ['payment', 'failover', 'checkout latency']),
        ('pricing-launch exclusion', ['pricing launch email', 'delay']),
    ]
    for group_label, tokens in groups:
        missing = [token for token in tokens if token not in norm]
        assert not missing, f'{label}: missing {group_label} tokens {missing!r}'
    assert 'd-16' not in norm or ('rejected' in norm or 'excluded' in norm or 'do not mention' in norm), (
        f'{label}: D-16 may only appear as rejected/excluded context'
    )


def _assert_message_facts(text, label):
    norm = normalize_text(text)
    for token in ['checkout latency', 'mitigated', 'payment failover', 'faq owner', 'luis romero']:
        assert token in norm, f'{label}: missing message token {token!r}'
    assert 'pricing launch email' in norm and ('do not mention' in norm or 'excluded' in norm or 'rejected' in norm), (
        f'{label}: missing pricing-launch exclusion note'
    )


def _rows_equivalent(actual, expected):
    actual_norm = _vopt_norm_rows(actual)
    expected_norm = _vopt_norm_rows(expected)
    if actual_norm == expected_norm:
        return True
    if len(actual_norm) != len(expected_norm):
        return False
    for actual_row, expected_row in zip(actual_norm, expected_norm):
        if len(actual_row) != len(expected_row):
            return False
        for idx, (actual_cell, expected_cell) in enumerate(zip(actual_row, expected_row)):
            if actual_cell == expected_cell:
                continue
            if idx == 3 and actual_row and actual_row[0] == 'D-14':
                if normalize_text('S1') in normalize_text(actual_cell) and '09:42' in actual_cell and normalize_text('S2') in normalize_text(actual_cell):
                    continue
            if idx == 3 and actual_row and actual_row[0] == 'D-15':
                if 'Luis Romero' in actual_cell:
                    continue
            return False
    return True


def _find_table_by_header(tables, expected_header):
    expected = [_vopt_norm_text(cell).lower() for cell in expected_header]
    for table in tables:
        if not table:
            continue
        header = [_vopt_norm_text(cell).lower() for cell in table[0]]
        if header == expected:
            return table
    raise AssertionError(f'Missing table with header {expected_header!r}')


def test_outputs_exist():
    deck = _resolve_output('deck')
    memo = _resolve_output('memo')
    assert deck.exists(), f'Missing deck output: {deck}'
    assert memo.exists(), f'Missing memo output: {memo}'
    assert deck.suffix.lower() == '.pptx'
    assert memo.suffix.lower() == '.docx'


def test_deck_content():
    prs = Presentation(_resolve_output('deck'))
    expected_titles = [title for title in EXPECT['titles_order'] if title != 'Executive Summary']
    assert _vopt_ordered_subset(expected_titles, _slide_titles(prs)), f"Unexpected slide order: {_slide_titles(prs)!r}"
    for title in EXPECT['absent_slide_titles']:
        assert title not in _slide_titles(prs), f'Obsolete slide still present: {title}'
    deck_text = _presentation_text(prs)
    _assert_release_facts(deck_text, 'deck summary')
    timeline = _table_rows(_first_table(_slide_by_title(prs, 'Timeline')))[1:]
    assert _vopt_rows_equal(timeline, EXPECT['timeline_rows']), f'Unexpected deck timeline rows: {timeline!r}'
    decisions = _table_rows(_first_table(_slide_by_title(prs, 'Decision Summary')))[1:]
    assert _rows_equivalent(decisions, EXPECT['decision_rows']), f'Unexpected deck decision rows: {decisions!r}'
    actions = _table_rows(_first_table(_slide_by_title(prs, 'Action Plan')))[1:]
    assert _vopt_rows_equal(actions, EXPECT['action_rows']), f'Unexpected deck action rows: {actions!r}'
    message = _slide_text(_slide_by_title(prs, 'Customer Message'))
    _assert_message_facts(message, 'deck message')
    forbid_any('\n'.join(_slide_text(slide) for slide in prs.slides), EXPECT['forbidden'], 'deck text')
    assert _table_rows(_first_table(_slide_by_title(prs, 'Reference - Action Table')))[0][0] == 'Action ID'
    require_all(_slide_text(_slide_by_title(prs, 'Reference - Message Layout')), ['LOCKED MESSAGE LAYOUT'], 'message reference')


def test_memo_content_and_structure():
    memo_path = _resolve_output('memo')
    doc = Document(memo_path)
    _assert_ordered_headings(doc, EXPECT['memo_heading_order'])
    text = _doc_text(doc)
    _assert_release_facts(text, 'memo text')
    _assert_message_facts(text, 'memo text')
    forbid_any(text, EXPECT['forbidden'], 'memo text')
    assert _has_toc_field(memo_path), 'Expected memo TOC field'
    assert _customer_message_later_section_or_page(memo_path, doc), 'Expected Customer Message to start in a later section/page'
    header_text = _safe_section_part_text(doc, 'header') + '\n' + _package_part_text(memo_path, 'word/header')
    footer_text = _safe_section_part_text(doc, 'footer') + '\n' + _package_part_text(memo_path, 'word/footer')
    header_norm = normalize_text(header_text)
    assert 'incident review' in header_norm and ('release' in header_norm or 'memo' in header_norm), (
        f'Expected incident-review release header, found {header_text!r}'
    )
    assert 'final' in normalize_text(footer_text), f'Expected final status in footer, found {footer_text!r}'
    tables = [_table_rows(table) for table in doc.tables]
    timeline = _find_table_by_header(tables, ['Time', 'Event', 'Phase', 'Severity'])
    decision = _find_table_by_header(tables, ['Decision ID', 'Decision', 'Outcome', 'Note'])
    action = _find_table_by_header(tables, ['Action ID', 'Action', 'Owner', 'Due Date', 'Status'])
    assert _vopt_rows_equal(timeline[1:], EXPECT['timeline_rows'])
    assert _rows_equivalent(decision[1:], EXPECT['decision_rows'])
    assert _vopt_rows_equal(action[1:], EXPECT['action_rows'])


def test_cross_output_consistency():
    prs = Presentation(_resolve_output('deck'))
    deck_timeline = _table_rows(_first_table(_slide_by_title(prs, 'Timeline')))[1:]
    deck_actions = _table_rows(_first_table(_slide_by_title(prs, 'Action Plan')))[1:]
    doc = Document(_resolve_output('memo'))
    tables = [_table_rows(table) for table in doc.tables]
    memo_timeline = _find_table_by_header(tables, ['Time', 'Event', 'Phase', 'Severity'])[1:]
    memo_actions = _find_table_by_header(tables, ['Action ID', 'Action', 'Owner', 'Due Date', 'Status'])[1:]
    assert deck_timeline == memo_timeline, 'Deck and memo timeline rows diverge'
    assert deck_actions == memo_actions, 'Deck and memo action rows diverge'
