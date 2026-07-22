import json
import os
import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.dml import MSO_FILL

from verifier_utils import normalize_text, run_preflight

META_PATH = Path(os.environ.get("TASK_METADATA_PATH", "/tests/task_metadata.json"))
if not META_PATH.exists():
    META_PATH = Path(__file__).parent / "task_metadata.json"
META = json.loads(META_PATH.read_text(encoding="utf-8"))
EXPECT = META["verifier_expectations"]
INPUT_PATH = Path(os.environ.get("INPUT_PATH", META["input_path"]))
OUTPUT_PATH = Path(os.environ.get("OUTPUT_PATH", META["output_path"]))


def rgb_str(rgb):
    if rgb is None:
        return None
    return str(rgb).upper()


def slide_text(slide, include_notes=True):
    parts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            txt = shape.text.strip()
            if txt:
                parts.append(txt)
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    txt = cell.text.strip()
                    if txt:
                        parts.append(txt)
    if include_notes and getattr(slide, "has_notes_slide", False):
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            parts.append(notes)
    return "\n".join(parts)


def norm_key(text):
    return re.sub(r"[^a-z0-9+$%]+", " ", str(text).lower()).strip()


def compact_key(text):
    return re.sub(r"[^a-z0-9]+", "", str(text).lower())


def all_text(prs):
    return "\n".join(slide_text(slide, include_notes=True) for slide in prs.slides)


def slide_contains_title(slide, expected_title):
    return compact_key(expected_title) in compact_key(slide_text(slide, include_notes=False))


def slide_title(slide, expected_title=None):
    if expected_title and slide_contains_title(slide, expected_title):
        return expected_title
    candidates = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = shape.text.strip()
            if text:
                candidates.append((shape.top, shape.left, text.splitlines()[0]))
    assert candidates, "Slide has no title text"
    return sorted(candidates, key=lambda item: (item[0], item[1]))[0][2]


def title_shape(slide, expected_title=None):
    title = slide_title(slide, expected_title)
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = shape.text.strip()
            if not text:
                continue
            if expected_title and compact_key(expected_title) in compact_key(text):
                return shape
            if expected_title:
                title_parts = [part for part in re.split(r"[:|]", expected_title) if compact_key(part)]
                if any(compact_key(part) in compact_key(text) for part in title_parts):
                    return shape
            if text.splitlines()[0] == title:
                return shape
    raise AssertionError(f"Could not find title shape for {title!r}")


def table_rows(slide):
    tables = []
    for shape in slide.shapes:
        if getattr(shape, "has_table", False) and shape.has_table:
            rows = []
            for row in shape.table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            tables.append(rows)
    return tables


def background_rgb(slide):
    fill = slide.background.fill
    if fill.type == MSO_FILL.SOLID:
        return rgb_str(fill.fore_color.rgb)
    return None


def has_accent_shape(slide):
    expected = EXPECT["style"]["accent_fill"]
    for shape in slide.shapes:
        fill = getattr(shape, "fill", None)
        if fill is None or fill.type != MSO_FILL.SOLID:
            continue
        if rgb_str(fill.fore_color.rgb) == expected:
            return True
    return False


def title_font_info(slide, expected_title=None):
    shape = title_shape(slide, expected_title)
    xml = shape._element.xml
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                try:
                    color = rgb_str(run.font.color.rgb)
                except Exception:
                    color = None
                return run.font.name, bool(run.font.bold), color, xml
    raise AssertionError(f"Title shape {shape.name!r} has no text run")


def speaker_notes(prs):
    notes = []
    for slide in prs.slides:
        if getattr(slide, "has_notes_slide", False):
            notes.append(slide.notes_slide.notes_text_frame.text.strip())
        else:
            notes.append("")
    return notes


def test_output_exists_and_is_pptx():
    run_preflight(META, INPUT_PATH, OUTPUT_PATH)
    assert OUTPUT_PATH != INPUT_PATH


def test_slide_order_and_private_material_removed():
    prs = Presentation(OUTPUT_PATH)
    assert len(prs.slides) == len(EXPECT["slide_titles"])
    for slide, expected_title in zip(prs.slides, EXPECT["slide_titles"]):
        assert slide_contains_title(slide, expected_title), f"Missing slide title {expected_title!r}"
    text = normalize_text(all_text(prs))
    hits = [p for p in EXPECT["forbidden_phrases"] if normalize_text(p) in text]
    assert not hits, f"Forbidden phrases still present: {hits}"


def require_parts(text, parts, label):
    norm = norm_key(text)
    missing = [part for part in parts if norm_key(part) not in norm]
    assert not missing, f"{label}: missing semantic parts {missing!r}"


def require_any_parts(text, part_groups, label):
    failures = []
    for parts in part_groups:
        try:
            require_parts(text, parts, label)
            return
        except AssertionError as exc:
            failures.append(str(exc))
    raise AssertionError(f"{label}: none of the accepted semantic variants matched: {failures!r}")


def test_required_public_content_and_corrections():
    prs = Presentation(OUTPUT_PATH)
    by_title = dict(zip(EXPECT["slide_titles"], prs.slides))
    title_text = slide_text(by_title["HarborLink Fare Pilot: Board Decision Packet"], include_notes=False)
    require_parts(title_text, ["Public board packet", "July 2026", "Prepared for the Mobility Committee"], "title slide")

    decision_text = slide_text(by_title["Decision Snapshot"], include_notes=False)
    require_parts(decision_text, ["recommended vote", "authorize", "12-month", "pilot"], "Decision Snapshot vote")
    require_parts(decision_text, ["East Bay Connector"], "Decision Snapshot corridor")
    require_any_parts(decision_text, [["September", "14", "2026"], ["Sep", "14", "2026"]], "Decision Snapshot launch")
    require_parts(decision_text, ["+9%", "ridership lift", "no increase", "unresolved complaints"], "Decision Snapshot success metric")

    action_text = slide_text(by_title["Board Actions"], include_notes=False)
    action_checks = [
        ["authorize", "12-month", "fare capping pilot", "East Bay Connector"],
        ["approve", "$2.96M", "public budget envelope"],
        ["direct", "monthly public dashboard", "October 30", "2026"],
        ["require", "equity", "accessibility checkpoint", "full rollout"],
    ]
    for idx, parts in enumerate(action_checks, start=1):
        require_parts(action_text, parts, f"Board Actions item {idx}")


def test_public_tables_rebuilt_exactly():
    prs = Presentation(OUTPUT_PATH)
    by_title = dict(zip(EXPECT["slide_titles"], prs.slides))
    checks = {
        "Pilot Corridor Baseline": [
            ["Pilot corridor", "East Bay Connector"],
            ["Baseline weekday ridership", "18,400", "average boardings"],
            ["Projected pilot ridership lift", "+9%"],
            ["Low-income pass enrollment", "12,750", "active riders"],
        ],
        "Equity Safeguards": [
            ["Zone A transfer credit", "Two-hour free transfer window", "Fare Policy"],
            ["Zone B subsidy cap", "$16", "monthly cap", "Equity Office"],
            ["Accessibility review", "Two ADA field audits before launch", "Customer Experience"],
            ["Community clinics", "Six multilingual enrollment clinics", "Outreach"],
        ],
        "Implementation Timeline": [
            ["Board authorization", "July 28", "2026"],
            ["Procurement gate", "May 22", "2026"],
            ["Operator training", "August 17", "2026"],
            ["Pilot launch", "September 14", "2026"],
            ["First public dashboard", "October 30", "2026"],
        ],
        "Budget Controls": [
            ["Fare capping technology", "$1.8M", "Fixed-price task order"],
            ["Enrollment support", "$620K", "Clinic staffing cap"],
            ["Accessibility testing", "$240K", "Independent QA signoff"],
            ["Contingency reserve", "$300K", "Board approval required to release"],
        ],
    }
    for title, groups in checks.items():
        text = slide_text(by_title[title], include_notes=False)
        for group in groups:
            if title == "Budget Controls" and group[0] == "Fare capping technology":
                require_any_parts(text, [group, ["Fare capping technology", "$1.80M", "Fixed-price task order"]], title)
            else:
                require_parts(text, group, title)


def test_speaker_notes_cleaned_and_preserved():
    prs = Presentation(OUTPUT_PATH)
    notes = speaker_notes(prs)
    assert len(notes) == len(EXPECT["slide_titles"])
    for idx, note in enumerate(notes, start=1):
        assert note.strip(), f"Slide {idx}: missing speaker notes"
    notes_text = normalize_text("\n".join(notes))
    hits = [p for p in EXPECT["forbidden_phrases"] if normalize_text(p) in notes_text]
    assert not hits, f"Speaker notes still contain forbidden phrases: {hits}"


def test_style_migration_applied():
    prs = Presentation(OUTPUT_PATH)
    style = EXPECT["style"]
    for idx, (slide, expected_title) in enumerate(zip(prs.slides, EXPECT["slide_titles"])):
        if idx == 0:
            assert background_rgb(slide) == style["title_background"]
            expected_title_color = style["title_color_title_slide"]
        else:
            assert background_rgb(slide) == style["content_background"]
            expected_title_color = style["title_color_content"]
        font_name, bold, color, title_xml = title_font_info(slide, expected_title)
        assert font_name == style["title_font"] or f'typeface="{style["title_font"]}"' in title_xml, f"Slide {idx + 1}: wrong title font"
        assert bold or ' b="1"' in title_xml or ' b="true"' in title_xml, f"Slide {idx + 1}: title should be bold"
        assert color == expected_title_color or expected_title_color in title_xml, f"Slide {idx + 1}: wrong title color"
        assert has_accent_shape(slide), f"Slide {idx + 1}: missing accent shape"
        assert style["footer_text"] in slide_text(slide, include_notes=False)


def test_source_artifact_was_not_modified():
    source = Presentation(INPUT_PATH)
    assert len(source.slides) > len(EXPECT["slide_titles"])
    text = normalize_text(all_text(source))
    assert "draft" in text
    assert "vendor dispute" in text
    assert "personal phone" in text
