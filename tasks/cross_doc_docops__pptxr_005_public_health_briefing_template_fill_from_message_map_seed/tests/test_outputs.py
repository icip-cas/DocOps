import json
import re
import sys
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

sys.path.insert(0, str(Path(__file__).parent))
from verifier_utils import *  # noqa: F401,F403

META = json.loads(Path('/tests/task_metadata.json').read_text(encoding='utf-8'))
INPUT_PATH = Path(META['input_path'])
OUTPUT_PATH = Path(META['output_path'])
EXPECT = META['verifier_expectations']


def sheet_row_highlighted(ws, row_idx, min_col, max_col):
    for col in range(min_col, max_col + 1):
        rgb = cell_fill_rgb(ws.cell(row_idx, col))
        if rgb and rgb not in ('FFFFFF', '000000', '00000000'):
            return True
    return False


def range_signatures(ws, refs):
    return [style_signature(ws[ref]) for ref in refs]


def heading_paragraphs(doc):
    out = []
    for p in doc.paragraphs:
        style = p.style.name if p.style else ''
        if p.text.strip() and style.startswith('Heading'):
            out.append((p.text.strip(), style))
    return out


def paragraphs_under_heading(doc, heading_text):
    paras = []
    capture = False
    for p in doc.paragraphs:
        text = p.text.strip()
        style = p.style.name if p.style else ''
        if text == heading_text:
            capture = True
            continue
        if capture and text and style.startswith('Heading'):
            break
        if capture and text:
            paras.append(text)
    return paras


def compact_text(text):
    text = normalize_text(str(text)).replace('\xa0', ' ')
    return re.sub(r'[^a-z0-9]+', ' ', text).strip()


def assert_contains_all(text, phrases, label):
    norm = compact_text(text)
    missing = [phrase for phrase in phrases if compact_text(phrase) not in norm]
    assert not missing, f'{label}: missing required anchors: {missing}'


def assert_contains_any_group(text, groups, label):
    norm = compact_text(text)
    for group in groups:
        if any(compact_text(phrase) in norm for phrase in group):
            return
    raise AssertionError(f'{label}: none of the accepted anchors were found: {groups}')


def has_page_break_before(doc, heading_text):
    prev = None
    for p in doc.paragraphs:
        if p.text.strip() == heading_text:
            if prev is None:
                return False
            return docx_has_page_break(prev) or docx_has_page_break(p)
        prev = p
    raise AssertionError(f'Heading not found for page-break check: {heading_text}')


def slide_titles(prs):
    titles = []
    for slide in prs.slides:
        title = None
        for shape in slide.shapes:
            if getattr(shape, 'has_text_frame', False):
                txt = ' '.join(shape.text.split())
                if txt:
                    title = txt.split('\n')[0]
                    break
        titles.append(title or '')
    return titles


def slide_by_title(prs, title):
    norm = normalize_text(title)
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, 'has_text_frame', False):
                txt = ' '.join(shape.text.split())
                if txt and normalize_text(txt.split('\n')[0]) == norm:
                    return slide
    raise AssertionError(f'Slide not found: {title}')


def slide_by_title_contains_any(prs, needles):
    for slide in prs.slides:
        title = ''
        for shape in slide.shapes:
            if getattr(shape, 'has_text_frame', False):
                txt = ' '.join(shape.text.split())
                if txt:
                    title = txt.split('\n')[0]
                    break
        if title and any(compact_text(needle) in compact_text(title) for needle in needles):
            return slide
    raise AssertionError(f'Slide not found with title containing any of: {needles}')


def slide_text_lines(slide, skip_title=True):
    lines = []
    title_text = None
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            txt = ' '.join(shape.text.split())
            if txt:
                title_text = txt.split('\n')[0]
                break
    for shape in slide.shapes:
        if not getattr(shape, 'has_text_frame', False):
            continue
        for para in shape.text_frame.paragraphs:
            text = ' '.join(para.text.split())
            if not text:
                continue
            if skip_title and title_text and normalize_text(text) == normalize_text(title_text):
                continue
            lines.append(text)
    return lines


def slide_table(slide):
    for shape in slide.shapes:
        if getattr(shape, 'has_table', False):
            return shape.table
    raise AssertionError('No table found on slide.')


def slide_picture_count(slide):
    return sum(1 for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE)


def page_titles_from_texts(page_texts):
    titles = []
    for text in page_texts:
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        titles.append(lines[0] if lines else '')
    return titles


def page_text_by_title(page_texts, title):
    titles = page_titles_from_texts(page_texts)
    for idx, current in enumerate(titles):
        if normalize_text(current) == normalize_text(title):
            return page_texts[idx]
    raise AssertionError(f'PDF page title not found: {title}')


def verify_excel():
    wb = load_workbook(OUTPUT_PATH)
    in_wb = load_workbook(INPUT_PATH)

    if 'sheet_order' in EXPECT:
        assert wb.sheetnames == EXPECT['sheet_order'], f'Unexpected sheet order: {wb.sheetnames!r}'

    for ref, expected in EXPECT.get('exact_cells', {}).items():
        sheet_name, cell_ref = ref.split('!')
        actual = wb[sheet_name][cell_ref].value
        actual = '' if actual is None else str(actual)
        assert actual == expected, f'{ref}: expected {expected!r}, found {actual!r}'

    for ref, expected in EXPECT.get('formula_cells', {}).items():
        sheet_name, cell_ref = ref.split('!')
        actual = wb[sheet_name][cell_ref].value
        assert actual == expected, f'{ref}: expected formula {expected!r}, found {actual!r}'

    for spec in EXPECT.get('header_style_matches', []):
        target = range_signatures(wb[spec['target_sheet']], spec['target_cells'])
        reference = range_signatures(wb[spec['reference_sheet']], spec['reference_cells'])
        assert target == reference, f"Style mismatch for {spec['target_sheet']} vs {spec['reference_sheet']}"

    for spec in EXPECT.get('highlight_rows', []):
        ws = wb[spec['sheet']]
        assert sheet_row_highlighted(ws, spec['row'], spec['min_col'], spec['max_col']), (
            f"Expected highlighted row {spec['sheet']}!{spec['row']}"
        )

    for spec in EXPECT.get('nonhighlight_rows', []):
        ws = wb[spec['sheet']]
        assert not sheet_row_highlighted(ws, spec['row'], spec['min_col'], spec['max_col']), (
            f"Row should not be highlighted: {spec['sheet']}!{spec['row']}"
        )

    for sheet_name in EXPECT.get('unchanged_sheets', []):
        assert workbook_values_signature(wb)[sheet_name] == workbook_values_signature(in_wb)[sheet_name], (
            f'Sheet values changed unexpectedly: {sheet_name}'
        )


def verify_word():
    doc = Document(OUTPUT_PATH)

    headings = heading_paragraphs(doc)
    if 'heading_order' in EXPECT:
        actual = [text for text, _style in headings]
        assert actual == EXPECT['heading_order'], f'Unexpected heading order: {actual!r}'

    for spec in EXPECT.get('heading_style_prefixes', []):
        match = [style for text, style in headings if text == spec['text']]
        assert match, f"Heading not found: {spec['text']}"
        assert match[0].startswith(spec['style_prefix']), (
            f"Heading style mismatch for {spec['text']}: {match[0]!r}"
        )

    for text in EXPECT.get('required_paragraphs', []):
        assert text in docx_texts(doc), f'Required paragraph missing: {text}'

    for text in EXPECT.get('forbidden_paragraphs', []):
        assert text not in docx_texts(doc), f'Forbidden paragraph present: {text}'

    for spec in EXPECT.get('paragraphs_under_headings', []):
        actual = paragraphs_under_heading(doc, spec['heading'])
        assert actual == spec['exact_lines'], (
            f"Unexpected content under heading {spec['heading']}: {actual!r}"
        )

    for spec in EXPECT.get('paragraph_anchor_groups_under_headings', []):
        actual = ' '.join(paragraphs_under_heading(doc, spec['heading']))
        assert_contains_all(actual, spec.get('required_all', []), f"heading {spec['heading']}")
        for groups in spec.get('required_any_groups', []):
            assert_contains_any_group(actual, groups, f"heading {spec['heading']}")

    for table_spec in EXPECT.get('tables', []):
        table = doc.tables[table_spec['index']]
        for key, expected in table_spec['cells'].items():
            row_idx, col_idx = [int(x) for x in key.split(',')]
            actual = table.cell(row_idx, col_idx).text.strip()
            assert actual == expected, f"table[{table_spec['index']}][{key}] expected {expected!r}, found {actual!r}"

    for table_spec in EXPECT.get('table_cell_anchor_groups', []):
        table = doc.tables[table_spec['index']]
        for key, spec in table_spec['cells'].items():
            row_idx, col_idx = [int(x) for x in key.split(',')]
            actual = table.cell(row_idx, col_idx).text.strip()
            assert_contains_all(actual, spec.get('required_all', []), f"table[{table_spec['index']}][{key}]")
            accepted = spec.get('accepted_all_groups', [])
            if accepted:
                matched = any(all(compact_text(phrase) in compact_text(actual) for phrase in group) for group in accepted)
                assert matched, f"table[{table_spec['index']}][{key}] did not match accepted anchors: {actual!r}"

    for spec in EXPECT.get('style_matches', []):
        target = docx_para_by_text(doc, spec['target_text'])
        reference = docx_para_by_text(doc, spec['reference_text'])
        assert docx_run_signature(target) == docx_run_signature(reference), (
            f"Style mismatch between {spec['target_text']!r} and {spec['reference_text']!r}"
        )

    for heading in EXPECT.get('page_break_before', []):
        assert has_page_break_before(doc, heading), f'Missing page break before heading: {heading}'


def verify_ppt():
    prs = Presentation(OUTPUT_PATH)

    if 'slide_count' in EXPECT:
        assert len(prs.slides) == EXPECT['slide_count'], f'Unexpected slide count: {len(prs.slides)}'

    if 'last_slide_title_contains' in EXPECT:
        actual = slide_titles(prs)
        assert actual, 'No slides found.'
        assert compact_text(EXPECT['last_slide_title_contains']) in compact_text(actual[-1]), (
            f"Last slide title {actual[-1]!r} does not contain {EXPECT['last_slide_title_contains']!r}"
        )

    if 'titles_order' in EXPECT:
        actual = slide_titles(prs)
        assert actual == EXPECT['titles_order'], f'Unexpected slide order: {actual!r}'

    for title in EXPECT.get('absent_slide_titles', []):
        assert title not in slide_titles(prs), f'Slide should be absent: {title}'

    for spec in EXPECT.get('slide_text', []):
        slide = slide_by_title(prs, spec['slide_title'])
        actual_lines = slide_text_lines(slide)
        if 'exact_lines' in spec:
            assert actual_lines == spec['exact_lines'], (
                f"Unexpected slide text on {spec['slide_title']}: {actual_lines!r}"
            )
        for line in spec.get('required_lines', []):
            assert line in actual_lines, f"Missing line on {spec['slide_title']}: {line}"

    for spec in EXPECT.get('slide_text_anchors', []):
        slide = slide_by_title_contains_any(prs, spec['title_contains_any'])
        actual_text = ' '.join(slide_text_lines(slide, skip_title=False))
        assert_contains_all(actual_text, spec.get('required_all', []), f"slide {spec['title_contains_any']}")
        for groups in spec.get('required_any_groups', []):
            assert_contains_any_group(actual_text, groups, f"slide {spec['title_contains_any']}")

    presentation_text = ' '.join(ppt_texts(prs))
    for spec in EXPECT.get('presentation_text_anchors', []):
        assert_contains_all(presentation_text, spec.get('required_all', []), 'presentation text')
        for groups in spec.get('required_any_groups', []):
            assert_contains_any_group(presentation_text, groups, 'presentation text')

    for spec in EXPECT.get('slide_tables', []):
        slide = slide_by_title(prs, spec['slide_title'])
        table = slide_table(slide)
        for key, expected in spec['cells'].items():
            row_idx, col_idx = [int(x) for x in key.split(',')]
            actual = table.cell(row_idx, col_idx).text.strip()
            assert actual == expected, (
                f"Unexpected table cell on {spec['slide_title']}[{key}]: {actual!r}"
            )

    for spec in EXPECT.get('picture_slides', []):
        slide = slide_by_title(prs, spec['slide_title'])
        assert slide_picture_count(slide) >= spec['min_pictures'], (
            f"Expected at least {spec['min_pictures']} pictures on {spec['slide_title']}"
        )

    for spec in EXPECT.get('style_matches', []):
        slide = slide_by_title(prs, spec['slide_title'])
        ref = slide_by_title(prs, spec['reference_title'])
        slide_lines = slide_text_lines(slide, skip_title=False)
        ref_lines = slide_text_lines(ref, skip_title=False)
        assert slide_lines, f'No text found on slide {spec["slide_title"]}'
        assert ref_lines, f'No text found on slide {spec["reference_title"]}'


def verify_pdf():
    page_texts = pdf_page_texts(OUTPUT_PATH)
    titles = page_titles_from_texts(page_texts)

    if 'page_titles_order' in EXPECT:
        assert titles == EXPECT['page_titles_order'], f'Unexpected page order: {titles!r}'

    for title in EXPECT.get('absent_page_titles', []):
        assert title not in titles, f'Page title should be absent: {title}'

    cover_text = page_texts[0] if page_texts else ''
    for line in EXPECT.get('cover_lines', []):
        assert normalize_text(line) in normalize_text(cover_text), f'Missing cover line: {line}'

    for spec in EXPECT.get('required_page_lines', []):
        text = page_text_by_title(page_texts, spec['page_title'])
        norm = normalize_text(text)
        for line in spec['lines']:
            assert normalize_text(line) in norm, f"Missing page line on {spec['page_title']}: {line}"

    for spec in EXPECT.get('forbidden_page_lines', []):
        text = page_text_by_title(page_texts, spec['page_title'])
        norm = normalize_text(text)
        for line in spec['lines']:
            assert normalize_text(line) not in norm, f"Forbidden page line remained on {spec['page_title']}: {line}"

    if 'outline' in EXPECT:
        outline = extract_pdf_outline_titles(OUTPUT_PATH)
        assert outline == [tuple(item) for item in EXPECT['outline']], f'Unexpected outline: {outline!r}'


def test_semantic_verifier():
    run_preflight(META, INPUT_PATH, OUTPUT_PATH)
    if META['doc_type'] == 'excel':
        verify_excel()
    elif META['doc_type'] == 'word':
        verify_word()
    elif META['doc_type'] == 'ppt':
        verify_ppt()
    elif META['doc_type'] == 'pdf':
        verify_pdf()
    else:
        raise AssertionError(f"Unsupported doc type: {META['doc_type']}")
