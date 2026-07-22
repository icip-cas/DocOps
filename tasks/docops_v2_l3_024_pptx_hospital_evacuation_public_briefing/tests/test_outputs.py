import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.enum.dml import MSO_FILL

from verifier_utils import normalize_text, run_preflight

META_PATH = Path(os.environ.get("TASK_METADATA_PATH", "/tests/task_metadata.json"))
if not META_PATH.exists():
    META_PATH = Path(__file__).parent / "task_metadata.json"
META = json.loads(META_PATH.read_text(encoding="utf-8"))
EXPECT = META["verifier_expectations"]
INPUT_PATH = Path(os.environ.get("INPUT_PATH", META["input_path"]))
OUTPUT_PATH = Path(os.environ.get("OUTPUT_PATH", META["output_path"]))


def rgb_str(rgb):
    return None if rgb is None else str(rgb).upper()


def slide_text(slide, include_notes=True):
    parts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            txt = shape.text.strip()
            if txt:
                parts.append(txt)
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    txt = cell.text.strip()
                    if txt:
                        parts.append(txt)
    if include_notes and getattr(slide, "has_notes_slide", False):
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            parts.append(notes)
    return "\n".join(parts)


def loose_text(slide, include_notes=False):
    return normalize_text(slide_text(slide, include_notes=include_notes))


def assert_terms(text, terms, label):
    missing = [term for term in terms if normalize_text(term) not in text]
    assert not missing, f"{label}: missing {missing}"


def all_text(prs):
    return "\n".join(slide_text(slide, include_notes=True) for slide in prs.slides)


def slide_title(slide):
    candidates = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = shape.text.strip()
            if text:
                candidates.append((shape.top, shape.left, text.splitlines()[0]))
    assert candidates, "Slide has no title text"
    return sorted(candidates, key=lambda item: (item[0], item[1]))[0][2]


def title_shape(slide):
    title = slide_title(slide)
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = shape.text.strip()
            if text and text.splitlines()[0] == title:
                return shape
    raise AssertionError(f"Could not find title shape for {title!r}")


def table_rows(slide):
    tables = []
    for shape in slide.shapes:
        if getattr(shape, "has_table", False) and shape.has_table:
            rows = []
            for row in shape.table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            tables.append(rows)
    return tables


def background_rgb(slide):
    fill = slide.background.fill
    if fill.type == MSO_FILL.SOLID:
        return rgb_str(fill.fore_color.rgb)
    slide_w = slide.part.package.presentation_part.presentation.slide_width
    slide_h = slide.part.package.presentation_part.presentation.slide_height
    for shape in slide.shapes:
        shape_fill = getattr(shape, "fill", None)
        if (
            shape_fill is not None
            and shape_fill.type == MSO_FILL.SOLID
            and shape.left <= 0
            and shape.top <= 0
            and shape.width >= slide_w * 0.95
            and shape.height >= slide_h * 0.95
        ):
            return rgb_str(shape_fill.fore_color.rgb)
    return None


def has_accent_shape(slide):
    expected = EXPECT["style"]["accent_fill"]
    for shape in slide.shapes:
        fill = getattr(shape, "fill", None)
        if fill is not None and fill.type == MSO_FILL.SOLID and rgb_str(fill.fore_color.rgb) == expected:
            return True
    return False


def title_font_info(slide):
    shape = title_shape(slide)
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                return run.font.name, bool(run.font.bold), rgb_str(run.font.color.rgb)
    raise AssertionError("Title has no text run")


def speaker_notes(prs):
    return [
        slide.notes_slide.notes_text_frame.text.strip() if getattr(slide, "has_notes_slide", False) else ""
        for slide in prs.slides
    ]


def native_table_text(slide):
    tables = table_rows(slide)
    assert tables, f"{slide_title(slide)}: expected a native public table"
    return normalize_text("\n".join(" | ".join(row) for table in tables for row in table))


def assert_table_header_style(slide):
    style = EXPECT["style"]
    for shape in slide.shapes:
        if getattr(shape, "has_table", False) and shape.has_table:
            for cell in shape.table.rows[0].cells:
                assert rgb_str(cell.fill.fore_color.rgb) == style["table_header_fill"]
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            assert run.font.bold
                            assert rgb_str(run.font.color.rgb) == style["table_header_font"]
                            return


def test_output_exists_and_is_pptx():
    run_preflight(META, INPUT_PATH, OUTPUT_PATH)
    assert OUTPUT_PATH != INPUT_PATH


def test_slide_order_and_private_material_removed():
    prs = Presentation(OUTPUT_PATH)
    assert [slide_title(slide) for slide in prs.slides] == EXPECT["slide_titles"]
    text = normalize_text(all_text(prs))
    hits = [p for p in EXPECT["forbidden_phrases"] if normalize_text(p) in text]
    assert not hits, f"Forbidden phrases still present: {hits}"


def test_required_public_content_and_tables():
    prs = Presentation(OUTPUT_PATH)
    by_title = {slide_title(slide): slide for slide in prs.slides}
    title_text = loose_text(by_title["Harborview Hospital Evacuation Exercise: Public Briefing"])
    assert_terms(title_text, ["public", "exercise", "briefing"], "Title slide")
    assert "2026-09-12" in title_text or "september 12, 2026" in title_text
    assert_terms(title_text, ["Regional Health Preparedness Council"], "Title slide")

    scope = loose_text(by_title["Exercise Scope and Assumptions"])
    assert_terms(scope, ["West Tower 5C", "patient movement", "transport", "family reunification", "communications"], "Exercise Scope")
    assert "clinical" in scope or "aggregate" in scope or "operational level" in scope

    timeline = native_table_text(by_title["Patient Movement Timeline"])
    assert_terms(timeline, ["08:10", "08:25", "09:05", "10:15", "11:40", "West Tower 5C", "Pediatric Unit"], "Patient Movement Timeline")

    unit = native_table_text(by_title["Unit Impact Matrix"])
    assert_terms(unit, ["West Tower 5C", "Partial evacuation", "18", "Pediatric Unit", "12", "Imaging", "0", "Emergency Department", "6"], "Unit Impact Matrix")

    transport = native_table_text(by_title["Transport Resource Plan"])
    assert_terms(transport, ["Ambulances", "18", "South ambulance bay", "Wheelchair vans", "6", "West Tower lobby", "9", "Bus support", "2"], "Transport Resource Plan")

    family = native_table_text(by_title["Family Reunification Workflow"])
    assert_terms(family, ["10:15", "Verify visitor identity", "30 minutes", "social work"], "Family Reunification Workflow")

    comms = native_table_text(by_title["Communications Cadence"])
    assert_terms(comms, ["Families", "Every 30 minutes", "Regional partners", "Hourly", "Public", "At status change", "Staff", "Every 45 minutes"], "Communications Cadence")

    commitments = native_table_text(by_title["Corrective Action Commitments"])
    assert_terms(commitments, ["West Tower 5C", "Facilities", "2026-09-30", "18", "Logistics", "2026-10-04", "Patient Experience", "2026-10-07", "Pediatric Unit", "Dr. Helena Mora", "2026-10-15"], "Corrective Action Commitments")

    style = loose_text(by_title["Publication Style Guide"])
    assert_terms(style, ["172A3A", "F8F4EC", "E4572E", "0B6E4F", "Georgia", EXPECT["style"]["footer_text"]], "Publication Style Guide")

    appendix = loose_text(by_title["Appendix Slide Index"])
    assert_terms(appendix, ["Exercise Scope", "Patient Movement Timeline", "Transport Resource Plan", "Corrective Action Commitments"], "Appendix Slide Index")


def test_speaker_notes_cleaned_and_preserved():
    prs = Presentation(OUTPUT_PATH)
    notes = speaker_notes(prs)
    assert len(notes) == len(EXPECT["slide_titles"])
    for idx, note in enumerate(notes, start=1):
        normalized = normalize_text(note)
        assert normalized, f"Slide {idx}: missing public speaker notes"
        assert len(note) <= 400, f"Slide {idx}: speaker notes should be concise"
        hits = [p for p in EXPECT["forbidden_phrases"] if normalize_text(p) in normalized]
        assert not hits, f"Slide {idx}: forbidden note material still present: {hits}"


def test_style_migration_applied():
    prs = Presentation(OUTPUT_PATH)
    style = EXPECT["style"]
    for idx, slide in enumerate(prs.slides):
        expected_bg = style["title_background"] if idx == 0 else style["content_background"]
        expected_title_color = style["title_color_title_slide"] if idx == 0 else style["title_color_content"]
        assert background_rgb(slide) == expected_bg
        font_name, bold, color = title_font_info(slide)
        assert font_name == style["title_font"]
        assert bold
        assert color == expected_title_color
        assert has_accent_shape(slide)
        assert style["footer_text"] in slide_text(slide, include_notes=False)
        if slide_title(slide) in EXPECT["tables"]:
            assert_table_header_style(slide)


def test_source_artifact_was_not_modified():
    source = Presentation(INPUT_PATH)
    assert len(source.slides) > len(EXPECT["slide_titles"])
    text = normalize_text(all_text(source))
    assert "draft" in text
    assert "patient name" in text
    assert "personal phone" in text
