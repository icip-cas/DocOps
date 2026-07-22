import json, os, re
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from verifier_utils import normalize_text, run_preflight
META_PATH=Path(os.environ.get('TASK_METADATA_PATH','/tests/task_metadata.json'))
if not META_PATH.exists(): META_PATH=Path(__file__).parent/'task_metadata.json'
META=json.loads(META_PATH.read_text()); EXPECT=META['verifier_expectations']; INPUT_PATH=Path(os.environ.get('INPUT_PATH',META['input_path'])); OUTPUT_PATH=Path(os.environ.get('OUTPUT_PATH',META['output_path']))
def rgb_str(rgb): return None if rgb is None else str(rgb).upper()
def slide_text(slide, include_notes=True):
    parts=[]
    for shape in slide.shapes:
        if getattr(shape,'has_text_frame',False) and shape.has_text_frame and shape.text.strip(): parts.append(shape.text.strip())
        if getattr(shape,'has_table',False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip(): parts.append(cell.text.strip())
    if include_notes and getattr(slide,'has_notes_slide',False):
        n=slide.notes_slide.notes_text_frame.text.strip()
        if n: parts.append(n)
    return '\n'.join(parts)
def all_text(prs): return '\n'.join(slide_text(s,True) for s in prs.slides)
def require_all(text, phrases, label):
    nt=normalize_text(text); missing=[p for p in phrases if normalize_text(p) not in nt]; assert not missing, f'{label} missing {missing}'
def require_any(text, phrases, label):
    nt=normalize_text(text); assert any(normalize_text(p) in nt for p in phrases), f'{label} expected one of {phrases}'
def cleanup_context(text, phrase):
    pattern=re.compile(rf'\\b{re.escape(phrase)}\\b', re.I)
    for m in pattern.finditer(text):
        window=text[max(0,m.start()-120):m.end()+180]
        if not re.search(r'\\b(no|not|removed|cleaned|cleared|excluded|without|public|superseded|corrected)\\b', window, re.I):
            return False
    return True
def assert_no_forbidden(text):
    nt=normalize_text(text)
    hard=['do not publish','peanut','FR-2407-D','2026-08-12','162 stores','555-0100','legal review hold','supplier dispute']
    hits=[p for p in hard if normalize_text(p) in nt]; assert not hits, hits
    contextual=[p for p in ['DRAFT','PRIVATE','INTERNAL'] if re.search(rf'\\b{re.escape(p)}\\b', text, re.I) and not cleanup_context(text,p)]
    assert not contextual, contextual
def slide_title(slide):
    cand=[]
    for shape in slide.shapes:
        if getattr(shape,'has_text_frame',False) and shape.has_text_frame and shape.text.strip(): cand.append((shape.top,shape.left,shape.text.strip().splitlines()[0]))
    assert cand, 'Slide has no title'; return sorted(cand,key=lambda x:(x[0],x[1]))[0][2]
def title_shape(slide):
    t=slide_title(slide)
    for shape in slide.shapes:
        if getattr(shape,'has_text_frame',False) and shape.has_text_frame:
            txt=shape.text.strip()
            if txt and txt.splitlines()[0]==t: return shape
    raise AssertionError('missing title shape')
def table_rows(slide):
    out=[]
    for shape in slide.shapes:
        if getattr(shape,'has_table',False) and shape.has_table:
            out.append([[cell.text.strip() for cell in row.cells] for row in shape.table.rows])
    return out
def background_rgb(slide):
    f=slide.background.fill; return rgb_str(f.fore_color.rgb) if f.type==MSO_FILL.SOLID else None
def title_font_info(slide):
    sh=title_shape(slide)
    for p in sh.text_frame.paragraphs:
        for r in p.runs:
            if r.text.strip(): return r.font.name,bool(r.font.bold),rgb_str(r.font.color.rgb)
    raise AssertionError('empty title')
def has_accent(slide):
    for sh in slide.shapes:
        f=getattr(sh,'fill',None)
        if f is not None and f.type==MSO_FILL.SOLID and rgb_str(f.fore_color.rgb)==EXPECT['style']['accent_fill']: return True
    return False
def notes(prs): return [s.notes_slide.notes_text_frame.text.strip() if getattr(s,'has_notes_slide',False) else '' for s in prs.slides]
def test_output_exists_and_is_pptx(): run_preflight(META,INPUT_PATH,OUTPUT_PATH); assert OUTPUT_PATH!=INPUT_PATH
def test_content_tables_notes_and_privacy():
    prs=Presentation(OUTPUT_PATH); assert len(prs.slides)==len(EXPECT['slide_titles'])
    slides=list(prs.slides)
    for slide, expected_title in zip(slides, EXPECT['slide_titles']):
        assert normalize_text(expected_title) in normalize_text(slide_text(slide, False)), f'missing slide title {expected_title}'
    raw_text=all_text(prs); assert_no_forbidden(raw_text)
    require_all(raw_text, ['2026-08-22','FR-2407-B','126','undeclared sesame','1-888-555-0134','https://health.example/recalls/FR-2407-B'], 'corrected recall values')
    by={expected_title: slide for expected_title, slide in zip(EXPECT['slide_titles'], slides)}
    require_all(slide_text(by['Public Food Recall Briefing'],False), ['2026-08-22','FR-2407-B','undeclared sesame'], 'title slide')
    require_all(slide_text(by['Correction Register Applied'],False), ['FR-2407-B','126','1-888-555-0134'], 'correction register')
    require_all(slide_text(by['Affected Product Summary'],False), ['undeclared sesame','126'], 'affected product summary')
    assert slide_text(by['Consumer Action Workflow'],False).strip()
    assert slide_text(by['Retailer Removal Checklist'],False).strip()
    require_all(slide_text(by['Communication Timeline'],False), ['2026-08-22'], 'communication timeline')
    require_any(slide_text(by['Final Release Checklist'],False), ['Correction register','corrected','complete','ready'], 'final checklist')
    for title in ['Correction Register Applied','Affected Product Summary']:
        assert table_rows(by[title]), f'{title} missing native public table'
    ns=notes(prs); assert len(ns)==len(EXPECT['slide_titles']); assert all(n.strip() for n in ns); assert_no_forbidden('\\n'.join(ns))
def test_style():
    prs=Presentation(OUTPUT_PATH); style=EXPECT['style']
    for i,s in enumerate(prs.slides):
        assert slide_text(s,False).strip()
        assert 'FR-2407-B' in slide_text(s,False) or 'Food Recall' in slide_text(s,False) or 'Recall' in slide_text(s,False)
        assert any(getattr(sh,'fill',None) is not None and getattr(sh,'fill',None).type==MSO_FILL.SOLID for sh in s.shapes), 'slide should include styled filled shapes'
def test_source_artifact_was_not_modified():
    prs=Presentation(INPUT_PATH); assert len(prs.slides)>len(EXPECT['slide_titles']); text=normalize_text(all_text(prs)); missing=[p for p in EXPECT['source_must_contain'] if normalize_text(p) not in text]; assert not missing, missing
