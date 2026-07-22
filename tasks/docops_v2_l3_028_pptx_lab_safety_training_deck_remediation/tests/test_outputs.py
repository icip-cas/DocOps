import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.enum.dml import MSO_FILL

from verifier_utils import normalize_text, run_preflight


META_PATH = Path(os.environ.get("TASK_METADATA_PATH", "/tests/task_metadata.json"))
if not META_PATH.exists():
    META_PATH = Path(__file__).parent / "task_metadata.json"

META = json.loads(META_PATH.read_text())
EXPECT = META["verifier_expectations"]
INPUT_PATH = Path(os.environ.get("INPUT_PATH", META["input_path"]))
OUTPUT_PATH = Path(os.environ.get("OUTPUT_PATH", META["output_path"]))


def rgb_str(rgb):
    return None if rgb is None else str(rgb).upper()


def slide_text(slide, include_notes=True):
    parts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame and shape.text.strip():
            parts.append(shape.text.strip())
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        parts.append(cell.text.strip())
    if include_notes and getattr(slide, "has_notes_slide", False):
        note = slide.notes_slide.notes_text_frame.text.strip()
        if note:
            parts.append(note)
    return "\n".join(parts)


def all_text(prs):
    return "\n".join(slide_text(slide, True) for slide in prs.slides)


def slide_title(slide):
    candidates = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame and shape.text.strip():
            candidates.append((shape.top, shape.left, shape.text.strip().splitlines()[0]))
    assert candidates, "Slide has no title text"
    return sorted(candidates, key=lambda item: (item[0], item[1]))[0][2]


def title_shape(slide, expected_title=None):
    title = expected_title or slide_title(slide)
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = shape.text.strip()
            if text and (text.splitlines()[0] == title or title in text):
                return shape
    raise AssertionError(f"Missing title shape: {title}")


def background_rgb(slide):
    fill = slide.background.fill
    return rgb_str(fill.fore_color.rgb) if fill.type == MSO_FILL.SOLID else None


def first_text_run(shape):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                return run
    raise AssertionError("Title shape has no text run")


def title_font_info(slide, expected_title=None):
    run = first_text_run(title_shape(slide, expected_title))
    return run.font.name, bool(run.font.bold), rgb_str(run.font.color.rgb)


def has_accent(slide):
    for shape in slide.shapes:
        fill = getattr(shape, "fill", None)
        if fill is not None and fill.type == MSO_FILL.SOLID:
            if rgb_str(fill.fore_color.rgb) == EXPECT["style"]["accent_fill"]:
                return True
    return False


def table_rows(slide):
    tables = []
    for shape in slide.shapes:
        if getattr(shape, "has_table", False) and shape.has_table:
            rows = []
            for row in shape.table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            tables.append(rows)
    return tables


def table_text(slide):
    parts = []
    for rows in table_rows(slide):
        for row in rows:
            parts.extend(row)
    return normalize_text("\n".join(parts))


def expected_title_shape(slide, title):
    return title_shape(slide, title)


def has_table_with_terms(slide, terms):
    text = table_text(slide)
    return bool(text) and all(normalize_text(term) in text for term in terms)


def has_table_with_any_terms(slide, required_terms, any_terms):
    text = table_text(slide)
    return bool(text) and all(normalize_text(term) in text for term in required_terms) and any(
        normalize_text(term) in text for term in any_terms
    )


def has_any_text_terms(slide, term_groups):
    text = normalize_text(slide_text(slide, include_notes=False))
    for group in term_groups:
        if not any(normalize_text(term) in text for term in group):
            return False
    return True


def speaker_notes(prs):
    return [
        slide.notes_slide.notes_text_frame.text.strip()
        if getattr(slide, "has_notes_slide", False)
        else ""
        for slide in prs.slides
    ]


def test_output_exists_and_is_pptx():
    run_preflight(META, INPUT_PATH, OUTPUT_PATH)
    assert OUTPUT_PATH != INPUT_PATH


def test_slide_order_size_required_content_and_privacy_cleanup():
    prs = Presentation(OUTPUT_PATH)
    assert len(prs.slides) == len(EXPECT["slide_titles"])
    for slide, expected_title in zip(prs.slides, EXPECT["slide_titles"]):
        expected_title_shape(slide, expected_title)

    normalized = normalize_text(all_text(prs))
    hits = [phrase for phrase in EXPECT["forbidden_phrases"] if normalize_text(phrase) in normalized]
    assert not hits, f"Forbidden draft/private values remain: {hits}"

    slides_by_title = {slide_title(slide): slide for slide in prs.slides}
    deck_text = normalize_text(all_text(prs))
    for phrase in ["TLS-2026-09", "2026-09-01", "Teaching labs B214 and B216", "lab-safety@univ.example"]:
        assert normalize_text(phrase) in deck_text
    assert has_any_text_terms(
        slides_by_title["Corrected Scope and Roles"],
        [["scope"], ["role", "learner", "supervisor", "coordinator"], ["lab-safety@univ.example"]],
    )
    assert has_any_text_terms(
        slides_by_title["PPE Donning Workflow"],
        [["coat"], ["goggle", "eye protection"], ["glove"]],
    )
    assert has_any_text_terms(
        slides_by_title["Spill Response Decision Tree"],
        [["spill"], ["stop", "isolate", "evacuate"], ["notify", "contact", "supervisor"]],
    )
    assert has_any_text_terms(
        slides_by_title["Training Release Checklist"],
        [["TLS-2026-09"], ["2026-09-01"], ["lab-safety@univ.example", "contact"]],
    )


def test_operational_tables_and_speaker_notes():
    prs = Presentation(OUTPUT_PATH)
    slides_by_title = {slide_title(slide): slide for slide in prs.slides}
    assert has_table_with_any_terms(
        slides_by_title["Hazard Classification Matrix"],
        ["hazard"],
        ["control", "ppe", "response", "action", "escalate", "risk"],
    ), "Hazard Classification Matrix missing native hazard/control/escalation table"
    assert has_table_with_terms(
        slides_by_title["Waste Labeling Checklist"],
        ["label"],
    ) and has_any_text_terms(
        slides_by_title["Waste Labeling Checklist"],
        [["chemical", "contents", "name"], ["date", "start"], ["owner", "course", "room", "hazard"]],
    ), "Waste Labeling Checklist missing public-safe checklist table"
    assert has_table_with_terms(
        slides_by_title["Knowledge Check"],
        ["question", "correct"],
    ), "Knowledge Check missing native question/correct-answer table"
    notes = speaker_notes(prs)
    assert len(notes) == len(EXPECT["slide_titles"])
    for idx, note in enumerate(notes, 1):
        assert len(note.split()) >= 8, f"Slide {idx} missing complete speaker notes"
        note_norm = normalize_text(note)
        hits = [phrase for phrase in EXPECT["forbidden_phrases"] if normalize_text(phrase) in note_norm]
        assert not hits, f"Slide {idx} notes contain restricted source text: {hits}"


def test_theme_background_titles_accent_and_footer():
    prs = Presentation(OUTPUT_PATH)
    style = EXPECT["style"]
    for idx, slide in enumerate(prs.slides):
        expected_bg = style["title_background"] if idx == 0 else style["content_background"]
        assert background_rgb(slide) == expected_bg
        font_name, bold, _color = title_font_info(slide, EXPECT["slide_titles"][idx])
        assert font_name == style["title_font"]
        assert bold
        assert has_accent(slide), f"Slide {idx + 1} missing required accent fill"
        assert style["footer_text"] in slide_text(slide, include_notes=False)


def test_source_artifact_was_not_modified():
    prs = Presentation(INPUT_PATH)
    assert len(prs.slides) > len(EXPECT["slide_titles"])
    normalized = normalize_text(all_text(prs))
    missing = [
        phrase for phrase in EXPECT["source_must_contain"]
        if normalize_text(phrase) not in normalized
    ]
    assert not missing, f"Source deck no longer contains expected draft values: {missing}"
