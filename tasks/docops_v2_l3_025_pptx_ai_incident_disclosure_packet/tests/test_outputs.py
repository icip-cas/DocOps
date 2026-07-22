import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.enum.dml import MSO_FILL

from verifier_utils import normalize_text, run_preflight

META_PATH = Path(os.environ.get("TASK_METADATA_PATH", "/tests/task_metadata.json"))
if not META_PATH.exists():
    META_PATH = Path(__file__).parent / "task_metadata.json"
META = json.loads(META_PATH.read_text(encoding="utf-8"))
EXPECT = META["verifier_expectations"]
INPUT_PATH = Path(os.environ.get("INPUT_PATH", META["input_path"]))
OUTPUT_PATH = Path(os.environ.get("OUTPUT_PATH", META["output_path"]))


def rgb_str(rgb):
    return None if rgb is None else str(rgb).upper()


def slide_text(slide, include_notes=True):
    parts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            txt = shape.text.strip()
            if txt:
                parts.append(txt)
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    txt = cell.text.strip()
                    if txt:
                        parts.append(txt)
    if include_notes and getattr(slide, "has_notes_slide", False):
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            parts.append(notes)
    return "\n".join(parts)


def all_text(prs):
    return "\n".join(slide_text(slide, include_notes=True) for slide in prs.slides)


def slide_norm(slide, include_notes=False):
    return normalize_text(slide_text(slide, include_notes=include_notes))


def table_text(slide):
    return normalize_text("\n".join(" | ".join(row) for table in table_rows(slide) for row in table))


def assert_terms(text, terms, label):
    missing = [term for term in terms if normalize_text(term) not in text]
    assert not missing, f"{label}: missing semantic terms {missing}"


def assert_any(text, terms, label):
    assert any(normalize_text(term) in text for term in terms), f"{label}: missing one of {terms}"


def slide_by_title(prs):
    slides = {slide_title(slide): slide for slide in prs.slides}
    missing = [title for title in EXPECT["slide_titles"] if title not in slides]
    assert not missing, f"Missing required slides: {missing}"
    return slides


def slide_title(slide):
    candidates = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = shape.text.strip()
            if text:
                candidates.append((shape.top, shape.left, text.splitlines()[0]))
    assert candidates, "Slide has no title text"
    return sorted(candidates, key=lambda item: (item[0], item[1]))[0][2]


def title_shape(slide):
    title = slide_title(slide)
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = shape.text.strip()
            if text and text.splitlines()[0] == title:
                return shape
    raise AssertionError(f"Could not find title shape for {title!r}")


def table_rows(slide):
    tables = []
    for shape in slide.shapes:
        if getattr(shape, "has_table", False) and shape.has_table:
            rows = []
            for row in shape.table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            tables.append(rows)
    return tables


def background_rgb(slide):
    fill = slide.background.fill
    return rgb_str(fill.fore_color.rgb) if fill.type == MSO_FILL.SOLID else None


def has_accent_shape(slide):
    expected = EXPECT["style"]["accent_fill"]
    for shape in slide.shapes:
        fill = getattr(shape, "fill", None)
        if fill is not None and fill.type == MSO_FILL.SOLID and rgb_str(fill.fore_color.rgb) == expected:
            return True
    return False


def title_font_info(slide):
    shape = title_shape(slide)
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                return run.font.name, bool(run.font.bold), rgb_str(run.font.color.rgb)
    raise AssertionError("Title has no text run")


def speaker_notes(prs):
    return [
        slide.notes_slide.notes_text_frame.text.strip() if getattr(slide, "has_notes_slide", False) else ""
        for slide in prs.slides
    ]


def assert_table_header_style(slide):
    style = EXPECT["style"]
    for shape in slide.shapes:
        if getattr(shape, "has_table", False) and shape.has_table:
            for cell in shape.table.rows[0].cells:
                assert rgb_str(cell.fill.fore_color.rgb) == style["table_header_fill"]
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            assert run.font.bold
                            assert rgb_str(run.font.color.rgb) == style["table_header_font"]
                            return


def test_output_exists_and_is_pptx():
    run_preflight(META, INPUT_PATH, OUTPUT_PATH)
    assert OUTPUT_PATH != INPUT_PATH


def test_slide_order_and_private_material_removed():
    prs = Presentation(OUTPUT_PATH)
    assert [slide_title(slide) for slide in prs.slides] == EXPECT["slide_titles"]
    text = normalize_text(all_text(prs))
    hits = [p for p in EXPECT["forbidden_phrases"] if normalize_text(p) in text]
    assert not hits, f"Forbidden phrases still present: {hits}"


def test_required_public_content_and_tables():
    prs = Presentation(OUTPUT_PATH)
    by_title = slide_by_title(prs)
    title_text = slide_norm(by_title["CivicHelp AI Incident Disclosure Packet"])
    assert_terms(title_text, ["CivicHelp", "AI", "Incident", "Disclosure"], "Title slide")
    assert_any(title_text, ["public disclosure packet", "public disclosure", "public packet", "public briefing"], "Title slide public status")
    assert_terms(title_text, ["2026-04-01", "2026-04-02"], "Title slide incident window")

    scope_text = slide_norm(by_title["Disclosure Scope"])
    assert_terms(scope_text, ["CivicHelp Chatbot v2.2"], "Disclosure Scope")
    assert_any(scope_text, ["1,240", "affected-session", "affected sessions"], "Disclosure Scope affected-session reference")
    assert_any(scope_text, ["public controls", "remediation", "transparency", "public notice", "service impact", "oversight"], "Disclosure Scope public workflow")
    assert_any(scope_text, ["excludes", "removed", "no confidential", "public-only", "omits", "non-public", "sensitive"], "Disclosure Scope exclusions")

    timeline = slide_norm(by_title["Corrected Incident Timeline"])
    assert_terms(timeline, ["2026-04-01 09:20", "CivicHelp Chatbot v2.2", "Contained"], "Timeline")
    assert_terms(timeline, ["1,240", "2026-04-02", "First public notice"], "Timeline correction facts")
    assert_any(timeline, ["Complete", "Published", "Ready"], "Timeline status")

    services = slide_norm(by_title["Affected Services Matrix"])
    assert_terms(services, ["Service triage", "CivicHelp Chatbot v2.2", "1,240", "Contained"], "Affected Services")
    assert_any(services, ["Misrouted", "affected", "public impact"], "Affected Services impact")

    controls = slide_norm(by_title["Remediation Controls"])
    assert_terms(controls, ["CTL-001", "Human review", "Jules Renner", "2026-04-04", "Active"], "Remediation Controls")

    transparency = slide_norm(by_title["Transparency Register"])
    assert_terms(transparency, ["TR-001", "First public notice posted", "2026-04-02", "Published"], "Transparency Register")
    assert_any(transparency, ["1,240", "disclosure packet", "transparency"], "Transparency Register follow-up")

    checklist = slide_norm(by_title["Public Disclosure Checklist"])
    assert_any(checklist, ["Correction", "Corrected", "Correct"], "Public Disclosure Checklist correction status")
    assert_any(checklist, ["1,240", "2026-04-02", "CivicHelp Chatbot v2.2", "Jules Renner"], "Public Disclosure Checklist corrected facts")
    assert_any(checklist, ["Complete", "Yes"], "Public Disclosure Checklist status")

    actions = slide_norm(by_title["Board Action Items"])
    assert_any(actions, ["Approve", "board", "follow-through", "action"], "Board Action Items action framing")
    assert_any(actions, ["disclosure", "publication", "public notice", "packet"], "Board Action Items disclosure follow-up")
    assert_terms(actions, ["control"], "Board Action Items control follow-up")
    assert_any(actions, ["transparency", "register", "update"], "Board Action Items follow-up")

    style = slide_norm(by_title["Publication Style Guide"])
    assert_terms(style, ["111827", "F9FAFB", "2563EB", "374151", "Aptos Display"], "Publication Style Guide")
    assert_terms(style, ["CivicHelp Public AI Disclosure | April 2026"], "Publication Style Guide footer")

    for title in [
        "Corrected Incident Timeline",
        "Affected Services Matrix",
        "Remediation Controls",
        "Transparency Register",
        "Public Disclosure Checklist",
        "Publication Style Guide",
    ]:
        assert table_rows(by_title[title]), f"{title}: expected a public table"


def test_speaker_notes_cleaned_and_preserved():
    prs = Presentation(OUTPUT_PATH)
    notes = speaker_notes(prs)
    assert len(notes) == len(EXPECT["slide_titles"])
    for idx, note in enumerate(notes, start=1):
        normalized = normalize_text(note)
        assert normalized, f"Slide {idx}: missing public speaker notes"
        assert len(note) <= 250, f"Slide {idx}: speaker notes should remain concise"
        hits = [p for p in EXPECT["forbidden_phrases"] if normalize_text(p) in normalized]
        assert not hits, f"Slide {idx}: forbidden note material still present: {hits}"


def test_style_migration_applied():
    prs = Presentation(OUTPUT_PATH)
    style = EXPECT["style"]
    for idx, slide in enumerate(prs.slides):
        expected_bg = style["title_background"] if idx == 0 else style["content_background"]
        expected_title_color = style["title_color_title_slide"] if idx == 0 else style["title_color_content"]
        assert background_rgb(slide) == expected_bg
        font_name, bold, color = title_font_info(slide)
        assert font_name == style["title_font"]
        assert bold
        assert color == expected_title_color
        assert has_accent_shape(slide)
        assert style["footer_text"] in slide_text(slide, include_notes=False)
        if slide_title(slide) in EXPECT["tables"]:
            assert_table_header_style(slide)


def test_source_artifact_was_not_modified():
    source = Presentation(INPUT_PATH)
    assert len(source.slides) > len(EXPECT["slide_titles"])
    text = normalize_text(all_text(source))
    assert "draft" in text
    assert "prompt transcript" in text
    assert "personal phone" in text
