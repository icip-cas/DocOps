import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.enum.dml import MSO_FILL

from verifier_utils import normalize_text, run_preflight

META_PATH = Path(os.environ.get("TASK_METADATA_PATH", "/tests/task_metadata.json"))
if not META_PATH.exists():
    META_PATH = Path(__file__).parent / "task_metadata.json"
META = json.loads(META_PATH.read_text(encoding="utf-8"))
EXPECT = META["verifier_expectations"]
INPUT_PATH = Path(os.environ.get("INPUT_PATH", META["input_path"]))
OUTPUT_PATH = Path(os.environ.get("OUTPUT_PATH", META["output_path"]))


def rgb_str(rgb):
    if rgb is None:
        return None
    return str(rgb).upper()


def slide_text(slide, include_notes=True):
    parts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            txt = shape.text.strip()
            if txt:
                parts.append(txt)
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    txt = cell.text.strip()
                    if txt:
                        parts.append(txt)
    if include_notes and getattr(slide, "has_notes_slide", False):
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            parts.append(notes)
    return "\n".join(parts)


def loose_text(slide, include_notes=False):
    return normalize_text(slide_text(slide, include_notes=include_notes))


def assert_has_terms(text, terms, message):
    missing = [term for term in terms if normalize_text(term) not in text]
    assert not missing, f"{message}: missing semantic terms {missing}"


def all_text(prs):
    return "\n".join(slide_text(slide, include_notes=True) for slide in prs.slides)


def slide_by_title(prs):
    slides = {slide_title(slide): slide for slide in prs.slides}
    missing = [title for title in EXPECT["slide_titles"] if title not in slides]
    assert not missing, f"Missing required slides: {missing}"
    return slides


def slide_title(slide):
    candidates = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = shape.text.strip()
            if text:
                candidates.append((shape.top, shape.left, text.splitlines()[0]))
    assert candidates, "Slide has no title text"
    return sorted(candidates, key=lambda item: (item[0], item[1]))[0][2]


def title_shape(slide):
    title = slide_title(slide)
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = shape.text.strip()
            if text and text.splitlines()[0] == title:
                return shape
    raise AssertionError(f"Could not find title shape for {title!r}")


def table_rows(slide):
    tables = []
    for shape in slide.shapes:
        if getattr(shape, "has_table", False) and shape.has_table:
            rows = []
            for row in shape.table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            tables.append(rows)
    return tables


def background_rgb(slide):
    fill = slide.background.fill
    if fill.type == MSO_FILL.SOLID:
        return rgb_str(fill.fore_color.rgb)
    slide_w = slide.part.package.presentation_part.presentation.slide_width
    slide_h = slide.part.package.presentation_part.presentation.slide_height
    for shape in slide.shapes:
        fill = getattr(shape, "fill", None)
        if (
            fill is not None
            and fill.type == MSO_FILL.SOLID
            and shape.left <= 0
            and shape.top <= 0
            and shape.width >= slide_w * 0.95
            and shape.height >= slide_h * 0.95
        ):
            return rgb_str(fill.fore_color.rgb)
    return None


def has_accent_shape(slide):
    expected = EXPECT["style"]["accent_fill"]
    for shape in slide.shapes:
        fill = getattr(shape, "fill", None)
        if fill is None or fill.type != MSO_FILL.SOLID:
            continue
        if rgb_str(fill.fore_color.rgb) == expected:
            return True
    return False


def title_font_info(slide):
    shape = title_shape(slide)
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                return run.font.name, bool(run.font.bold), rgb_str(run.font.color.rgb)
    raise AssertionError(f"Title shape {shape.name!r} has no text run")


def speaker_notes(prs):
    notes = []
    for slide in prs.slides:
        if getattr(slide, "has_notes_slide", False):
            notes.append(slide.notes_slide.notes_text_frame.text.strip())
        else:
            notes.append("")
    return notes


def native_table_text(slide):
    tables = table_rows(slide)
    assert tables, f"{slide_title(slide)}: expected a native table"
    return normalize_text("\n".join(" | ".join(row) for table in tables for row in table))


def test_output_exists_and_is_pptx():
    run_preflight(META, INPUT_PATH, OUTPUT_PATH)
    assert OUTPUT_PATH != INPUT_PATH


def test_slide_order_and_private_material_removed():
    prs = Presentation(OUTPUT_PATH)
    assert [slide_title(slide) for slide in prs.slides] == EXPECT["slide_titles"]
    text = normalize_text(all_text(prs))
    hits = [p for p in EXPECT["forbidden_phrases"] if normalize_text(p) in text]
    assert not hits, f"Forbidden phrases still present: {hits}"


def test_required_public_content_and_corrections():
    prs = Presentation(OUTPUT_PATH)
    by_title = slide_by_title(prs)
    title_text = loose_text(by_title["Port Azure Cyber Tabletop Exercise: Public After-Action Briefing"])
    assert_has_terms(title_text, ["public", "after-action", "briefing"], "Title slide")
    assert (
        "2026-08-03" in title_text or "august 3, 2026" in title_text
    ), "Title slide: missing exercise date"
    assert_has_terms(title_text, ["Port Resilience Committee"], "Title slide")

    scope_text = loose_text(by_title["Exercise Scope"])
    assert_has_terms(
        scope_text,
        ["degraded", "terminal", "scheduling", "public communications"],
        "Exercise Scope scenario",
    )
    assert_has_terms(
        scope_text,
        ["Terminal 2", "gate scheduling", "appointment restoration", "external communications"],
        "Exercise Scope included scope",
    )
    assert (
        (
            "exclude" in scope_text
            or "outside" in scope_text
            or "not included" in scope_text
            or "out of scope" in scope_text
        )
        and ("nonpublic" in scope_text or "confidential" in scope_text)
        and (
            "threat attribution" in scope_text
            or "investigative" in scope_text
            or "law enforcement" in scope_text
            or "technical investigation" in scope_text
        )
    ), "Exercise Scope: missing public-disclosure boundary"

    comm_text = loose_text(by_title["Public Communications"])
    assert "11:05" in comm_text and (
        "public update" in comm_text or "verified update" in comm_text or "first update" in comm_text
    )
    assert_has_terms(comm_text, ["advisory banner", "appointment"], "Public Communications")
    assert "2026-08-10" in comm_text and "dashboard" in comm_text

    commitments_text = loose_text(by_title["After-Action Commitments"])
    assert_has_terms(commitments_text, ["identity reset", "2026-08-04"], "After-Action Commitments")
    assert "update approval" in commitments_text
    assert "20 minutes" in commitments_text or "20-minute" in commitments_text
    assert_has_terms(commitments_text, ["gate scheduling drill"], "After-Action Commitments")
    assert ("each quarter" in commitments_text or "quarterly" in commitments_text)
    assert_has_terms(commitments_text, ["dashboard", "2026-08-10"], "After-Action Commitments")


def test_public_tables_rebuilt_exactly():
    prs = Presentation(OUTPUT_PATH)
    by_title = slide_by_title(prs)
    timeline = native_table_text(by_title["Incident Timeline"])
    for term in ["08:40", "09:15", "10:05", "11:05", "14:30", "Terminal 2"]:
        assert term.lower() in timeline, f"Incident Timeline: missing {term}"
    assert_has_terms(timeline, ["phishing report", "gate scheduling", "manual confirmation"], "Incident Timeline")
    assert "47 minutes" in timeline or "47-minute" in timeline or "stabilized" in timeline

    decisions = native_table_text(by_title["Decision Log"])
    for term in ["DEC-01", "DEC-02", "DEC-03", "DEC-04"]:
        assert term.lower() in decisions, f"Decision Log: missing {term}"
    assert_has_terms(decisions, ["manual queue", "advisory banner", "appointment restoration"], "Decision Log")
    assert (
        "threat" in decisions
        or "investigative" in decisions
        or "investigation" in decisions
        or "nonpublic" in decisions
        or "verified impacts" in decisions
        or "public reporting" in decisions
    ), "Decision Log: missing public-disclosure decision"

    impact = native_table_text(by_title["Operational Impact Matrix"])
    assert_has_terms(impact, ["Gate scheduling", "Terminal 2", "47 minutes", "Operations"], "Impact matrix")
    assert_has_terms(impact, ["Container release", "Manual verification", "2 hours", "Port Police"], "Impact matrix")
    assert_has_terms(impact, ["Reefer monitoring", "No public service interruption", "0 minutes", "Facilities"], "Impact matrix")
    assert_has_terms(impact, ["Public website", "Advisory banner", "3 hours", "Communications"], "Impact matrix")

    priorities = native_table_text(by_title["Recovery Priorities"])
    for term in ["P1", "P2", "P3", "P4"]:
        assert term.lower() in priorities, f"Recovery Priorities: missing {term}"
    assert_has_terms(priorities, ["identity reset", "CISO Dana Reed", "2026-08-04"], "Recovery Priorities")
    assert_has_terms(priorities, ["public update approval", "20 minutes", "Communications", "2026-08-07"], "Recovery Priorities")
    assert_has_terms(priorities, ["gate scheduling drill", "Operations", "2026-08-14"], "Recovery Priorities")
    assert_has_terms(priorities, ["recovery progress dashboard", "Technology", "2026-08-10"], "Recovery Priorities")


def test_speaker_notes_cleaned_and_preserved():
    prs = Presentation(OUTPUT_PATH)
    notes = speaker_notes(prs)
    assert len(notes) == len(EXPECT["slide_titles"])
    for idx, note in enumerate(notes, start=1):
        normalized = normalize_text(note)
        assert normalized, f"Slide {idx}: missing public speaker notes"
        assert len(note) <= 350, f"Slide {idx}: speaker notes should remain concise"
        hits = [p for p in EXPECT["forbidden_phrases"] if normalize_text(p) in normalized]
        assert not hits, f"Slide {idx}: forbidden note material still present: {hits}"


def test_style_migration_applied():
    prs = Presentation(OUTPUT_PATH)
    style = EXPECT["style"]
    for idx, slide in enumerate(prs.slides):
        if idx == 0:
            assert background_rgb(slide) == style["title_background"]
        else:
            assert background_rgb(slide) == style["content_background"]
        font_name, bold, color = title_font_info(slide)
        assert font_name == style["title_font"], f"Slide {idx + 1}: wrong title font"
        assert bold, f"Slide {idx + 1}: title should be bold"
        assert has_accent_shape(slide), f"Slide {idx + 1}: missing accent shape"
        assert style["footer_text"] in slide_text(slide, include_notes=False)


def test_source_artifact_was_not_modified():
    source = Presentation(INPUT_PATH)
    assert len(source.slides) > len(EXPECT["slide_titles"])
    text = normalize_text(all_text(source))
    assert "draft" in text
    assert "ransomware attribution" in text
    assert "personal phone" in text
