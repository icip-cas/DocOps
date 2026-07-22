import json
import os
import re
import sys
import zipfile
from pathlib import Path
from datetime import date, datetime

from docx import Document
from docx.enum.section import WD_ORIENT
from pptx import Presentation

sys.path.insert(0, str(Path(__file__).parent))
from verifier_utils import *  # noqa: F401,F403

META_PATH = Path(os.environ.get('TASK_METADATA_PATH', '/tests/task_metadata.json'))
META = json.loads(META_PATH.read_text(encoding='utf-8'))
EXPECT = META['verifier_expectations']


def _vopt_norm_text(value):
    if value is None:
        return ''
    if isinstance(value, datetime):
        if value.hour or value.minute or value.second:
            return value.strftime('%Y-%m-%d %H:%M')
        return value.strftime('%Y-%m-%d')
    if isinstance(value, date):
        return value.strftime('%Y-%m-%d')
    if isinstance(value, float) and value.is_integer():
        value = int(value)
    text = str(value)
    text = text.replace('\u2013', '-').replace('\u2014', '-').replace('\u2212', '-')
    text = text.replace('\xa0', ' ')
    text = re.sub(r'\s+', ' ', text).strip()
    return text


def _vopt_norm_formula(value):
    return re.sub(r'\s+', '', _vopt_norm_text(value)).upper()


def _vopt_assert_formula(actual, expected, label='formula'):
    assert _vopt_norm_formula(actual) == _vopt_norm_formula(expected), f"{label}: expected {expected!r}, found {actual!r}"


def _vopt_norm_rows(rows):
    return [[_vopt_norm_text(cell) for cell in row] for row in rows]


def _vopt_rows_equal(actual, expected):
    return _vopt_norm_rows(actual) == _vopt_norm_rows(expected)


def _vopt_rows_text(rows):
    return '\n'.join('|'.join(_vopt_norm_text(cell) for cell in row) for row in rows)


def _vopt_norm_key(value):
    return re.sub(r'[^a-z0-9]+', '', _vopt_norm_text(value).lower())


def _vopt_header_has(header, required):
    actual = {_vopt_norm_text(cell).lower() for cell in header}
    missing = [cell for cell in required if _vopt_norm_text(cell).lower() not in actual]
    assert not missing, f"Missing expected header columns {missing!r}; actual={header!r}"


def _vopt_section_text(section_part):
    return '\n'.join(p.text for p in getattr(section_part, 'paragraphs', []) if p.text is not None)


def _vopt_header_text(doc):
    return '\n'.join(_vopt_section_text(section.header) for section in doc.sections)


def _vopt_footer_text(doc):
    return '\n'.join(_vopt_section_text(section.footer) for section in doc.sections)


def _vopt_table_ref(ws, expected_name=None):
    tables = ws.tables
    if expected_name:
        for name in tables.keys():
            if str(name).lower() == str(expected_name).lower():
                return tables[name].ref
    vals = list(tables.values())
    assert vals, f"No Excel native table found on sheet {ws.title!r}"
    return vals[0].ref


def _vopt_clean_area(value):
    text = _vopt_norm_text(value).replace("'", '').replace('$', '').replace(' ', '')
    if '!' in text:
        text = text.split('!', 1)[1]
    return text.upper()


def _vopt_assert_print_area(ws, expected):
    actual = _vopt_clean_area(ws.print_area)
    target = _vopt_clean_area(expected)
    assert target in actual or actual in target, f"{ws.title}: expected print area {expected!r}, found {ws.print_area!r}"


def _vopt_freeze_pane(value):
    return getattr(value, 'coordinate', value)


def _vopt_validated_cells(ws):
    cells = set()
    for dv in ws.data_validations.dataValidation:
        for rng in dv.cells.ranges:
            text = str(rng).replace('$', '')
            try:
                min_col, min_row, max_col, max_row = range_boundaries(text)
            except ValueError:
                cells.add(text)
                continue
            if max_row > 10000:
                max_row = max(min_row, 200)
            for row in range(min_row, max_row + 1):
                for col in range(min_col, max_col + 1):
                    cells.add(f"{get_column_letter(col)}{row}")
    return cells


class _VoptValidationRanges(list):
    def __init__(self, ws, ranges):
        super().__init__(ranges)
        self.ws = ws

    def __contains__(self, expected):
        expected = str(expected).replace('$', '')
        if super().__contains__(expected):
            return True
        try:
            min_col, min_row, max_col, max_row = range_boundaries(expected)
        except ValueError:
            return super().__contains__(expected)
        target = {f"{get_column_letter(col)}{row}" for row in range(min_row, max_row + 1) for col in range(min_col, max_col + 1)}
        return target.issubset(_vopt_validated_cells(self.ws))


def _vopt_dv_ranges(ws):
    ranges = []
    for dv in ws.data_validations.dataValidation:
        ranges.extend(str(rng).replace('$', '') for rng in dv.cells.ranges)
    return _VoptValidationRanges(ws, ranges)


def _vopt_ordered_subset(expected, actual):
    actual_iter = iter([_vopt_norm_text(item) for item in actual])
    for item in [_vopt_norm_text(value) for value in expected]:
        for candidate in actual_iter:
            if item == candidate:
                break
        else:
            return False
    return True



def _doc_path():
    return Path(os.environ.get('DOCX_OUTPUT_PATH', EXPECT['document_output']))


def _ppt_path():
    return Path(os.environ.get('PPT_OUTPUT_PATH', EXPECT['deck_output']))


def _has_toc(path):
    with zipfile.ZipFile(path) as zf:
        xml = zf.read('word/document.xml').decode('utf-8', errors='ignore')
    return bool(re.search(r'TOC\s+(?:\\)?o|TOC\s*(?:&quot;|")', xml))


def _headings(doc):
    out = []
    for p in doc.paragraphs:
        style = p.style.name if p.style else ''
        if p.text.strip() and (style.startswith('Heading') or style == 'Title'):
            out.append(p.text.strip())
    return out


def _doc_text(doc):
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append('|'.join(cell.text.strip() for cell in row.cells))
    return '\n'.join(parts)


def _table_rows(table):
    return [[cell.text.strip() for cell in row.cells] for row in table.rows]


def _highlighted(doc, text):
    for p in doc.paragraphs:
        if _vopt_norm_text(text) in _vopt_norm_text(p.text):
            return any(run.font.highlight_color is not None for run in p.runs) or '<w:highlight' in p._p.xml
    return False


def _has_highlighted_terms(doc, terms):
    for p in doc.paragraphs:
        text = _vopt_norm_text(p.text).lower()
        if not text:
            continue
        has_highlight = any(run.font.highlight_color is not None for run in p.runs) or '<w:highlight' in p._p.xml
        if has_highlight and all(term.lower() in text for term in terms):
            return True
    return False


def _assert_heading_sequence(doc):
    actual = [_vopt_norm_key(item) for item in _headings(doc)]
    expected = [_vopt_norm_key(item) for item in EXPECT['heading_order']]
    doc_text = _vopt_norm_key(_doc_text(doc))
    pos = 1 if expected and expected[0] in doc_text else 0
    for item in actual:
        if pos < len(expected) and (item == expected[pos] or item.endswith(expected[pos])):
            pos += 1
    assert pos == len(expected), f"Missing or out-of-order required headings: {_headings(doc)!r}"


def _assert_word_semantics(text):
    for rid in EXPECT['public_roles']:
        assert rid in text, f'docx: missing public role {rid}'
    for pid in EXPECT['public_projects']:
        assert pid in text, f'docx: missing public project {pid}'


def _assert_ppt_semantics(text):
    norm = _vopt_norm_text(text).lower()
    assert 'graphite' in norm and 'copper' in norm, 'pptx: missing graphite/copper system reference'
    assert 'locked reference' in norm or ('locked' in norm and 'reference' in norm), 'pptx: missing locked reference marker'


def _row_texts_from_table(table):
    return [' | '.join(_vopt_norm_text(cell) for cell in row) for row in _table_rows(table)]


def _assert_rows_present_in_any_table(tables, expected_rows, label):
    remaining = {row[0]: row for row in expected_rows}
    for table in tables:
        text_rows = _row_texts_from_table(table)
        for row_id, expected in list(remaining.items()):
            for actual in text_rows:
                if row_id in actual and all(_vopt_norm_text(token) in actual for token in expected):
                    remaining.pop(row_id)
                    break
    assert not remaining, f"Missing {label} rows: {list(remaining.values())!r}"


def _norm_lower(value):
    return _vopt_norm_text(value).lower()


def _has_terms(text, terms):
    lowered = _norm_lower(text)
    return all(term.lower() in lowered for term in terms)


def _slide_title(slide):
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            text = ' '.join(shape.text.split())
            if text:
                return text
    return ''


def _slide_text(slide):
    parts = []
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            text = ' '.join(shape.text.split())
            if text:
                parts.append(text)
        if getattr(shape, 'has_table', False):
            for row in shape.table.rows:
                parts.append('|'.join(cell.text.strip() for cell in row.cells))
    return '\n'.join(parts)


def _first_table_by_title(prs, title):
    for slide in prs.slides:
        if _vopt_norm_text(title) in _vopt_norm_text(_slide_text(slide)):
            for shape in slide.shapes:
                if getattr(shape, 'has_table', False):
                    return shape.table
    raise AssertionError(f'Missing table for slide {title}')


def _assert_slide_sequence(prs):
    slide_texts = [_slide_text(slide) for slide in prs.slides]
    assert len(slide_texts) >= len(EXPECT['ppt_titles']), f'Expected at least {len(EXPECT["ppt_titles"])} slides, found {len(slide_texts)}'
    expected_terms = [
        ['shortlist', 'interview', 'deck'],
        ['storyline', 'team'],
        ['project', 'proof'],
        ['risk', 'delivery'],
        ['locked', 'interview', 'grid'],
    ]
    for terms, text in zip(expected_terms, slide_texts):
        assert _has_terms(text, terms), f'Missing PPT slide terms {terms!r}'


def _assert_rows_present_in_ppt_tables(prs, expected_rows, label):
    found_texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, 'has_table', False):
                for row in shape.table.rows:
                    found_texts.append(' | '.join(_vopt_norm_text(cell.text) for cell in row.cells))
    full_text = '\n'.join([_slide_text(slide) for slide in prs.slides] + found_texts)
    missing = []
    for row in expected_rows:
        if not all(_vopt_norm_text(token) in full_text for token in row):
            missing.append(row)
    assert not missing, f'Missing {label} rows in PPTX: {missing!r}'


def test_outputs_exist():
    assert _doc_path().exists()
    assert _ppt_path().exists()
    assert _doc_path().suffix.lower() == '.docx'
    assert _ppt_path().suffix.lower() == '.pptx'


def test_word_package_format_and_boundary():
    doc = Document(_doc_path())
    _assert_heading_sequence(doc)
    assert _has_toc(_doc_path())
    assert len(doc.sections) >= EXPECT['section_count_min']
    assert doc.sections[EXPECT['landscape_section_index']].orientation == WD_ORIENT.LANDSCAPE
    text = _doc_text(doc)
    _assert_word_semantics(text)
    forbid_any(text, EXPECT['forbidden_text'], 'docx')
    for pid in EXPECT['confidential_projects']:
        assert pid not in text, f'{pid} leaked to Word output'
    assert EXPECT['header_contains'] in _vopt_header_text(doc)
    assert EXPECT['footer_contains'] in _vopt_footer_text(doc)
    _assert_rows_present_in_any_table(doc.tables, EXPECT['team_rows'], 'team')
    _assert_rows_present_in_any_table(doc.tables, EXPECT['project_rows'], 'project')
    assert _has_highlighted_terms(doc, ['format exception', 'confidential', 'omitted'])


def test_ppt_interview_deck_format_and_tables():
    prs = Presentation(_ppt_path())
    _assert_slide_sequence(prs)
    full = '\n'.join(_slide_text(slide) for slide in prs.slides)
    _assert_ppt_semantics(full)
    forbid_any(full, EXPECT['ppt_forbidden'], 'pptx')
    for pid in EXPECT['confidential_projects']:
        assert pid not in full, f'{pid} leaked to PPTX output'
    _assert_rows_present_in_ppt_tables(prs, EXPECT['team_rows'], 'team')
    _assert_rows_present_in_ppt_tables(prs, EXPECT['project_rows'], 'project')


def test_cross_output_role_project_consistency():
    doc_text = _doc_text(Document(_doc_path()))
    ppt_text = '\n'.join(_slide_text(slide) for slide in Presentation(_ppt_path()).slides)
    for rid in EXPECT['public_roles']:
        assert rid in doc_text and rid in ppt_text
    for pid in EXPECT['public_projects']:
        assert pid in doc_text and pid in ppt_text
    for pid in EXPECT['confidential_projects']:
        assert pid not in ppt_text
        assert pid not in doc_text
