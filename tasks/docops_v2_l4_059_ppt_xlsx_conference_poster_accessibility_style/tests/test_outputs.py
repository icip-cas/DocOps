import json
import os
import re
import sys
from pathlib import Path
from datetime import date, datetime

from openpyxl import load_workbook
from openpyxl.utils.cell import get_column_letter, range_boundaries
from pptx import Presentation

sys.path.insert(0, str(Path(__file__).parent))
from verifier_utils import *  # noqa: F401,F403

META_PATH = Path(os.environ.get('TASK_METADATA_PATH', '/tests/task_metadata.json'))
META = json.loads(META_PATH.read_text(encoding='utf-8'))
EXPECT = META['verifier_expectations']


def _vopt_norm_text(value):
    if value is None:
        return ''
    if isinstance(value, datetime):
        if value.hour or value.minute or value.second:
            return value.strftime('%Y-%m-%d %H:%M')
        return value.strftime('%Y-%m-%d')
    if isinstance(value, date):
        return value.strftime('%Y-%m-%d')
    if isinstance(value, float) and value.is_integer():
        value = int(value)
    text = str(value)
    text = text.replace('\u2013', '-').replace('\u2014', '-').replace('\u2212', '-')
    text = text.replace('\xa0', ' ')
    text = re.sub(r'\s+', ' ', text).strip()
    return text


def _vopt_norm_formula(value):
    return re.sub(r'\s+', '', _vopt_norm_text(value)).upper()


def _vopt_assert_formula(actual, expected, label='formula'):
    assert _vopt_norm_formula(actual) == _vopt_norm_formula(expected), f"{label}: expected {expected!r}, found {actual!r}"


def _vopt_norm_rows(rows):
    return [[_vopt_norm_text(cell) for cell in row] for row in rows]


def _vopt_rows_equal(actual, expected):
    return _vopt_norm_rows(actual) == _vopt_norm_rows(expected)


def _vopt_norm_key(value):
    return re.sub(r'[^a-z0-9]+', '', _vopt_norm_text(value).lower())


def _vopt_rows_text(rows):
    return '\n'.join('|'.join(_vopt_norm_text(cell) for cell in row) for row in rows)


def _vopt_header_has(header, required):
    actual = {_vopt_norm_text(cell).lower() for cell in header}
    missing = [cell for cell in required if _vopt_norm_text(cell).lower() not in actual]
    assert not missing, f"Missing expected header columns {missing!r}; actual={header!r}"


def _vopt_section_text(section_part):
    return '\n'.join(p.text for p in getattr(section_part, 'paragraphs', []) if p.text is not None)


def _vopt_header_text(doc):
    return '\n'.join(_vopt_section_text(section.header) for section in doc.sections)


def _vopt_footer_text(doc):
    return '\n'.join(_vopt_section_text(section.footer) for section in doc.sections)


def _vopt_table_ref(ws, expected_name=None):
    tables = ws.tables
    if expected_name:
        for name in tables.keys():
            if str(name).lower() == str(expected_name).lower():
                return tables[name].ref
    vals = list(tables.values())
    assert vals, f"No Excel native table found on sheet {ws.title!r}"
    return vals[0].ref


def _vopt_clean_area(value):
    text = _vopt_norm_text(value).replace("'", '').replace('$', '').replace(' ', '')
    if '!' in text:
        text = text.split('!', 1)[1]
    return text.upper()


def _vopt_assert_print_area(ws, expected):
    actual = _vopt_clean_area(ws.print_area)
    target = _vopt_clean_area(expected)
    assert target in actual or actual in target, f"{ws.title}: expected print area {expected!r}, found {ws.print_area!r}"


def _vopt_freeze_pane(value):
    return getattr(value, 'coordinate', value)


def _vopt_validated_cells(ws):
    cells = set()
    for dv in ws.data_validations.dataValidation:
        for rng in dv.cells.ranges:
            text = str(rng).replace('$', '')
            try:
                min_col, min_row, max_col, max_row = range_boundaries(text)
            except ValueError:
                cells.add(text)
                continue
            if max_row > 10000:
                max_row = max(min_row, 200)
            for row in range(min_row, max_row + 1):
                for col in range(min_col, max_col + 1):
                    cells.add(f"{get_column_letter(col)}{row}")
    return cells


class _VoptValidationRanges(list):
    def __init__(self, ws, ranges):
        super().__init__(ranges)
        self.ws = ws

    def __contains__(self, expected):
        expected = str(expected).replace('$', '')
        if super().__contains__(expected):
            return True
        try:
            min_col, min_row, max_col, max_row = range_boundaries(expected)
        except ValueError:
            return super().__contains__(expected)
        target = {f"{get_column_letter(col)}{row}" for row in range(min_row, max_row + 1) for col in range(min_col, max_col + 1)}
        return target.issubset(_vopt_validated_cells(self.ws))


def _vopt_dv_ranges(ws):
    ranges = []
    for dv in ws.data_validations.dataValidation:
        ranges.extend(str(rng).replace('$', '') for rng in dv.cells.ranges)
    return _VoptValidationRanges(ws, ranges)


def _vopt_ordered_subset(expected, actual):
    actual_iter = iter([_vopt_norm_text(item) for item in actual])
    for item in [_vopt_norm_text(value) for value in expected]:
        for candidate in actual_iter:
            if item == candidate:
                break
        else:
            return False
    return True



def _deck_path():
    return Path(os.environ.get('PPT_OUTPUT_PATH', EXPECT['deck_output']))


def _workbook_path():
    return Path(os.environ.get('XLSX_OUTPUT_PATH', EXPECT['workbook_output']))


def _slide_title(slide):
    title_shape = getattr(slide.shapes, 'title', None)
    if title_shape is not None and getattr(title_shape, 'has_text_frame', False):
        text = ' '.join(title_shape.text.split())
        if text:
            return text
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            text = ' '.join(shape.text.split())
            if text:
                return text
    return ''


def _slide_text(slide):
    parts = []
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False):
            text = ' '.join(shape.text.split())
            if text:
                parts.append(text)
        if getattr(shape, 'has_table', False):
            for row in shape.table.rows:
                parts.append('|'.join(cell.text.strip() for cell in row.cells))
    return '\n'.join(parts)


def _first_table(slide):
    for shape in slide.shapes:
        if getattr(shape, 'has_table', False):
            return shape.table
    raise AssertionError('Missing native PPT table')


def _table_rows(table):
    return [[cell.text.strip() for cell in row.cells] for row in table.rows]


def _status_equivalent(actual, expected):
    actual_key = _vopt_norm_key(actual)
    expected_key = _vopt_norm_key(expected)
    if actual_key == expected_key:
        return True
    if expected_key == 'needscaption':
        return 'caption' in actual_key and ('need' in actual_key or 'hold' in actual_key or 'required' in actual_key)
    if expected_key == 'ready':
        return actual_key in {'ready', 'captioned', 'ok'} or 'ready' in actual_key
    return False


def _assert_session_rows(rows):
    actual_by_id = {_vopt_norm_text(row[0]): row for row in rows if row}
    for expected in EXPECT['session_rows']:
        pid = expected[0]
        assert pid in actual_by_id, f'Missing poster row {pid}'
        row = actual_by_id[pid]
        for idx in range(1, 4):
            assert _vopt_norm_text(row[idx]) == _vopt_norm_text(expected[idx]), (
                f'{pid}: expected column {idx + 1}={expected[idx]!r}, found {row[idx]!r}'
            )
        assert _status_equivalent(row[4], expected[4]), f'{pid}: expected status {expected[4]!r}, found {row[4]!r}'


def _row_record_from_headers(ws, row):
    header = [_vopt_norm_text(ws.cell(2, col).value) for col in range(1, ws.max_column + 1)]
    return {_vopt_norm_key(header[col - 1]): ws.cell(row, col).value for col in range(1, ws.max_column + 1)}


def _assert_formula_contains(cell, tokens, context):
    formula = _vopt_norm_formula(cell.value)
    assert formula.startswith('='), f'{context}: expected live formula, found {cell.value!r}'
    missing = [token for token in tokens if token.upper() not in formula]
    assert not missing, f'{context}: formula missing {missing!r}; found {cell.value!r}'


def _all_table_rows(prs):
    rows = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, 'has_table', False):
                rows.extend(_table_rows(shape.table))
    return rows


def _assert_reference_grid(slide_texts):
    matches = [text for text in slide_texts if 'reference' in _vopt_norm_text(text).lower() and 'grid' in _vopt_norm_text(text).lower()]
    assert matches, 'Missing locked/reference poster grid slide'
    assert any('not export' in _vopt_norm_text(text).lower() or 'do not export' in _vopt_norm_text(text).lower() for text in matches), (
        'Reference grid slide must be marked not for export'
    )


def _spec_sheet(wb):
    candidates = []
    for ws in wb.worksheets:
        if ws.sheet_state != 'visible':
            continue
        text = _worksheet_text(ws)
        if ('caption' in text or 'accessibility' in text or 'print' in text) and ('export' in text or 'hold' in text):
            candidates.append(ws)
    assert candidates, f'No visible print/accessibility spec sheet found; sheets={wb.sheetnames!r}'
    return candidates[0]


def _workbook_text(wb):
    return '\n'.join(_worksheet_text(ws) for ws in wb.worksheets)


def _worksheet_text(ws):
    values = []
    for row in ws.iter_rows(values_only=True):
        for value in row:
            if value is not None:
                values.append(_vopt_norm_text(value))
    return '\n'.join(values).lower()


def _assert_workbook_semantics(ws, workbook_text):
    visible_text = _worksheet_text(ws)
    if all(pid.lower() in visible_text for pid in ['P-104', 'P-118', 'P-126']):
        for expected in EXPECT['session_rows']:
            pid, study, presenter, track, status = expected
            for term in [pid, study, presenter, track]:
                assert term.lower() in visible_text, f'{term} missing from workbook spec'
    else:
        assert 'slide 3' in visible_text and 'caption' in visible_text, 'Slide-level print spec must include caption repair row'
        assert 'slide 4' in visible_text and ('reference' in visible_text or 'do not export' in visible_text), 'Slide-level print spec must include reference layout row'
    assert 'p-118' in workbook_text and 'caption' in workbook_text and ('hold' in workbook_text or 'needs caption' in workbook_text), (
        'P-118 caption hold missing from workbook'
    )
    assert 'do not export' in workbook_text or 'not export' in workbook_text, 'Reference grid do-not-export rule missing from workbook'


def _assert_any_formulas(ws):
    formulas = [cell.value for row in ws.iter_rows() for cell in row if isinstance(cell.value, str) and cell.value.startswith('=')]
    assert formulas, 'Workbook must include live formulas'
    joined = ' '.join(_vopt_norm_formula(formula) for formula in formulas)
    assert 'HOLD' in joined and ('EXPORT' in joined or 'POSTER' in joined), 'Workbook formulas must compute export/hold status'


def test_outputs_exist():
    assert _deck_path().exists()
    assert _workbook_path().exists()
    assert _deck_path().suffix.lower() == '.pptx'
    assert _workbook_path().suffix.lower() == '.xlsx'


def test_ppt_slide_size_order_content_and_tables():
    prs = Presentation(_deck_path())
    assert abs(prs.slide_width - EXPECT['slide_size'][0]) < 1000
    assert abs(prs.slide_height - EXPECT['slide_size'][1]) < 1000
    assert len(prs.slides) >= 2, 'Deck should contain poster content and retained reference grid'
    slide_texts = [_slide_text(slide) for slide in prs.slides]
    full = '\n'.join(slide_texts)
    forbid_any(full, EXPECT['deck_forbidden'], 'deck')
    assert 'poster' in full.lower() and ('accessibility' in full.lower() or 'caption' in full.lower()), 'Missing poster accessibility content'
    assert 'P-118' in full and 'caption' in full.lower() and ('hold' in full.lower() or 'needs caption' in full.lower()), 'Missing P-118 caption hold in deck'
    _assert_reference_grid(slide_texts)
    rows = _all_table_rows(prs)
    _assert_session_rows(rows)


def test_workbook_print_spec_controls():
    wb = load_workbook(_workbook_path(), data_only=False)
    assert 'Rules' in wb.sheetnames, f'Missing Rules sheet; sheets={wb.sheetnames!r}'
    ws = _spec_sheet(wb)
    assert ws.merged_cells.ranges, 'Workbook must include a merged title area'
    workbook_text = _workbook_text(wb)
    _assert_workbook_semantics(ws, workbook_text)
    assert _vopt_table_ref(ws), 'Workbook must include a native table'
    assert ws.print_area, 'Workbook must define a print area'
    assert ws.freeze_panes, 'Workbook must freeze panes'
    _assert_any_formulas(ws)
    actual_ranges = _vopt_dv_ranges(ws)
    assert actual_ranges, 'Workbook must include data validation'
    assert len(ws.conditional_formatting) >= 1
    for sheet in EXPECT['hidden_sheets']:
        assert wb[sheet].sheet_state in ('hidden', 'veryHidden')


def test_cross_output_caption_gate():
    prs = Presentation(_deck_path())
    deck_text = '\n'.join(_slide_text(slide) for slide in prs.slides)
    wb = load_workbook(_workbook_path(), data_only=False)
    ws = _spec_sheet(wb)
    workbook_text = _workbook_text(wb)
    assert 'p-118' in workbook_text and 'caption' in workbook_text and ('hold' in workbook_text or 'needs caption' in workbook_text)
    assert 'P-118' in deck_text
    assert 'caption' in deck_text.lower() and ('hold' in deck_text.lower() or 'required before export' in deck_text.lower() or 'needs caption' in deck_text.lower())
