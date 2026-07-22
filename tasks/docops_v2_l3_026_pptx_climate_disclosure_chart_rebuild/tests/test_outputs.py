import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_FILL

from verifier_utils import normalize_text, run_preflight

META_PATH = Path(os.environ.get("TASK_METADATA_PATH", "/tests/task_metadata.json"))
if not META_PATH.exists():
    META_PATH = Path(__file__).parent / "task_metadata.json"
META = json.loads(META_PATH.read_text(encoding="utf-8"))
EXPECT = META["verifier_expectations"]
INPUT_PATH = Path(os.environ.get("INPUT_PATH", META["input_path"]))
OUTPUT_PATH = Path(os.environ.get("OUTPUT_PATH", META["output_path"]))


def rgb_str(rgb):
    return None if rgb is None else str(rgb).upper()


def slide_text(slide, include_notes=True):
    parts = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            txt = shape.text.strip()
            if txt:
                parts.append(txt)
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    txt = cell.text.strip()
                    if txt:
                        parts.append(txt)
    if include_notes and getattr(slide, "has_notes_slide", False):
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            parts.append(notes)
    return "\n".join(parts)


def all_text(prs):
    return "\n".join(slide_text(slide, include_notes=True) for slide in prs.slides)


def slide_title(slide):
    candidates = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = shape.text.strip()
            if text:
                candidates.append((shape.top, shape.left, text.splitlines()[0]))
    assert candidates, "Slide has no title text"
    return sorted(candidates, key=lambda item: (item[0], item[1]))[0][2]


def title_shape(slide):
    title = slide_title(slide)
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = shape.text.strip()
            if text and text.splitlines()[0] == title:
                return shape
    raise AssertionError(f"Could not find title shape for {title!r}")


def table_rows(slide):
    tables = []
    for shape in slide.shapes:
        if getattr(shape, "has_table", False) and shape.has_table:
            rows = []
            for row in shape.table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            tables.append(rows)
    return tables


def has_table(slide):
    return any(getattr(shape, "has_table", False) and shape.has_table for shape in slide.shapes)


def assert_terms(title, text, terms):
    normalized = normalize_text(text)
    missing = [term for term in terms if normalize_text(term) not in normalized]
    assert not missing, f"{title}: missing required public terms {missing}"


def assert_any_terms(title, text, groups):
    normalized = normalize_text(text)
    missing = []
    for label, choices in groups.items():
        if not any(normalize_text(choice) in normalized for choice in choices):
            missing.append(label)
    assert not missing, f"{title}: missing required public concept(s) {missing}"


def background_rgb(slide):
    fill = slide.background.fill
    return rgb_str(fill.fore_color.rgb) if fill.type == MSO_FILL.SOLID else None


def has_accent_shape(slide):
    expected = EXPECT["style"]["accent_fill"]
    for shape in slide.shapes:
        fill = getattr(shape, "fill", None)
        if fill is not None and fill.type == MSO_FILL.SOLID and rgb_str(fill.fore_color.rgb) == expected:
            return True
    return False


def title_font_info(slide):
    shape = title_shape(slide)
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                return run.font.name, bool(run.font.bold), rgb_str(run.font.color.rgb)
    raise AssertionError("Title has no text run")


def speaker_notes(prs):
    return [
        slide.notes_slide.notes_text_frame.text.strip() if getattr(slide, "has_notes_slide", False) else ""
        for slide in prs.slides
    ]


def assert_table_header_style(slide):
    style = EXPECT["style"]
    for shape in slide.shapes:
        if getattr(shape, "has_table", False) and shape.has_table:
            for cell in shape.table.rows[0].cells:
                assert rgb_str(cell.fill.fore_color.rgb) == style["table_header_fill"]
                runs = [
                    run
                    for paragraph in cell.text_frame.paragraphs
                    for run in paragraph.runs
                    if run.text.strip()
                ]
                assert runs, "Table header cell has no text run"
                assert all(run.font.bold for run in runs)
                assert all(rgb_str(run.font.color.rgb) == style["table_header_font"] for run in runs)


def chart_on_slide(slide):
    charts = [shape.chart for shape in slide.shapes if getattr(shape, "has_chart", False)]
    assert len(charts) == 1, f"Expected exactly one chart, found {len(charts)}"
    return charts[0]


def slides_by_expected_title(prs):
    by_title = {slide_title(slide): slide for slide in prs.slides}
    missing = [title for title in EXPECT["slide_titles"] if title not in by_title]
    assert not missing, f"Missing required public slide(s): {missing}"
    return by_title


def test_output_exists_and_is_pptx():
    run_preflight(META, INPUT_PATH, OUTPUT_PATH)
    assert OUTPUT_PATH != INPUT_PATH


def test_slide_order_and_private_material_removed():
    prs = Presentation(OUTPUT_PATH)
    assert [slide_title(slide) for slide in prs.slides] == EXPECT["slide_titles"]
    text = normalize_text(all_text(prs))
    hits = [p for p in EXPECT["forbidden_phrases"] if normalize_text(p) in text]
    assert not hits, f"Forbidden phrases still present: {hits}"


def test_required_public_content_and_tables():
    prs = Presentation(OUTPUT_PATH)
    by_title = slides_by_expected_title(prs)
    slide_requirements = {
        "Northline Climate Disclosure Briefing": ["Climate Disclosure", "FY2026"],
        "Disclosure Boundary": ["Boundary", "2024"],
        "Corrected Emissions Inventory": ["Scope 2", "13,480", "Renewable", "48%", "Submitted"],
        "Emissions Trend Chart": ["Emissions"],
        "Assurance Status": ["Kara Saito", "Submitted"],
        "Transition Action Plan": ["Action"],
        "Supplier Engagement Register": ["Supplier"],
        "Board Actions": ["Board"],
        "Publication Style Guide": ["Aptos", "accent"],
    }
    for title, terms in slide_requirements.items():
        assert_terms(title, slide_text(by_title[title], include_notes=False), terms)
    assert_any_terms(
        "Transition Action Plan",
        slide_text(by_title["Transition Action Plan"], include_notes=False),
        {"transition work": ["Renewable", "emissions", "dashboard", "supplier", "efficiency"]},
    )
    assert_any_terms(
        "Board Actions",
        slide_text(by_title["Board Actions"], include_notes=False),
        {"board actions": ["Approve", "Receive", "Review", "Confirm", "governance"]},
    )
    for title in [
        "Corrected Emissions Inventory",
        "Assurance Status",
        "Transition Action Plan",
        "Supplier Engagement Register",
    ]:
        assert has_table(by_title[title]), f"{title}: expected a rebuilt native table"


def test_clustered_column_chart_rebuilt_from_corrected_values():
    prs = Presentation(OUTPUT_PATH)
    by_title = slides_by_expected_title(prs)
    spec = EXPECT["chart"]
    chart = chart_on_slide(by_title[spec["slide_title"]])
    assert chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED
    assert len(chart.series) == 1
    assert chart.series[0].name == spec["series_name"]
    assert [str(category.label) for category in chart.plots[0].categories] == spec["categories"]
    assert [round(float(value), 2) for value in chart.series[0].values] == [
        round(float(value), 2) for value in spec["values"]
    ]


def test_speaker_notes_cleaned_and_preserved():
    prs = Presentation(OUTPUT_PATH)
    notes = speaker_notes(prs)
    assert len(notes) == len(EXPECT["slide_titles"])
    forbidden = [normalize_text(p) for p in EXPECT["forbidden_phrases"]]
    for idx, note in enumerate(notes, 1):
        normalized = normalize_text(note)
        assert normalized, f"Slide {idx}: missing public speaker notes"
        assert len(note) <= 300, f"Slide {idx}: speaker notes are not concise"
        hits = [p for p in forbidden if p in normalized]
        assert not hits, f"Slide {idx}: speaker notes still contain forbidden/private material {hits}"


def test_style_migration_applied():
    prs = Presentation(OUTPUT_PATH)
    style = EXPECT["style"]
    for idx, slide in enumerate(prs.slides):
        expected_bg = style["title_background"] if idx == 0 else style["content_background"]
        expected_title_color = style["title_color_title_slide"] if idx == 0 else style["title_color_content"]
        assert background_rgb(slide) == expected_bg
        font_name, bold, color = title_font_info(slide)
        assert font_name == style["title_font"]
        assert bold
        assert color == expected_title_color
        assert has_accent_shape(slide)
        assert style["footer_text"] in slide_text(slide, include_notes=False)
        if slide_title(slide) in EXPECT["tables"]:
            assert_table_header_style(slide)


def test_source_artifact_was_not_modified():
    source = Presentation(INPUT_PATH)
    assert len(source.slides) > len(EXPECT["slide_titles"])
    text = normalize_text(all_text(source))
    missing = [p for p in EXPECT["source_must_contain"] if normalize_text(p) not in text]
    assert not missing, f"Source artifact no longer contains expected seed defects: {missing}"
